# -*- coding: utf-8 -*-
import glob
import os
from pptx import Presentation

# Find AHL files on desktop
desktop = os.path.join(os.path.expanduser("~"), "Desktop")
ahl_files = glob.glob(os.path.join(desktop, "AHL*.pptx"))

print("Found AHL files:")
for f in ahl_files:
    print(f"  {os.path.basename(f)}")

# Find V6.0 file
v6_file = None
for f in ahl_files:
    if "V6.0" in f or "V6" in f:
        v6_file = f
        break

if v6_file:
    print(f"\nReading: {v6_file}")
    try:
        prs = Presentation(v6_file)
        print(f"Total slides: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides[:30], 1):
            title_text = ""
            content_text = ""
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if text:
                        # Check if it's a title (usually first shape or short)
                        if not title_text and len(text) < 100:
                            title_text = text
                        elif len(text) > 100:
                            content_text = text[:200]
            
            if title_text or content_text:
                print(f"\n--- Slide {i} ---")
                if title_text:
                    print(f"Title: {title_text[:100]}")
                if content_text:
                    print(f"Content: {content_text[:150]}...")
    except Exception as e:
        print(f"Error: {e}")
else:
    print("No V6.0 file found")
