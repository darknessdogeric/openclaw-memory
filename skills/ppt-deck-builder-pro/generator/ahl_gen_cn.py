# -*- coding: utf-8 -*-
"""AHL PPT Generator - Chinese version"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATES = {
    "premium": {"primary": "#1e3a5f", "accent": "#d4af37", "background": "#0a0f1a"},
    "corporate": {"primary": "#1e40af", "accent": "#06b6d4", "background": "#0a0a0f"},
    "startup": {"primary": "#7c3aed", "accent": "#22d3ee", "background": "#0f0f1a"},
    "tech": {"primary": "#0c4a6e", "accent": "#38bdf8", "background": "#0c1929"}
}

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

class Generator:
    def __init__(self, template="premium"):
        self.t = TEMPLATES.get(template, TEMPLATES["premium"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def bg(self, slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*hex_to_rgb(self.t["background"]))
    
    def text(self, slide, txt, x, y, w, h, size=12, bold=False, color="#ffffff", align="center"):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        p = tb.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        return tb
    
    def title_slide(self, title, subtitle, presenter):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 2.5, 12, 1.5, 44, True, self.t["accent"])
        self.text(s, subtitle, 0.5, 4, 12, 1, 24, False, "#888888")
        self.text(s, presenter, 0.5, 5.5, 12, 1.5, 14, False, "#666666")
        return s
    
    def content_slide(self, title, bullets):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        for i, b in enumerate(bullets):
            self.text(s, "• " + b, 0.7, 1.3 + i * 0.55, 11.5, 0.5, 18, False, "#ffffff")
        return s
    
    def stats_slide(self, title, stats):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        for i, st in enumerate(stats):
            x = 0.7 + i * 3.2
            sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.8), Inches(1.8))
            sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(30,30,40)
            sh.line.color.rgb = RGBColor(*hex_to_rgb(self.t["accent"]))
            self.text(s, st["value"], x, 1.8, 2.8, 0.8, 36, True, self.t["accent"])
            self.text(s, st["label"], x, 2.6, 2.8, 0.5, 14, False, "#888888")
        return s
    
    def two_col_slide(self, title, lt, lb, rt, rb):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        self.text(s, lt, 0.5, 1.2, 5.5, 0.5, 20, True, "#ef4444")
        for i, b in enumerate(lb): self.text(s, "• " + b, 0.7, 1.8+i*0.5, 5, 0.4, 14, False, "#cccccc")
        self.text(s, rt, 7, 1.2, 5.5, 0.5, 20, True, "#10b981")
        for i, b in enumerate(rb): self.text(s, "• " + b, 7.2, 1.8+i*0.5, 5, 0.4, 14, False, "#cccccc")
        return s
    
    def timeline_slide(self, title, phases):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        for i, p in enumerate(phases):
            x = 0.5 + i * 3.1
            sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.8), Inches(4.5))
            sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(25,25,35)
            sh.line.color.rgb = RGBColor(*hex_to_rgb(self.t["accent"]))
            self.text(s, p["phase"], x, 2, 2.8, 0.4, 14, True, self.t["accent"])
            self.text(s, p["time"], x, 2.4, 2.8, 0.4, 12, False, "#888888")
            self.text(s, p["desc"], x, 2.9, 2.6, 2, 11, False, "#cccccc")
        return s
    
    def team_slide(self, title, members):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        for i, m in enumerate(members):
            r, c = i // 3, i % 3
            x, y = 0.5 + c * 4.2, 1.3 + r * 2.8
            sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.5))
            sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(30,30,40)
            self.text(s, m["name"], x, y+0.3, 3.8, 0.5, 20, True, "#ffffff")
            self.text(s, m["role"], x, y+0.8, 3.8, 0.4, 14, False, self.t["accent"])
            self.text(s, m["desc"], x+0.2, y+1.3, 3.4, 1, 11, False, "#aaaaaa")
        return s
    
    def table_slide(self, title, headers, rows):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        tbl = s.shapes.add_table(len(rows)+1, len(headers), Inches(0.5), Inches(1.3), Inches(12), Inches(0.5)).table
        for i, h in enumerate(headers):
            tbl.cell(0,i).text = h
            tbl.cell(0,i).fill.solid()
            tbl.cell(0,i).fill.fore_color.rgb = RGBColor(*hex_to_rgb(self.t["primary"]))
        for r_idx, r in enumerate(rows):
            for c_idx, v in enumerate(r):
                tbl.cell(r_idx+1, c_idx).text = str(v)
        return s
    
    def closing_slide(self, title, contact):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        self.text(s, title, 0.5, 2.5, 12, 1.5, 44, True, self.t["accent"])
        if contact: self.text(s, contact, 0.5, 4.5, 12, 1.5, 18, False, "#888888")
        return s
    
    def save(self, path): self.prs.save(path)


# 中文内容
CONTENT = {
    "title": "AHL去中心化住宿业交易生态协议",
    "subtitle": "人工智能 + Web3.0 融合创新项目",
    "presenter": "AHL智能科技有限公司（筹）\n项目负责人：张实\n电话：17760348653",
    "slides": [
        {"type": "title", "title": "AHL去中心化住宿业交易生态协议", "subtitle": "人工智能 + Web3.0 融合创新项目"},
        {"type": "content", "title": "项目概述 - 核心定位", "bullets": ["构建全球首个基于大语言模型的去中心化住宿业交易生态协议", "从货架经济向客户经济的范式革命", "从人找货到双向智能匹配", "年节省行业费用1000亿+"]},
        {"type": "stats", "title": "行业痛点 - 三大顽疾", "stats": [{"value": "20-30%", "label": "OTA垄断佣金"}, {"value": "1200亿", "label": "年损失佣金"}, {"value": "<3%", "label": "匹配转化率"}, {"value": "600亿", "label": "营销浪费"}]},
        {"type": "two_col", "title": "新旧模式对比", "lt": "OTA货架模式", "lb": ["商家把产品摆到货架", "消费者自己挑选", "信息过载匹配低效", "佣金20-30%"], "rt": "AHL客户模式", "rb": ["AI理解双方需求", "主动精准匹配", "双向智能匹配", "效率费3-5%"]},
        {"type": "content", "title": "技术架构 - 双AGENT+多SKILL", "bullets": ["C端AI管家：客户需求理解、偏好学习", "B端AI运营官：产品服务封装、动态定价", "匹配引擎：双向向量语义匹配", "9大AGENT x 87个SKILL可插拔模块", "向量匹配准确率95%+"]},
        {"type": "content", "title": "C2B直连交易闭环", "bullets": ["需求表达：客户自然语言描述", "AI理解：意图识别+实体提取", "向量匹配：检索Top-10匹配商家", "智能响应：实时房态+动态报价", "支付确认：直连支付"]},
        {"type": "timeline", "title": "四年发展规划", "phases": [{"phase": "Phase 1", "time": "0-6月", "desc": "协议核心研发\nAHL V1.0"}, {"phase": "Phase 2", "time": "6-12月", "desc": "AGENT矩阵\n100家商家"}, {"phase": "Phase 3", "time": "12-24月", "desc": "生态规模化\n500家商家"}, {"phase": "Phase 4", "time": "24-36月", "desc": "全球化拓展\n海外试点"}]},
        {"type": "team", "title": "核心团队 - 黄金三角", "members": [{"name": "张实", "role": "项目总控", "desc": "24年酒店业经验"}, {"name": "CTO待补充", "role": "AI技术架构", "desc": "华科AI博士"}, {"name": "CSO待补充", "role": "战略生态", "desc": "世界500强高管"}]},
        {"type": "content", "title": "申请支持事项", "bullets": ["算力支持：800万/年 - GPU集群74台", "办公场地：108万 - 500-800平米", "硬件设备：575万 - A100服务器"]},
        {"type": "table", "title": "经济效益预测", "headers": ["年份", "GMV", "平台收入", "行业节省", "净利润"], "rows": [["Year1", "5000万", "250万", "-", "-500万"], ["Year2", "3亿", "1500万", "150万", "300万"], ["Year3", "10亿", "5000万", "600万", "1500万"], ["Year4", "50亿", "2.5亿", "3000万", "8000万"]]},
        {"type": "content", "title": "社会效益 - 核心价值", "bullets": ["产业变革：打破OTA垄断", "技术创新：全球首个AI交易协议", "数据安全：商家数据自主可控", "就业：直接100人带动500人"]},
        {"type": "content", "title": "投资预算 - 四年总投入8000万", "bullets": ["算力成本：3000万(38%)", "研发投入：2000万(25%)", "硬件设备：1150万(14%)", "市场推广：800万(10%)", "运营生态：1050万(13%)"]},
        {"type": "table", "title": "风险分析与对策", "headers": ["风险类型", "风险描述", "等级", "应对措施"], "rows": [["技术风险", "大模型效果不达预期", "中", "多模型融合"], ["市场风险", "商家接受度低", "中", "免费试用"], ["竞争风险", "OTA巨头反击", "高", "差异化定位"], ["政策风险", "数据监管趋严", "低", "合规先行"], ["资金风险", "融资不及预期", "中", "控制烧钱"]]},
        {"type": "closing", "title": "恳请政府支持", "contact": "联系人：张实\n电话：17760348653\n邮箱：ericzhangshi@163.com"},
    ]
}


def generate(template="premium", output_dir="ppt_output"):
    os.makedirs(output_dir, exist_ok=True)
    g = Generator(template)
    
    # Cover
    g.title_slide(CONTENT["title"], CONTENT["subtitle"], CONTENT["presenter"])
    
    # Slides
    for sl in CONTENT["slides"]:
        if sl["type"] == "content": g.content_slide(sl["title"], sl["bullets"])
        elif sl["type"] == "stats": g.stats_slide(sl["title"], sl["stats"])
        elif sl["type"] == "two_col": g.two_col_slide(sl["title"], sl["lt"], sl["lb"], sl["rt"], sl["rb"])
        elif sl["type"] == "timeline": g.timeline_slide(sl["title"], sl["phases"])
        elif sl["type"] == "team": g.team_slide(sl["title"], sl["members"])
        elif sl["type"] == "table": g.table_slide(sl["title"], sl["headers"], sl["rows"])
        elif sl["type"] == "closing": g.closing_slide(sl["title"], sl["contact"])
    
    path = os.path.join(output_dir, "AHL-商业计划书-{}.pptx".format(template))
    g.save(path)
    return path


if __name__ == "__main__":
    t = sys.argv[1] if len(sys.argv) > 1 else "premium"
    print("Generated:", generate(t))
