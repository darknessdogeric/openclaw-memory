# -*- coding: utf-8 -*-
"""
AHL PPT Generator - 支持中文PPTX输出
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============== 模板定义 ==============
TEMPLATES = {
    "premium": {
        "name": "高端奢华",
        "primary": "#1e3a5f",
        "accent": "#d4af37",
        "background": "#0a0f1a",
    },
    "corporate": {
        "name": "企业商务",
        "primary": "#1e40af", 
        "accent": "#06b6d4",
        "background": "#0a0a0f",
    },
    "startup": {
        "name": "创业路演",
        "primary": "#7c3aed",
        "accent": "#22d3ee", 
        "background": "#0f0f1a",
    },
    "tech": {
        "name": "科技蓝",
        "primary": "#0c4a6e",
        "accent": "#38bdf8", 
        "background": "#0c1929",
    }
}

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

class AHLPPTGENERATOR:
    def __init__(self, template="premium"):
        self.template = TEMPLATES.get(template, TEMPLATES["premium"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def add_title_slide(self, title, subtitle="", presenter=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(2.5), 
                      Inches(12), Inches(1.5), font_size=44, bold=True,
                      color=self.template["accent"])
        
        if subtitle:
            self._add_text(slide, subtitle, Inches(0.5), Inches(4),
                          Inches(12), Inches(1), font_size=24, color="#888888")
        
        if presenter:
            self._add_text(slide, presenter, Inches(0.5), Inches(5.5),
                          Inches(12), Inches(1.5), font_size=14, color="#666666")
        
        return slide
    
    def add_content_slide(self, title, bullets):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        for i, bullet in enumerate(bullets):
            self._add_text(slide, "• " + bullet, Inches(0.7),
                          Inches(1.3 + i * 0.55), Inches(11.5), Inches(0.5),
                          font_size=18, color="#ffffff")
        
        return slide
    
    def add_two_column_slide(self, title, left_title, left_bullets, right_title, right_bullets):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        self._add_text(slide, left_title, Inches(0.5), Inches(1.2),
                      Inches(5.5), Inches(0.5), font_size=20, bold=True, color="#ef4444")
        
        for i, bullet in enumerate(left_bullets):
            self._add_text(slide, "• " + bullet, Inches(0.7),
                          Inches(1.8 + i * 0.5), Inches(5), Inches(0.4),
                          font_size=14, color="#cccccc")
        
        self._add_text(slide, right_title, Inches(7), Inches(1.2),
                      Inches(5.5), Inches(0.5), font_size=20, bold=True, color="#10b981")
        
        for i, bullet in enumerate(right_bullets):
            self._add_text(slide, "• " + bullet, Inches(7.2),
                          Inches(1.8 + i * 0.5), Inches(5), Inches(0.4),
                          font_size=14, color="#cccccc")
        
        return slide
    
    def add_stats_slide(self, title, stats):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        card_width = 2.8
        gap = 0.4
        start_x = 0.7
        
        for i, stat in enumerate(stats):
            x = Inches(start_x + i * (card_width + gap))
            y = Inches(1.5)
            
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                Inches(card_width), Inches(1.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 30, 40)
            shape.line.color.rgb = RGBColor(*hex_to_rgb(self.template["accent"]))
            
            self._add_text(slide, stat.get("value", ""), 
                          x, Inches(1.8), Inches(card_width), Inches(0.8),
                          font_size=36, bold=True, color=self.template["accent"])
            
            self._add_text(slide, stat.get("label", ""),
                          x, Inches(2.6), Inches(card_width), Inches(0.5),
                          font_size=14, color="#888888")
        
        return slide
    
    def add_timeline_slide(self, title, phases):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        phase_width = 2.8
        gap = 0.3
        start_x = 0.5
        
        for i, phase in enumerate(phases):
            x = Inches(start_x + i * (phase_width + gap))
            y = Inches(1.8)
            
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                Inches(phase_width), Inches(4.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(25, 25, 35)
            shape.line.color.rgb = RGBColor(*hex_to_rgb(self.template["accent"]))
            
            self._add_text(slide, phase.get("phase", ""),
                          x, Inches(2), Inches(phase_width), Inches(0.4),
                          font_size=14, bold=True, color=self.template["accent"])
            
            self._add_text(slide, phase.get("time", ""),
                          x, Inches(2.4), Inches(phase_width), Inches(0.4),
                          font_size=12, color="#888888")
            
            self._add_text(slide, phase.get("desc", ""),
                          x, Inches(2.9), Inches(phase_width - 0.2), Inches(2),
                          font_size=11, color="#cccccc")
        
        return slide
    
    def add_team_slide(self, title, members):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        member_width = 3.8
        gap = 0.4
        
        for i, member in enumerate(members):
            row = i // 3
            col = i % 3
            x = Inches(0.5 + col * (member_width + gap))
            y = Inches(1.3 + row * 2.8)
            
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                Inches(member_width), Inches(2.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 30, 40)
            
            self._add_text(slide, member.get("name", ""),
                          x, Inches(y.inches + 0.3), Inches(member_width), Inches(0.5),
                          font_size=20, bold=True, color="#ffffff")
            
            self._add_text(slide, member.get("role", ""),
                          x, Inches(y.inches + 0.8), Inches(member_width), Inches(0.4),
                          font_size=14, color=self.template["accent"])
            
            self._add_text(slide, member.get("desc", ""),
                          x + Inches(0.2), Inches(y.inches + 1.3), 
                          Inches(member_width - 0.4), Inches(1),
                          font_size=11, color="#aaaaaa")
        
        return slide
    
    def add_table_slide(self, title, headers, rows):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        table_rows = len(rows) + 1
        table_cols = len(headers)
        
        table = slide.shapes.add_table(table_rows, table_cols, 
            Inches(0.5), Inches(1.3), Inches(12), Inches(0.5)).table
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*hex_to_rgb(self.template["primary"]))
        
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
        
        return slide
    
    def add_closing_slide(self, title, contact=""):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(2.5),
                      Inches(12), Inches(1.5), font_size=44, bold=True,
                      color=self.template["accent"])
        
        if contact:
            self._add_text(slide, contact, Inches(0.5), Inches(4.5),
                          Inches(12), Inches(1.5), font_size=18, color="#888888")
        
        return slide
    
    def _set_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*hex_to_rgb(self.template["background"]))
    
    def _add_text(self, slide, text, x, y, width, height, 
                  font_size=12, bold=False, color="#ffffff", align="center"):
        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "left":
            p.alignment = PP_ALIGN.LEFT
        
        return textbox
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


# ============== 中文内容数据 ==============
AHL_CONTENT_CN = {
    "title": "AHL去中心化住宿业交易生态协议",
    "subtitle": "人工智能 + Web3.0 融合创新项目",
    "presenter": "AHL智能科技有限公司（筹）\n项目负责人：张实\n电话：17760348653",
    
    "slides": [
        {
            "type": "title",
            "title": "AHL去中心化住宿业交易生态协议",
            "subtitle": "人工智能 + Web3.0 融合创新项目",
        },
        {
            "type": "content",
            "title": "项目概述 - 核心定位",
            "bullets": [
                "构建全球首个基于大语言模型的去中心化住宿业交易生态协议",
                "从"货架经济"向"客户经济"的范式革命",
                "从"人找货"到"双向智能匹配"",
                "年节省行业费用¥1000亿+",
            ]
        },
        {
            "type": "stats",
            "title": "行业痛点 - 三大顽疾",
            "stats": [
                {"value": "20-30%", "label": "OTA垄断佣金"},
                {"value": "¥1200亿", "label": "年损失佣金"},
                {"value": "<3%", "label": "匹配转化率"},
                {"value": "¥600亿", "label": "营销浪费"},
            ]
        },
        {
            "type": "two_column",
            "title": "新旧模式对比",
            "left_title": "❌ OTA货架模式",
            "left_bullets": [
                "商家把产品摆到货架",
                "消费者自己挑选",
                "信息过载、匹配低效",
                "佣金20-30%",
            ],
            "right_title": "✓ AHL客户模式",
            "right_bullets": [
                "AI理解双方需求",
                "主动精准匹配",
                "双向智能匹配",
                "效率费3-5%",
            ]
        },
        {
            "type": "content",
            "title": "技术架构 - 双AGENT+多SKILL",
            "bullets": [
                "C端AI管家：客户需求理解、偏好学习、智能推荐",
                "B端AI运营官：产品服务封装、动态定价、智能营销",
                "匹配引擎：双向向量语义匹配、实时推荐",
                "9大AGENT × 87个SKILL：可插拔模块，按需订阅",
                "向量匹配准确率>95%，检索延迟<100ms",
            ]
        },
        {
            "type": "content",
            "title": "C2B直连交易闭环",
            "bullets": [
                "需求表达：客户自然语言描述需求",
                "AI理解：意图识别+实体提取+向量化",
                "向量匹配：检索Top-10匹配商家",
                "智能响应：实时房态+动态报价",
                "支付确认：直连支付，智能合约保障",
            ]
        },
        {
            "type": "timeline",
            "title": "四年发展规划",
            "phases": [
                {"phase": "Phase 1", "time": "0-6月", "desc": "协议核心研发\nAHL Protocol V1.0\n双AGENT MVP"},
                {"phase": "Phase 2", "time": "6-12月", "desc": "AGENT矩阵建设\n9大AGENT×87SKILL\n100家商家接入"},
                {"phase": "Phase 3", "time": "12-24月", "desc": "生态规模化\n500家商家\n月GMV>5000万"},
                {"phase": "Phase 4", "time": "24-36月", "desc": "全球化拓展\n海外试点\n国际标准"},
            ]
        },
        {
            "type": "team",
            "title": "核心团队 - 黄金三角",
            "members": [
                {"name": "张实", "role": "项目总控/发起人", "desc": "24年酒店业老兵\n战略/资源/政府关系"},
                {"name": "CTO [待补充]", "role": "AI技术架构构建者", "desc": "华科AI博士\n10年+AI研发经验"},
                {"name": "CSO [待补充]", "role": "战略与生态推动者", "desc": "世界500强高管\nPE/VC投资经验"},
            ]
        },
        {
            "type": "content",
            "title": "申请支持事项",
            "bullets": [
                "💻 算力支持：¥800万/年 - GPU集群74台",
                "🏢 办公场地：¥108万 - 500-800平米",
                "🖥️ 硬件设备：¥575万 - A100服务器",
            ]
        },
        {
            "type": "table",
            "title": "经济效益预测",
            "headers": ["年份", "GMV", "平台收入", "为行业节省", "净利润"],
            "rows": [
                ["Year 1", "¥5,000万", "¥250万", "-", "-¥500万"],
                ["Year 2", "¥3亿", "¥1,500万", "¥150万", "¥300万"],
                ["Year 3", "¥10亿", "¥5,000万", "¥600万", "¥1,500万"],
                ["Year 4", "¥50亿", "¥2.5亿", "¥3,000万", "¥8,000万"],
            ]
        },
        {
            "type": "content",
            "title": "社会效益 - 核心价值",
            "bullets": [
                "🏆 产业变革：打破OTA垄断，建立公平透明交易环境",
                "💡 技术创新：全球首个住宿业AI交易协议",
                "🔒 数据安全：商家数据自主可控，打破平台数据垄断",
                "👥 就业：直接就业100人，带动产业链500人",
            ]
        },
        {
            "type": "content",
            "title": "投资预算 - 四年总投入¥8000万",
            "bullets": [
                "💰 算力成本：¥3000万 (38%)",
                "💰 研发投入：¥2000万 (25%)",
                "💰 硬件设备：¥1150万 (14%)",
                "💰 市场推广：¥800万 (10%)",
                "💰 运营+生态：¥1050万 (13%)",
            ]
        },
        {
            "type": "table",
            "title": "风险分析与对策",
            "headers": ["风险类型", "风险描述", "等级", "应对措施"],
            "rows": [
                ["技术风险", "大模型效果不达预期", "中", "多模型融合"],
                ["市场风险", "商家接受度低", "中", "免费试用"],
                ["竞争风险", "OTA巨头反击", "高", "差异化定位"],
                ["政策风险", "数据监管趋严", "低", "合规先行"],
                ["资金风险", "融资不及预期", "中", "控制烧钱"],
            ]
        },
        {
            "type": "closing",
            "title": "恳请政府支持",
            "contact": "联系人：张实\n电话/微信：17760348653\n邮箱：ericzhangshi@163.com",
        },
    ]
}


def generate_ahl_pptx_cn(template="premium", output_dir="ppt_output"):
    os.makedirs(output_dir, exist_ok=True)
    
    generator = AHLPPTGENERATOR(template)
    content = AHL_CONTENT_CN
    
    # 封面
    generator.add_title_slide(
        content["title"],
        content["subtitle"],
        content["presenter"]
    )
    
    # 内容页
    for slide_data in content["slides"]:
        if slide_data["type"] == "content":
            generator.add_content_slide(slide_data["title"], slide_data["bullets"])
        elif slide_data["type"] == "stats":
            generator.add_stats_slide(slide_data["title"], slide_data["stats"])
        elif slide_data["type"] == "two_column":
            generator.add_two_column_slide(
                slide_data["title"],
                slide_data["left_title"],
                slide_data["left_bullets"],
                slide_data["right_title"],
                slide_data["right_bullets"]
            )
        elif slide_data["type"] == "timeline":
            generator.add_timeline_slide(slide_data["title"], slide_data["phases"])
        elif slide_data["type"] == "team":
            generator.add_team_slide(slide_data["title"], slide_data["members"])
        elif slide_data["type"] == "table":
            generator.add_table_slide(slide_data["title"], slide_data["headers"], slide_data["rows"])
        elif slide_data["type"] == "closing":
            generator.add_closing_slide(slide_data["title"], slide_data.get("contact", ""))
    
    output_path = os.path.join(output_dir, "AHL-商业计划书-{}.pptx".format(template))
    generator.save(output_path)
    
    return output_path


if __name__ == "__main__":
    template = sys.argv[1] if len(sys.argv) > 1 else "premium"
    output = generate_ahl_pptx_cn(template)
    print("Generated: {}".format(output))
