# -*- coding: utf-8 -*-
"""
高级PPT生成器 - NotebookLM/Swiss风格深度优化版
- 丰富背景设计
- 深度排版系统
- 详细文字描述
- 图形化数据展示
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# ==================== 颜色系统 ====================
COLORS = {
    # 主色
    "primary": "#111111",      # 纯黑
    "secondary": "#333333",    # 深灰
    "accent": "#FFCC00",       # 亮黄
    "accent_red": "#FF3333",   # 警示红
    
    # 背景
    "bg_white": "#FFFFFF",
    "bg_gray": "#F5F5F5",
    "bg_light": "#FAFAFA",
    "bg_dark": "#1A1A1A",
    
    # 文字
    "text_primary": "#111111",
    "text_secondary": "#666666",
    "text_tertiary": "#999999",
    "text_light": "#FFFFFF",
    
    # 功能色
    "success": "#10B981",
    "warning": "#F59E0B",
    "danger": "#EF4444",
    "info": "#3B82F6",
}

def hex_rgb(h):
    if not h or len(h) < 6:
        return (0, 0, 0)
    h = h.lstrip('#')
    if len(h) != 6:
        h = "000000"
    try:
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
    except:
        return (0, 0, 0)


# ==================== 背景系统 ====================
class BackgroundSystem:
    """高级背景系统"""
    
    @staticmethod
    def solid(slide, color=COLORS["bg_white"]):
        """纯色背景"""
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*hex_rgb(color))
    
    @staticmethod
    def gradient(slide, colors=["#111111", "#333333"]):
        """渐变背景"""
        # 简化版渐变 - 用深色填充
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*hex_rgb(colors[0]))
    
    @staticmethod
    def grid_pattern(slide, color="#000000", opacity=0.03):
        """网格图案"""
        # 添加网格线 - 通过形状模拟
        for i in range(12):
            left = Inches(i * 1.1)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, Inches(0), Inches(0.02), Inches(7.5))
            shape.fill.background()
            shape.line.color.rgb = RGBColor(*hex_rgb(color))
            shape.line.transparency = 100 - int(opacity * 100)
    
    @staticmethod
    def diagonal_lines(slide, color="#000000", spacing=20):
        """对角线纹理"""
        pass  # 可扩展
    
    @staticmethod
    def accent_bar(slide, position="top", height=0.05, color=COLORS["accent"]):
        """强调色条"""
        if position == "top":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(height))
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.5-height), Inches(13.333), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_rgb(color))


# ==================== 排版系统 ====================
class TypographySystem:
    """高级排版系统"""
    
    @staticmethod
    def title_main(slide, text, x, y, size=48, color=COLORS["text_primary"], width=12):
        """主标题 - 超大字号"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial Black"
        p.alignment = PP_ALIGN.LEFT
        return tb
    
    @staticmethod
    def title_section(slide, text, x, y, size=14, color=COLORS["text_tertiary"]):
        """章节标题 - 标签风格"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT
        return tb
    
    @staticmethod
    def headline(slide, text, x, y, size=28, color=COLORS["text_primary"], width=10):
        """副标题/段落标题"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT
        return tb
    
    @staticmethod
    def body(slide, text, x, y, size=12, color=COLORS["text_secondary"], width=10, line_spacing=0.5):
        """正文文本"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT
        return tb
    
    @staticmethod
    def label(slide, text, x, y, size=10, color=COLORS["text_tertiary"]):
        """标签/分类"""
        return TypographySystem.title_section(slide, text, x, y, size, color)
    
    @staticmethod
    def big_number(slide, text, x, y, size=56, color=COLORS["accent"]):
        """大数字强调"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3), Inches(1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial Black"
        p.alignment = PP_ALIGN.LEFT
        return tb
    
    @staticmethod
    def bullet_list(slide, items, x, y, size=12, color=COLORS["text_secondary"], 
                   bullet="•", indent=0.3, spacing=0.45):
        """项目符号列表"""
        for i, item in enumerate(items):
            tb = slide.shapes.add_textbox(Inches(x + indent), Inches(y + i * spacing), 
                                         Inches(10), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{bullet} {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = RGBColor(*hex_rgb(color))
            p.font.name = "Arial"
            p.alignment = PP_ALIGN.LEFT
        return y + len(items) * spacing


# ==================== 图形元素 ====================
class GraphicElements:
    """图形元素系统"""
    
    @staticmethod
    def card(slide, x, y, w, h, color=COLORS["bg_gray"], border_color=None, radius=0.1):
        """卡片"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_rgb(color))
        if border_color:
            shape.line.color.rgb = RGBColor(*hex_rgb(border_color))
            shape.line.width = Pt(1)
        return shape
    
    @staticmethod
    def stat_card(slide, x, y, w, h, number, label, number_size=36, label_size=11, 
                  number_color=COLORS["text_primary"], label_color=COLORS["text_tertiary"]):
        """数据卡片"""
        # 背景
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_rgb("#F8F8F8"))
        
        # 数字
        TypographySystem.big_number(slide, number, x + 0.2, y + 0.3, number_size, number_color)
        
        # 标签
        TypographySystem.body(slide, label, x + 0.2, y + h - 0.5, label_size, label_color)
        
        return shape
    
    @staticmethod
    def comparison_box(slide, x, y, w, h, title, items, is_positive=False, title_size=14):
        """对比框"""
        # 背景色
        bg_color = "#E8F5E9" if is_positive else "#FFEBEE"
        text_color = COLORS["success"] if is_positive else COLORS["danger"]
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_rgb(bg_color))
        
        # 标题
        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(w - 0.4), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb(text_color))
        
        # 项目
        for i, item in enumerate(items):
            TypographySystem.body(slide, f"• {item}", x + 0.2, y + 0.7 + i * 0.35, 10, COLORS["text_secondary"])
        
        return shape
    
    @staticmethod
    def timeline_item(slide, x, y, phase, time, description, is_active=True):
        """时间线项目"""
        # 圆点
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(y), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(*hex_rgb(COLORS["accent"]))
        
        # 阶段
        TypographySystem.title_section(slide, phase, x, y + 0.2, 11, COLORS["text_tertiary"])
        
        # 时间
        TypographySystem.body(slide, time, x, y + 0.5, 10, COLORS["text_tertiary"])
        
        # 描述
        TypographySystem.body(slide, description, x, y + 0.8, 10, COLORS["text_secondary"], width=2.5)
        
        return dot
    
    @staticmethod
    def progress_bar(slide, x, y, w, h, percent, color=COLORS["accent"]):
        """进度条"""
        # 背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*hex_rgb("#E0E0E0"))
        
        # 进度
        if percent > 0:
            fill_w = w * percent / 100
            fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(fill_w), Inches(h))
            fill.fill.solid()
            fill.fill.fore_color.rgb = RGBColor(*hex_rgb(color))
        
        return bg


# ==================== 主PPT生成器 ====================
class AdvancedPPTGEN:
    """高级PPT生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.bg = BackgroundSystem()
        self.typo = TypographySystem()
        self.Graphic = GraphicElements()
    
    def save(self, path):
        self.prs.save(path)
        return path
    
    # ---------- 封面页 ----------
    def add_cover_deep(self, title, subtitle, punchline, category="PROJECT"):
        """深度封面页"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 背景
        self.bg.solid(s, COLORS["bg_white"])
        self.bg.accent_bar(s, "top", 0.08, COLORS["accent"])
        
        # 分类标签
        self.typo.label(s, f"{category}", 0.5, 0.35)
        
        # 主标题
        self.typo.title_main(s, title, 0.5, 0.8, 64, COLORS["primary"])
        
        # 副标题
        self.typo.body(s, subtitle, 0.5, 2.6, 20, COLORS["text_secondary"], width=8)
        
        # Punchline - 超大黄色数字
        self.typo.big_number(s, punchline, 0.5, 4.2, 80, COLORS["accent"])
        
        # 底部装饰线
        line = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.8), Inches(3), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*hex_rgb(COLORS["accent"]))
        
        return s
    
    # ---------- 章节页 ----------
    def add_section_deep(self, num, title, subtitle=""):
        """深度章节页"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self.bg.solid(s, COLORS["bg_white"])
        
        # 章节号 - 大字
        self.typo.big_number(s, f"0{num}", 0.5, 1.5, 72, COLORS["accent"])
        
        # 标题
        self.typo.title_main(s, title, 0.5, 3.2, 42, COLORS["primary"])
        
        # 副标题
        if subtitle:
            self.typo.body(s, subtitle, 0.5, 4.2, 16, COLORS["text_secondary"])
        
        return s
    
    # ---------- 统计页 ----------
    def add_stats_deep(self, title, stats_data):
        """深度统计页"""
        s = self.add_section_deep("1", "PROBLEM", title.replace("问题", "").replace("数据", ""))
        
        # 副标题
        self.typo.headline(s, title, 0.5, 1.8, 24)
        
        # 统计卡片
        for i, (num, label, desc) in enumerate(stats_data):
            x = 0.5 + i * 3.2
            self.Graphic.stat_card(s, x, 2.6, 2.8, 1.8, num, label, 32, 10)
            # 描述
            self.typo.body(s, desc, x, 4.5, 9, COLORS["text_tertiary"])
        
        return s
    
    # ---------- 对比页 ----------
    def add_comparison_deep(self, title, left_title, left_items, right_title, right_items):
        """深度对比页"""
        s = self.add_section_deep("2", "SOLUTION", "解决方案")
        
        self.typo.headline(s, title, 0.5, 1.8, 24)
        
        # 左侧对比框
        self.Graphic.comparison_box(s, 0.5, 2.5, 5.8, 4.2, left_title, left_items, False)
        
        # 右侧对比框
        self.Graphic.comparison_box(s, 7, 2.5, 5.8, 4.2, right_title, right_items, True)
        
        return s
    
    # ---------- 技术架构页 ----------
    def add_tech_deep(self):
        """深度技术页"""
        s = self.add_section_deep("3", "TECHNOLOGY", "技术架构")
        
        self.typo.headline(s, "双AGENT智能架构", 0.5, 1.8, 24)
        
        # 三个层次
        layers = [
            ("01", "C端AI管家", "需求理解\n偏好学习\n智能推荐\n向量画像"),
            ("02", "匹配引擎", "意图识别\n语义检索\n实时推荐\n95%+准确率"),
            ("03", "B端AI运营官", "产品封装\n动态定价\n智能营销\n运营自动化"),
        ]
        
        for i, (num, name, desc) in enumerate(layers):
            x = 0.5 + i * 4.3
            
            # 数字标签
            self.typo.big_number(s, num, x, 2.5, 28, COLORS["accent"])
            
            # 名称
            self.typo.headline(s, name, x, 3.0, 16)
            
            # 描述
            self.typo.body(s, desc, x, 3.5, 10, COLORS["text_secondary"], width=3.8)
        
        return s
    
    # ---------- 产品矩阵页 ----------
    def add_product_deep(self):
        """深度产品页"""
        s = self.add_section_deep("4", "PRODUCT", "产品矩阵")
        
        self.typo.headline(s, "9大AGENT × 87 SKILL", 0.5, 1.8, 24)
        
        # 9个AGENT卡片
        agents = [
            ("运营AGENT", "客房/入住/客情维护"),
            ("营销AGENT", "获客/线索/RFP应答"),
            ("收益AGENT", "动态定价/需求预测"),
            ("宴会AGENT", "婚宴顾问/MICE管理"),
            ("工程AGENT", "预测维护/能源管理"),
            ("安防AGENT", "AI视频/应急处理"),
            ("财务AGENT", "自动对账/智能报表"),
            ("人资AGENT", "智能排班/培训助手"),
            ("供应AGENT", "采购/库存管理"),
        ]
        
        for i, (name, desc) in enumerate(agents):
            row, col = i // 3, i % 3
            x = 0.5 + col * 4.3
            y = 2.5 + row * 1.55
            
            # 卡片
            self.Graphic.card(s, x, y, 4, 1.3, COLORS["bg_gray"], "#DDDDDD")
            
            # 名称
            self.typo.headline(s, name, x + 0.15, y + 0.15, 14)
            
            # 描述
            self.typo.body(s, desc, x + 0.15, y + 0.65, 10, COLORS["text_tertiary"])
        
        # SKILL说明
        self.typo.body(s, "× 87个可插拔SKILL模块 | 按需订阅，灵活组合", 0.5, 7, 11, COLORS["text_tertiary"])
        
        return s
    
    # ---------- 市场页 ----------
    def add_market_deep(self):
        """深度市场页"""
        s = self.add_section_deep("5", "MARKET", "市场机会")
        
        self.typo.headline(s, "万亿级住宿业市场", 0.5, 1.8, 24)
        
        # 数据展示
        market_data = [
            ("¥6000亿", "住宿业市场规模", "年市场规模"),
            ("70%", "在线预订占比", "数字化程度"),
            ("¥1200亿", "OTA佣金池", "年佣金支出"),
            ("¥1000亿+", "AHL可释放价值", "节省+效率提升"),
        ]
        
        for i, (num, label, sub) in enumerate(market_data):
            x = 0.5 + i * 3.2
            
            # 数字
            self.Graphic.stat_card(s, x, 2.5, 2.9, 1.6, num, label, 26, 10)
            
            # 副标签
            self.typo.body(s, sub, x, 4.2, 9, COLORS["accent"])
        
        # 市场说明
        self.typo.body(s, 
            "中国住宿业在线渗透率持续提升，AI技术将重构交易链路，",
            0.5, 5.2, 12, COLORS["text_secondary"])
        self.typo.body(s, 
            "AHL协议可帮助酒店节省80%佣金成本，提升3倍转化效率。",
            0.5, 5.6, 12, COLORS["text_secondary"])
        
        return s
    
    # ---------- 商业模式页 ----------
    def add_business_deep(self):
        """深度商业模式页"""
        s = self.add_section_deep("6", "BUSINESS", "商业模式")
        
        # 左侧 - 收入
        self.typo.headline(s, "收入结构", 0.5, 2.3, 18)
        
        revenue = [("技术服务费 2-3%", "60%", "主要收入"),
                  ("AI增值服务", "25%", "增值订阅"),
                  ("数据服务", "10%", "数据洞察"),
                  ("推广广告", "5%", "精准推荐")]
        
        for i, (item, percent, note) in enumerate(revenue):
            y = 2.7 + i * 0.45
            self.typo.body(s, f"• {item}", 0.5, y, 11, COLORS["text_primary"])
            self.typo.body(s, percent, 4.5, y, 11, COLORS["accent"])
        
        # 右侧 - 成本
        self.typo.headline(s, "成本构成", 7, 2.3, 18)
        
        costs = [("服务器/云", "40%", "基础设施"),
                ("AI模型调用", "30%", "技术成本"),
                ("运维人员", "20%", "运营"),
                ("其他", "10%", "行政杂项")]
        
        for i, (item, percent, note) in enumerate(costs):
            y = 2.7 + i * 0.45
            self.typo.body(s, f"• {item}", 7, y, 11, COLORS["text_primary"])
            self.typo.body(s, percent, 10.5, y, 11, COLORS["info"])
        
        # 底部说明
        self.Graphic.card(s, 0.5, 5.8, 12.3, 1.2, COLORS["bg_gray"])
        self.typo.body(s, "AHL协议将传统OTA 20-30%的佣金成本降至2-3%，", 0.7, 5.95, 11, COLORS["text_primary"])
        self.typo.body(s, "为酒店业每年节省超1000亿元费用。", 0.7, 6.25, 11, COLORS["accent"])
        
        return s
    
    # ---------- 发展规划页 ----------
    def add_roadmap_deep(self):
        """深度发展规划页"""
        s = self.add_section_deep("7", "ROADMAP", "发展规划")
        
        self.typo.headline(s, "四年战略路线图", 0.5, 1.8, 24)
        
        # 时间线
        phases = [
            ("PHASE 1", "0-6月", "协议核心研发\nAHL V1.0发布\n50家测试商家"),
            ("PHASE 2", "6-12月", "AGENT矩阵建设\n9大AGENT×87SKILL\n500家商家接入"),
            ("PHASE 3", "12-24月", "生态规模化\n5000家商家\n月GMV>5000万"),
            ("PHASE 4", "24-36月", "全球化拓展\n东南亚/日本试点\n行业标准制定"),
        ]
        
        for i, (phase, time, desc) in enumerate(phases):
            x = 0.5 + i * 3.2
            self.Graphic.timeline_item(s, x, 2.5, phase, time, desc, i <= 1)
        
        return s
    
    # ---------- 团队页 ----------
    def add_team_deep(self):
        """深度团队页"""
        s = self.add_section_deep("8", "TEAM", "核心团队")
        
        self.typo.headline(s, "黄金三角组合", 0.5, 1.8, 24)
        
        members = [
            ("张实", "项目总控/发起人", "24年酒店业老兵\n战略/资源/政府关系\n曾任多家酒店集团高管", "创始人"),
            ("CTO", "AI技术架构构建者", "华科AI博士\n10年+AI研发经验\n20+AI工程落地", "联合创始人"),
            ("CSO", "战略与生态推动者", "世界500强高管\nPE/VC投资经验\nIPO全流程经验", "联合创始人"),
        ]
        
        for i, (name, role, desc, tag) in enumerate(members):
            x = 0.5 + i * 4.3
            
            # 头像圆
            avatar = s.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 1.3), Inches(2.5), Inches(1.4), Inches(1.4))
            avatar.fill.solid()
            avatar.fill.fore_color.rgb = RGBColor(*hex_rgb(COLORS["bg_gray"]))
            
            # 名字
            self.typo.headline(s, name, x + 0.5, 4.1, 20)
            
            # 角色
            self.typo.body(s, role, x + 0.5, 4.5, 11, COLORS["accent"])
            
            # 描述
            self.typo.body(s, desc, x + 0.5, 5.0, 10, COLORS["text_secondary"])
            
            # 标签
            self.typo.label(s, tag, x + 0.5, 6.3, 9, COLORS["text_tertiary"])
        
        return s
    
    # ---------- 融资页 ----------
    def add_funding_deep(self):
        """深度融资页"""
        s = self.add_section_deep("9", "FUNDING", "融资计划")
        
        # 融资信息框 - 黑色背景
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.3), Inches(5.5), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*hex_rgb(COLORS["primary"]))
        
        # 标签
        self.typo.label(s, "SEED ROUND", 0.7, 2.5, 10, COLORS["accent"])
        
        # 金额
        self.typo.title_main(s, "¥500-800万", 0.7, 2.9, 42, COLORS["text_light"])
        
        # 详情
        self.typo.body(s, "出让股份 15-20%", 0.7, 4.3, 12, COLORS["text_tertiary"])
        self.typo.body(s, "投后估值 3000-5000万", 0.7, 4.6, 12, COLORS["text_tertiary"])
        
        # 资金用途
        self.typo.headline(s, "资金用途", 6.5, 2.5, 16)
        uses = [("技术团队 40%", "核心研发"),
                ("数据体系 30%", "向量数据库"),
                ("市场拓展 20%", "商家拓展"),
                ("合规储备 10%", "法务/合规")]
        
        for i, (item, note) in enumerate(uses):
            self.typo.body(s, f"• {item}", 6.5, 2.9 + i * 0.4, 11, COLORS["text_primary"])
        
        # 联系人
        self.typo.headline(s, "联系人", 6.5, 4.8, 16)
        self.typo.body(s, "张实", 6.5, 5.2, 12, COLORS["text_primary"])
        self.typo.body(s, "17760348653", 6.5, 5.5, 11, COLORS["text_secondary"])
        self.typo.body(s, "ericzhangshi@163.com", 6.5, 5.8, 11, COLORS["text_secondary"])
        
        return s
    
    # ---------- 结尾页 ----------
    def add_closing_deep(self):
        """深度结尾页"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self.bg.solid(s, COLORS["primary"])
        
        # 标题
        self.typo.title_main(s, "加入AHL生态", 0.5, 2, 36, COLORS["text_light"])
        
        # 副标题
        self.typo.body(s, "打破OTA垄断 · 建立公平交易 · 共享AI红利", 
                     0.5, 3, 16, COLORS["text_tertiary"])
        
        # 共赢大字
        self.typo.big_number(s, "共赢", 0.5, 4, 72, COLORS["accent"])
        
        # 行动号召
        cta = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(4), Inches(0.8))
        cta.fill.solid()
        cta.fill.fore_color.rgb = RGBColor(*hex_rgb(COLORS["accent"]))
        
        self.typo.body(s, "联系我们开始合作", 1, 6.35, 14, COLORS["primary"])
        
        return s


def generate_advanced_ppt(output="ppt_output/AHL-Advanced.pptx"):
    """生成高级PPT"""
    ppt = AdvancedPPTGEN()
    
    # 1. 封面
    ppt.add_cover_deep(
        "AHL",
        "去中心化住宿业交易生态协议",
        "范式革命",
        "PROJECT BLUEPRINT"
    )
    
    # 2. 问题
    ppt.add_stats_deep(
        "行业三大顽疾",
        [
            ("20-30%", "OTA垄断佣金", "携程/美团占据80%市场"),
            ("¥1200亿", "年损失佣金", "酒店业每年损失"),
            ("<3%", "匹配转化率", "传统搜索效率低"),
            ("¥600亿", "营销浪费", "获客成本持续攀升"),
        ]
    )
    
    # 3. 解决方案
    ppt.add_comparison_deep(
        "新旧模式对比",
        "❌ OTA货架模式",
        ["中心化垄断，数据不透明", "佣金20-30%，成本高昂", "人找货，匹配效率低", "客户数据归平台所有"],
        "✓ AHL客户模式",
        ["去中心化，直连交易", "效率费2-3%，成本革命", "AI双向匹配，95%+准确率", "数据商户自主可控"]
    )
    
    # 4. 技术
    ppt.add_tech_deep()
    
    # 5. 产品
    ppt.add_product_deep()
    
    # 6. 市场
    ppt.add_market_deep()
    
    # 7. 商业
    ppt.add_business_deep()
    
    # 8. 规划
    ppt.add_roadmap_deep()
    
    # 9. 团队
    ppt.add_team_deep()
    
    # 10. 融资
    ppt.add_funding_deep()
    
    # 11. 结尾
    ppt.add_closing_deep()
    
    return ppt.save(output)


if __name__ == "__main__":
    output = generate_advanced_ppt()
    print(f"Generated: {output}")
