# -*- coding: utf-8 -*-
"""
AHL PPT Generator - 支持HTML和PPTX双输出
支持多种模板风格，内容深度丰富
"""
import os
import sys
import codecs
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============== 模板定义 ==============
TEMPLATES = {
    "corporate": {
        "name": "企业商务",
        "primary": "#1e40af",
        "secondary": "#3b82f6", 
        "accent": "#06b6d4",
        "background": "#0a0a0f",
        "font_title": "Arial",
        "font_body": "Arial",
        "style": "现代商务"
    },
    "premium": {
        "name": "高端奢华",
        "primary": "#1e3a5f",
        "secondary": "#2d5a87",
        "accent": "#d4af37",
        "background": "#0a0f1a",
        "font_title": "Arial Black",
        "font_body": "Arial",
        "style": "奢华金边"
    },
    "startup": {
        "name": "创业路演",
        "primary": "#7c3aed",
        "secondary": "#a855f7",
        "accent": "#22d3ee",
        "background": "#0f0f1a",
        "font_title": "Arial",
        "font_body": "Arial",
        "style": "活力渐变"
    },
    "minimal": {
        "name": "极简白",
        "primary": "#1f2937",
        "secondary": "#374151",
        "accent": "#10b981",
        "background": "#ffffff",
        "font_title": "Arial",
        "font_body": "Arial",
        "style": "纯净极简"
    },
    "tech": {
        "name": "科技蓝",
        "primary": "#0c4a6e",
        "secondary": "#0369a1",
        "accent": "#38bdf8",
        "background": "#0c1929",
        "font_title": "Consolas",
        "font_body": "Consolas",
        "style": "科技感"
    }
}

# ============== 颜色转换 ==============
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# ============== PPT生成器 ==============
class AHLPPTGENERATOR:
    def __init__(self, template="corporate"):
        self.template = TEMPLATES.get(template, TEMPLATES["corporate"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def add_title_slide(self, title, subtitle="", presenter=""):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(2.5), 
                      Inches(12), Inches(1.5), font_size=44, bold=True,
                      color=self.template["accent"])
        
        if subtitle:
            self._add_text(slide, subtitle, Inches(0.5), Inches(4),
                          Inches(12), Inches(1), font_size=24,
                          color=self.template.get("secondary", "#666666"))
        
        if presenter:
            self._add_text(slide, presenter, Inches(0.5), Inches(6),
                          Inches(12), Inches(1), font_size=14,
                          color="#888888")
        
        return slide
    
    def add_content_slide(self, title, bullets):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        for i, bullet in enumerate(bullets):
            self._add_text(slide, "• " + bullet, Inches(0.7),
                          Inches(1.3 + i * 0.55), Inches(11.5), Inches(0.5),
                          font_size=18, color="#ffffff")
        
        return slide
    
    def add_two_column_slide(self, title, left_title, left_bullets, right_title, right_bullets):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        self._add_text(slide, left_title, Inches(0.5), Inches(1.2),
                      Inches(5.5), Inches(0.5), font_size=20, bold=True,
                      color="#ef4444")
        
        for i, bullet in enumerate(left_bullets):
            self._add_text(slide, "• " + bullet, Inches(0.7),
                          Inches(1.8 + i * 0.5), Inches(5), Inches(0.4),
                          font_size=14, color="#cccccc")
        
        self._add_text(slide, right_title, Inches(7), Inches(1.2),
                      Inches(5.5), Inches(0.5), font_size=20, bold=True,
                      color="#10b981")
        
        for i, bullet in enumerate(right_bullets):
            self._add_text(slide, "• " + bullet, Inches(7.2),
                          Inches(1.8 + i * 0.5), Inches(5), Inches(0.4),
                          font_size=14, color="#cccccc")
        
        return slide
    
    def add_stats_slide(self, title, stats):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        card_width = 2.8
        gap = 0.4
        start_x = 0.7
        
        for i, stat in enumerate(stats):
            x = Inches(start_x + i * (card_width + gap))
            y = Inches(1.5)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                Inches(card_width), Inches(1.8)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 30, 40)
            shape.line.color.rgb = RGBColor(*hex_to_rgb(self.template["accent"]))
            
            self._add_text(slide, stat.get("value", ""), 
                          x, Inches(1.8), Inches(card_width), Inches(0.8),
                          font_size=36, bold=True,
                          color=self.template["accent"])
            
            self._add_text(slide, stat.get("label", ""),
                          x, Inches(2.6), Inches(card_width), Inches(0.5),
                          font_size=14, color="#888888")
        
        return slide
    
    def add_timeline_slide(self, title, phases):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        phase_width = 2.8
        gap = 0.3
        start_x = 0.5
        
        for i, phase in enumerate(phases):
            x = Inches(start_x + i * (phase_width + gap))
            y = Inches(1.8)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                Inches(phase_width), Inches(4.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(25, 25, 35)
            shape.line.color.rgb = RGBColor(*hex_to_rgb(self.template["accent"]))
            
            self._add_text(slide, phase.get("phase", ""),
                          x, Inches(2), Inches(phase_width), Inches(0.4),
                          font_size=14, bold=True,
                          color=self.template["accent"])
            
            self._add_text(slide, phase.get("time", ""),
                          x, Inches(2.4), Inches(phase_width), Inches(0.4),
                          font_size=12, color="#888888")
            
            self._add_text(slide, phase.get("desc", ""),
                          x, Inches(2.9), Inches(phase_width - 0.2), Inches(2),
                          font_size=11, color="#cccccc")
        
        return slide
    
    def add_team_slide(self, title, members):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        member_width = 3.8
        gap = 0.4
        cols = 3
        
        for i, member in enumerate(members):
            row = i // cols
            col = i % cols
            x = Inches(0.5 + col * (member_width + gap))
            y = Inches(1.3 + row * 2.8)
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
                Inches(member_width), Inches(2.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 30, 40)
            
            self._add_text(slide, member.get("name", ""),
                          x, Inches(y.inches + 0.3), Inches(member_width), Inches(0.5),
                          font_size=20, bold=True, color="#ffffff")
            
            self._add_text(slide, member.get("role", ""),
                          x, Inches(y.inches + 0.8), Inches(member_width), Inches(0.4),
                          font_size=14, color=self.template["accent"])
            
            self._add_text(slide, member.get("desc", ""),
                          x + Inches(0.2), Inches(y.inches + 1.3), 
                          Inches(member_width - 0.4), Inches(1),
                          font_size=11, color="#aaaaaa")
        
        return slide
    
    def add_table_slide(self, title, headers, rows):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(0.3),
                      Inches(12), Inches(0.8), font_size=32, bold=True,
                      color=self.template["accent"])
        
        table_rows = len(rows) + 1
        table_cols = len(headers)
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(12)
        height = Inches(0.5)
        
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
        
        col_width = width / table_cols
        for i in range(table_cols):
            table.cell(0, i).width = Inches(12.0/table_cols)
            # table.columns[i].width = col_width
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*hex_to_rgb(self.template["primary"]))
        
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = cell_data
        
        return slide
    
    def add_closing_slide(self, title, contact=""):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_background(slide)
        
        self._add_text(slide, title, Inches(0.5), Inches(2.5),
                      Inches(12), Inches(1.5), font_size=44, bold=True,
                      color=self.template["accent"])
        
        if contact:
            self._add_text(slide, contact, Inches(0.5), Inches(4.5),
                          Inches(12), Inches(1), font_size=18,
                          color="#888888")
        
        return slide
    
    def _set_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*hex_to_rgb(self.template["background"]))
    
    def _add_text(self, slide, text, x, y, width, height, 
                  font_size=12, bold=False, color="#ffffff", align="center"):
        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "left":
            p.alignment = PP_ALIGN.LEFT
        
        return textbox
    
    def save(self, filename):
        self.prs.save(filename)
        return filename


# ============== AHL内容数据 ==============
AHL_CONTENT = {
    "title": "AHL去中心化住宿业交易生态协议",
    "subtitle": "人工智能 + Web3.0 融合创新项目",
    "presenter": "AHL Intelligent Technology Ltd.\nProject Lead: Zhang Shi\nPhone: 17760348653",
    
    "slides": [
        {
            "type": "title",
            "title": "AHL Decentralized Hospitality Transaction Protocol",
            "subtitle": "AI + Web3.0 Innovation Project",
        },
        {
            "type": "content",
            "title": "Project Overview - Core Positioning",
            "bullets": [
                "World's first LLM-based decentralized hospitality transaction protocol",
                "Paradigm shift from 'Shelf Economy' to 'Customer Economy'",
                "From 'People finding products' to 'Bidirectional AI Matching'",
                "Save industry costs of 100B+ RMB annually",
            ]
        },
        {
            "type": "stats",
            "title": "Industry Pain Points",
            "stats": [
                {"value": "20-30%", "label": "OTA Commission"},
                {"value": "120B RMB", "label": "Annual Commission Loss"},
                {"value": "<3%", "label": "Match Conversion Rate"},
                {"value": "60B RMB", "label": "Marketing Waste"},
            ]
        },
        {
            "type": "two_column",
            "title": "Old vs New Model",
            "left_title": "OTA Shelf Model",
            "left_bullets": [
                "Merchants list products on shelves",
                "Consumers browse and choose",
                "Information overload, low efficiency",
                "Commission 20-30%",
            ],
            "right_title": "AHL Customer Model",
            "right_bullets": [
                "AI understands both parties",
                "Proactive precise matching",
                "Bidirectional intelligent matching",
                "Efficiency fee 3-5%",
            ]
        },
        {
            "type": "content",
            "title": "Technical Architecture",
            "bullets": [
                "C-End AI Butler: Customer demand understanding, preference learning",
                "B-End AI Operator: Product service packaging, dynamic pricing",
                "Matching Engine: Bidirectional vector semantic matching",
                "9 Agents x 87 Skills: Pluggable modules, subscribe as needed",
                "Vector matching accuracy >95%, latency <100ms",
            ]
        },
        {
            "type": "content",
            "title": "C2B Direct Transaction Loop",
            "bullets": [
                "Demand Expression: Natural language customer requirements",
                "AI Understanding: Intent recognition + entity extraction",
                "Vector Matching: Retrieve Top-10 matched merchants",
                "Smart Response: Real-time room availability + dynamic pricing",
                "Payment Confirmation: Direct payment, smart contract security",
            ]
        },
        {
            "type": "timeline",
            "title": "Four-Year Development Plan",
            "phases": [
                {"phase": "Phase 1", "time": "0-6 Months", "desc": "Protocol Core R&D\nAHL Protocol V1.0\nDual Agent MVP"},
                {"phase": "Phase 2", "time": "6-12 Months", "desc": "Agent Matrix\n9 Agents x 87 Skills\n100 Hotels"},
                {"phase": "Phase 3", "time": "12-24 Months", "desc": "Ecosystem Scale\n500 Hotels\nGMV > 50M RMB/month"},
                {"phase": "Phase 4", "time": "24-36 Months", "desc": "Global Expansion\nOverseas Pilot\nInternational Standards"},
            ]
        },
        {
            "type": "team",
            "title": "Core Team",
            "members": [
                {"name": "Zhang Shi", "role": "Project Lead", "desc": "24 years hospitality industry\nStrategy/Resources/Gov Relations"},
                {"name": "CTO [TBD]", "role": "AI Architecture", "desc": "AI PhD\n10+ years AI R&D"},
                {"name": "CSO [TBD]", "role": "Strategy & Eco", "desc": "Fortune 500 Exec\nPE/VC Experience"},
            ]
        },
        {
            "type": "content",
            "title": "Support Request",
            "bullets": [
                "Computing Power: 8M RMB/year - 74 GPU servers",
                "Office Space: 1.08M RMB - 500-800 sqm in digital economy park",
                "Hardware Equipment: 5.75M RMB - A100 servers, 50% subsidy",
            ]
        },
        {
            "type": "table",
            "title": "Economic Forecast",
            "headers": ["Year", "GMV", "Platform Revenue", "Industry Savings", "Net Profit"],
            "rows": [
                ["Year 1", "50M RMB", "2.5M RMB", "-", "-5M RMB"],
                ["Year 2", "300M RMB", "15M RMB", "1.5M RMB", "3M RMB"],
                ["Year 3", "1B RMB", "50M RMB", "6M RMB", "15M RMB"],
                ["Year 4", "5B RMB", "250M RMB", "30M RMB", "80M RMB"],
            ]
        },
        {
            "type": "content",
            "title": "Social Benefits",
            "bullets": [
                "Industry Transformation: Break OTA monopoly, fair transparent trading",
                "Technology Innovation: World's first hospitality AI transaction protocol",
                "Data Security: Merchant data self-controlled, break platform monopoly",
                "Employment: 100 direct jobs, 500 indirect jobs",
            ]
        },
        {
            "type": "content",
            "title": "Investment Budget - 80M RMB Total",
            "bullets": [
                "Computing Cost: 30M RMB (38%)",
                "R&D Investment: 20M RMB (25%)",
                "Hardware Equipment: 11.5M RMB (14%)",
                "Marketing: 8M RMB (10%)",
                "Operations + Ecosystem: 10.5M RMB (13%)",
            ]
        },
        {
            "type": "table",
            "title": "Risk Analysis",
            "headers": ["Risk Type", "Description", "Level", "Mitigation"],
            "rows": [
                ["Technical", "LLM effect below expectation", "Medium", "Multi-model fusion"],
                ["Market", "Low merchant adoption", "Medium", "Free trial, guarantee"],
                ["Competition", "OTA giants fight back", "High", "Differentiation, alliance"],
                ["Policy", "Stricter data regulation", "Low", "Compliance first"],
                ["Funding", "Insufficient financing", "Medium", "Control burn rate"],
            ]
        },
        {
            "type": "closing",
            "title": "Seeking Government Support",
            "contact": "Contact: Zhang Shi\nPhone/WeChat: 17760348653\nEmail: ericzhangshi@163.com",
        },
    ]
}


def generate_ahl_pptx(template="premium", output_dir="ppt_output"):
    os.makedirs(output_dir, exist_ok=True)
    
    generator = AHLPPTGENERATOR(template)
    content = AHL_CONTENT
    
    generator.add_title_slide(
        content["title"],
        content["subtitle"],
        content["presenter"]
    )
    
    for slide_data in content["slides"]:
        if slide_data["type"] == "content":
            generator.add_content_slide(
                slide_data["title"],
                slide_data["bullets"]
            )
        elif slide_data["type"] == "stats":
            generator.add_stats_slide(
                slide_data["title"],
                slide_data["stats"]
            )
        elif slide_data["type"] == "two_column":
            generator.add_two_column_slide(
                slide_data["title"],
                slide_data["left_title"],
                slide_data["left_bullets"],
                slide_data["right_title"],
                slide_data["right_bullets"]
            )
        elif slide_data["type"] == "timeline":
            generator.add_timeline_slide(
                slide_data["title"],
                slide_data["phases"]
            )
        elif slide_data["type"] == "team":
            generator.add_team_slide(
                slide_data["title"],
                slide_data["members"]
            )
        elif slide_data["type"] == "table":
            generator.add_table_slide(
                slide_data["title"],
                slide_data["headers"],
                slide_data["rows"]
            )
        elif slide_data["type"] == "closing":
            generator.add_closing_slide(
                slide_data["title"],
                slide_data.get("contact", "")
            )
    
    output_path = os.path.join(output_dir, "ahl-business-plan-{}.pptx".format(template))
    generator.save(output_path)
    
    return output_path


if __name__ == "__main__":
    template = sys.argv[1] if len(sys.argv) > 1 else "premium"
    output = generate_ahl_pptx(template)
    print("Generated: {}".format(output))
