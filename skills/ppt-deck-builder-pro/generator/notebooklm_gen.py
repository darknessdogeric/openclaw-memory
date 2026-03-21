# -*- coding: utf-8 -*-
"""
NotebookLM风格 PPT生成器 - Swiss/Bauhaus设计
现代报纸风格，极简，高对比度
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# NotebookLM/Swiss风格模板
TEMPLATE = {
    "primary": "#111111",    # 黑色
    "accent": "#FFCC00",     # 黄色强调
    "bg": "#FFFFFF",         # 白色背景
    "bg_gray": "#F5F5F5",   # 浅灰背景
}

def hex_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

class NotebookLMPPT:
    """NotebookLM/Swiss风格PPT生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def bg(self, slide, color="#FFFFFF"):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*hex_rgb(color))
    
    def title_text(self, slide, text, x, y, size, bold=True, color="#111111"):
        """大标题字体"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(12), Inches(2))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial"
        return tb
    
    def body_text(self, slide, text, x, y, size=14, color="#666666"):
        """正文字体"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(11), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial"
        return tb
    
    def label_text(self, slide, text, x, y, size=12, color="#999999"):
        """标签字体"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = text.upper()
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.font.name = "Arial"
        return tb
    
    def add_cover(self):
        """封面页"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        
        # 副标题标签
        self.label_text(s, "Project Blueprint", 0.5, 0.3, 10)
        
        # 主标题 - 超大
        self.title_text(s, "AHL", 0.5, 0.8, 72, True, "#111111")
        
        # 副标题
        self.body_text(s, "去中心化住宿业交易生态协议", 0.5, 2.5, 24, "#666666")
        
        # Punchline - 黄色大字号
        tb = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2))
        p = tb.text_frame.paragraphs[0]
        p.text = "范式革命"
        p.font.size = Pt(96)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
        p.font.name = "Arial Black"
        
        return s
    
    def add_section(self, num, title):
        """章节页"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        
        # 标签
        self.label_text(s, f"0{num} / {title.upper()}", 0.5, 0.3)
        
        # 标题
        self.title_text(s, title, 0.5, 0.8, 48, True, "#111111")
        
        return s
    
    def add_problem(self):
        """问题页 - 三大顽疾"""
        s = self.add_section("1", "Problem")
        s = self.prs.slides[-1]
        
        # 副标题
        self.title_text(s, "行业三大顽疾", 0.5, 1.8, 32, True, "#111111")
        
        # 4个统计卡片
        stats = [
            ("20-30%", "OTA垄断佣金"),
            ("¥1200亿", "年损失佣金"),
            ("<3%", "匹配转化率"),
            ("¥600亿", "营销浪费"),
        ]
        
        for i, (num, label) in enumerate(stats):
            x = 0.5 + i * 3.2
            # 数字
            tb = s.shapes.add_textbox(Inches(x), Inches(2.8), Inches(2.8), Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = num
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_rgb("#111111"))
            p.font.name = "Arial Black"
            # 标签
            self.body_text(s, label, x, 3.8, 12, "#666666")
        
        return s
    
    def add_solution(self):
        """解决方案对比页"""
        s = self.add_section("2", "Solution")
        s = self.prs.slides[-1]
        
        self.title_text(s, "解决方案", 0.5, 1.8, 32, True, "#111111")
        
        # 左侧 - OTA
        tb = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "❌ OTA模式"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb("#666666"))
        
        ota_points = [
            "中心化垄断，数据不透明",
            "佣金20-30%，成本高昂",
            "人找货，匹配效率低"
        ]
        for i, pt in enumerate(ota_points):
            self.body_text(s, f"× {pt}", 0.5, 3.1 + i * 0.5, 14, "#999999")
        
        # 右侧 - AHL
        tb = s.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "✓ AHL协议"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
        
        ahl_points = [
            "去中心化，直连交易",
            "效率费2-3%，成本革命",
            "AI双向匹配，95%+准确率"
        ]
        for i, pt in enumerate(ahl_points):
            self.body_text(s, f"✓ {pt}", 7, 3.1 + i * 0.5, 14, "#111111")
        
        return s
    
    def add_tech(self):
        """技术架构"""
        s = self.add_section("3", "Technology")
        
        self.title_text(s, "双AGENT技术架构", 0.5, 1.8, 32, True, "#111111")
        
        # 三个层次
        points = [
            ("01", "C端AI管家", "需求理解 · 偏好学习 · 智能推荐"),
            ("02", "向量匹配引擎", "语义检索 · 实时推荐 · 95%+准确率"),
            ("03", "B端AI运营官", "产品封装 · 动态定价 · 智能营销"),
        ]
        
        for i, (num, title, desc) in enumerate(points):
            y = 2.5 + i * 1.4
            # 数字
            tb = self.prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(y), Inches(0.8), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = num
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
            p.font.name = "Arial Black"
            
            # 标题
            self.title_text(self.prs.slides[-1], title, 1.4, y, 20, True, "#111111")
            # 描述
            self.body_text(self.prs.slides[-1], desc, 1.4, y + 0.4, 12, "#666666")
        
        return s
    
    def add_product(self):
        """产品矩阵"""
        s = self.add_section("4", "Product")
        
        self.title_text(s, "9大AGENT产品矩阵", 0.5, 1.8, 32, True, "#111111")
        
        # 9个AGENT
        agents = [
            ("运营", "客房/入住/客情"),
            ("营销", "获客/线索/RFP"),
            ("收益", "动态定价/预测"),
            ("宴会", "婚宴/MICE"),
            ("工程", "维保/能耗"),
            ("安防", "监控/预警"),
            ("财务", "对账/报表"),
            ("人资", "排班/招聘"),
            ("供应", "采购/库存"),
        ]
        
        for i, (name, desc) in enumerate(agents):
            row, col = i // 3, i % 3
            x = 0.5 + col * 4.3
            y = 2.5 + row * 1.5
            
            # 卡片背景
            sh = self.prs.slides[-1].shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4), Inches(1.2))
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor(*hex_rgb("#F5F5F5"))
            sh.line.color.rgb = RGBColor(*hex_rgb("#DDDDDD"))
            
            self.title_text(self.prs.slides[-1], name, x + 0.2, y + 0.2, 16, True, "#111111")
            self.body_text(self.prs.slides[-1], desc, x + 0.2, y + 0.7, 10, "#666666")
        
        # SKILL说明
        self.body_text(self.prs.slides[-1], "× 87个可插拔SKILL模块", 0.5, 6.8, 14, "#999999")
        
        return s
    
    def add_market(self):
        """市场机会"""
        s = self.add_section("5", "Market")
        
        self.title_text(s, "万亿级市场机会", 0.5, 1.8, 32, True, "#111111")
        
        stats = [
            ("¥6000亿", "住宿业市场规模"),
            ("¥4200亿", "在线预订(70%)"),
            ("¥1200亿", "OTA佣金"),
            ("¥1000亿+", "AHL可释放价值"),
        ]
        
        for i, (num, label) in enumerate(stats):
            x = 0.5 + i * 3.2
            # 数字 - 黄色
            tb = self.prs.slides[-1].shapes.add_textbox(Inches(x), Inches(2.8), Inches(2.8), Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = num
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
            p.font.name = "Arial Black"
            
            self.body_text(self.prs.slides[-1], label, x, 3.8, 12, "#666666")
        
        return s
    
    def add_business(self):
        """商业模式"""
        s = self.add_section("6", "Business")
        
        self.title_text(s, "商业模式", 0.5, 1.8, 32, True, "#111111")
        
        # 收入
        self.title_text(s, "收入来源", 0.5, 2.5, 18, True, "#111111")
        revenue = [("技术服务费 60%", "AI增值服务 25%", "数据服务 10%", "推广广告 5%")]
        for i, rev in enumerate(revenue):
            self.body_text(s, f"• {rev}", 0.5, 3.0 + i * 0.4, 14, "#333333")
        
        # 成本
        self.title_text(s, "成本构成", 7, 2.5, 18, True, "#111111")
        cost = [("服务器/云 40%", "AI模型调用 30%", "运维人员 20%", "其他 10%")]
        for i, c in enumerate(cost):
            self.body_text(s, f"• {c}", 7, 3.0 + i * 0.4, 14, "#333333")
        
        return s
    
    def add_roadmap(self):
        """发展规划"""
        s = self.add_section("7", "Roadmap")
        
        self.title_text(s, "四年发展规划", 0.5, 1.8, 32, True, "#111111")
        
        phases = [
            ("PHASE 1", "0-6月", "协议核心研发\nV1.0发布\n50家测试"),
            ("PHASE 2", "6-12月", "AGENT矩阵\n87个SKILL\n500家"),
            ("PHASE 3", "12-24月", "生态规模化\n5000家\n月GMV>5000万"),
            ("PHASE 4", "24-36月", "全球化\n海外试点\n行业标准"),
        ]
        
        for i, (phase, time, desc) in enumerate(phases):
            x = 0.5 + i * 3.2
            
            # 阶段
            tb = self.prs.slides[-1].shapes.add_textbox(Inches(x), Inches(2.5), Inches(2.8), Inches(0.4))
            p = tb.text_frame.paragraphs[0]
            p.text = phase
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
            
            # 时间
            self.body_text(self.prs.slides[-1], time, x, 2.9, 11, "#999999")
            
            # 描述
            self.body_text(self.prs.slides[-1], desc, x, 3.4, 11, "#333333")
        
        return s
    
    def add_team(self):
        """团队介绍"""
        s = self.add_section("8", "Team")
        
        self.title_text(s, "核心团队", 0.5, 1.8, 32, True, "#111111")
        
        members = [
            ("张实", "项目总控/发起人", "24年酒店业老兵\n战略/资源/政府关系"),
            ("CTO[待]", "AI技术架构", "华科AI博士\n10年+AI研发"),
            ("CSO[待]", "战略与生态", "世界500强高管\nPE/VC经验"),
        ]
        
        for i, (name, role, desc) in enumerate(members):
            x = 0.5 + i * 4.3
            
            # 头像占位
            sh = self.prs.slides[-1].shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 1.3), Inches(2.5), Inches(1.4), Inches(1.4))
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor(*hex_rgb("#F5F5F5"))
            
            # 姓名
            self.title_text(self.prs.slides[-1], name, x + 0.5, 4.1, 20, True, "#111111")
            
            # 角色 - 黄色
            tb = self.prs.slides[-1].shapes.add_textbox(Inches(x), Inches(4.5), Inches(3.8), Inches(0.4))
            p = tb.text_frame.paragraphs[0]
            p.text = role
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
            
            # 描述
            self.body_text(self.prs.slides[-1], desc, x, 5.0, 11, "#666666")
        
        return s
    
    def add_funding(self):
        """融资计划"""
        s = self.add_section("9", "Funding")
        
        self.title_text(s, "融资计划", 0.5, 1.8, 32, True, "#111111")
        
        # 融资信息框 - 黑色背景
        sh = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.5), Inches(6), Inches(3))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(*hex_rgb("#111111"))
        
        tb = s.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(5.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "种子轮"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
        
        tb = s.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(5.5), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = "¥500-800万"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb("#FFFFFF"))
        
        tb = s.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(5.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "出让股份 15-20% | 投后估值 3000-5000万"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(*hex_rgb("#999999"))
        
        # 资金用途
        self.title_text(s, "资金用途", 7, 2.5, 16, True, "#111111")
        self.body_text(s, "技术团队 40%", 7, 3.0, 13, "#333333")
        self.body_text(s, "数据体系 30%", 7, 3.4, 13, "#333333")
        self.body_text(s, "市场拓展 20%", 7, 3.8, 13, "#333333")
        self.body_text(s, "合规储备 10%", 7, 4.2, 13, "#333333")
        
        # 联系人
        self.title_text(s, "联系人", 7, 5.0, 16, True, "#111111")
        self.body_text(s, "张实 17760348653", 7, 5.5, 12, "#666666")
        self.body_text(s, "ericzhangshi@163.com", 7, 5.9, 12, "#666666")
        
        return s
    
    def add_closing(self):
        """结尾页"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        
        self.title_text(s, "加入AHL生态", 0.5, 2, 36, True, "#111111")
        
        self.body_text(s, "打破OTA垄断 · 建立公平交易 · 共享AI红利", 0.5, 3, 18, "#666666")
        
        # 共赢大字
        tb = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(2))
        p = tb.text_frame.paragraphs[0]
        p.text = "共赢"
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*hex_rgb("#FFCC00"))
        p.font.name = "Arial Black"
        
        return s
    
    def save(self, path):
        self.prs.save(path)
        return path


def generate_notebooklm_ppt(output="ppt_output/AHL-NotebookLM.pptx"):
    """生成NotebookLM风格PPT"""
    ppt = NotebookLMPPT()
    
    # 1. 封面
    ppt.add_cover()
    
    # 2. 问题
    ppt.add_problem()
    
    # 3. 解决方案
    ppt.add_solution()
    
    # 4. 技术
    ppt.add_tech()
    
    # 5. 产品
    ppt.add_product()
    
    # 6. 市场
    ppt.add_market()
    
    # 7. 商业
    ppt.add_business()
    
    # 8. 规划
    ppt.add_roadmap()
    
    # 9. 团队
    ppt.add_team()
    
    # 10. 融资
    ppt.add_funding()
    
    # 11. 结尾
    ppt.add_closing()
    
    return ppt.save(output)


if __name__ == "__main__":
    output = generate_notebooklm_ppt()
    print(f"Generated: {output}")
