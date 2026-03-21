# -*- coding: utf-8 -*-
"""
PPT智能生成器 - 组装模块
将大纲和配图组装成完整PPT
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 模板配置
TEMPLATES = {
    "premium": {
        "name": "高端奢华",
        "primary": "#1e3a5f", "accent": "#d4af37", "bg": "#0a0f1a"
    },
    "corporate": {
        "name": "企业商务", 
        "primary": "#1e40af", "accent": "#06b6d4", "bg": "#0a0a0f"
    },
    "startup": {
        "name": "创业路演",
        "primary": "#7c3aed", "accent": "#22d3ee", "bg": "#0f0f1a"
    },
    "tech": {
        "name": "科技蓝",
        "primary": "#0c4a6e", "accent": "#38bdf8", "bg": "#0c1929"
    },
    "minimal": {
        "name": "极简白",
        "primary": "#1f2937", "accent": "#10b981", "bg": "#ffffff"
    }
}

def hex_rgb(h):
    if not h or len(h) < 6:
        h = "000000"
    h = h.lstrip('#')
    if len(h) != 6:
        h = "000000"
    try:
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
    except:
        return (0, 0, 0)

class PPTAssembler:
    """PPT组装器"""
    
    def __init__(self, template="premium"):
        self.t = TEMPLATES.get(template, TEMPLATES["premium"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.images = []  # 配图信息
    
    def set_images(self, images):
        """设置配图"""
        self.images = images
    
    def bg(self, slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*hex_rgb(self.t["bg"]))
    
    def txt(self, slide, t, x, y, w, h, size=12, bold=False, color="#fff", align="center"):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        p = tb.text_frame.paragraphs[0]
        p.text = t
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*hex_rgb(color))
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        return tb
    
    def add_title_slide(self, title, subtitle=""):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        
        # 标题
        self.txt(s, title, 0.5, 2.5, 12, 1.5, 44, True, self.t["accent"])
        
        # 副标题
        if subtitle:
            self.txt(s, subtitle, 0.5, 4.2, 12, 1, 24, False, "#888888")
        
        return s
    
    def add_content_slide(self, title, bullets, image_info=None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        
        # 标题
        self.txt(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        
        # 内容
        for i, b in enumerate(bullets[:6]):
            self.txt(s, "• " + b, 0.7, 1.2 + i * 0.6, 11.5, 0.5, 16, False, "#ddd")
        
        # 配图占位（右侧）
        if image_info:
            # 创建图片占位框
            shape = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8), Inches(2), Inches(4.5), Inches(4)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(40, 40, 50)
            shape.line.color.rgb = RGBColor(*hex_rgb(self.t["accent"]))
            
            self.txt(s, f"[配图]\n{image_info.get('title', '')}", 
                    8, 3.5, 4.5, 2, 12, False, "#666")
        
        return s
    
    def add_image_slide(self, title, image_info):
        """带配图的幻灯片"""
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        
        self.txt(s, title, 0.5, 0.3, 12, 0.8, 32, True, self.t["accent"])
        
        # 配图区域
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(1.5), Inches(11.3), Inches(5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 30, 40)
        
        self.txt(s, f"[Image: {image_info.get('prompt', '')[:50]}...]",
                1, 3.5, 11.3, 1.5, 14, False, "#888")
        
        return s
    
    def add_closing_slide(self, title, contact=""):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.bg(s)
        
        self.txt(s, title, 0.5, 2.5, 12, 1.5, 44, True, self.t["accent"])
        
        if contact:
            self.txt(s, contact, 0.5, 4.5, 12, 1.5, 16, False, "#888")
        
        return s
    
    def assemble(self, outline, output_path):
        """组装PPT"""
        # 封面
        self.add_title_slide(
            outline.get("title", "演示文稿"),
            outline.get("subtitle", "")
        )
        
        # 内容页
        for section in outline.get("sections", []):
            title = section.get("title", "")
            content = section.get("content", "")
            
            # 获取对应配图
            img_info = None
            for img in self.images:
                if img.get("slide_num") == section.get("num"):
                    img_info = img
                    break
            
            # 生成要点
            bullets = [
                content,
                f"关键点: {', '.join(outline.get('keywords', [])[:3])}"
            ]
            
            self.add_content_slide(title, bullets, img_info)
        
        # 结尾页
        self.add_closing_slide("感谢关注", "联系方式：请替换")
        
        # 保存
        self.prs.save(output_path)
        return output_path


def assemble_ppt(outline, images, template="premium", output="output.pptx"):
    """组装PPT的入口函数"""
    assembler = PPTAssembler(template)
    assembler.set_images(images)
    return assembler.assemble(outline, output)


if __name__ == "__main__":
    # 测试
    test_outline = {
        "title": "测试演示",
        "sections": [
            {"num": 1, "title": "背景", "content": "项目背景介绍"},
            {"num": 2, "title": "方案", "content": "解决方案说明"},
            {"num": 3, "title": "团队", "content": "核心团队介绍"}
        ],
        "keywords": ["AI", "技术", "创新"]
    }
    
    test_images = [
        {"slide_num": 1, "title": "背景图", "prompt": "business office"},
        {"slide_num": 2, "title": "方案图", "prompt": "tech innovation"},
        {"slide_num": 3, "title": "团队图", "prompt": "team meeting"}
    ]
    
    output = assemble_ppt(test_outline, test_images, "premium", "test_output.pptx")
    print(f"Generated: {output}")
