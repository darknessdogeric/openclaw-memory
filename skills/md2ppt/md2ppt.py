#!/usr/bin/env python3
"""
MD2PPT Converter - Markdown转PPT工具
将AHL项目文档转换为路演演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import re
import markdown

class MD2PPTConverter:
    """Markdown转PPT转换器"""
    
    def __init__(self):
        self.prs = Presentation()
        # 设置幻灯片尺寸为16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # 定义AHL品牌色
        self.colors = {
            'primary': RGBColor(26, 84, 144),      # AHL蓝
            'secondary': RGBColor(44, 90, 160),    # 浅蓝
            'accent': RGBColor(255, 127, 39),      # 橙色强调
            'dark': RGBColor(51, 51, 51),          # 深灰
            'light': RGBColor(240, 240, 240),      # 浅灰
            'white': RGBColor(255, 255, 255),      # 白色
        }
    
    def add_title_slide(self, title, subtitle=""):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景色块
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.CENTER
        
        # 底部装饰
        shape = slide.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['accent']
        shape.line.fill.background()
        
        return slide
    
    def add_content_slide(self, title, content_lines, layout_type="bullet"):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # 页面标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # 内容区域
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(content_lines[:8]):  # 每页最多8条
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 清理markdown标记
            clean_line = re.sub(r'\*\*|\*|__|_|`', '', line)
            clean_line = re.sub(r'^[-*+]\s*', '• ', clean_line)
            clean_line = re.sub(r'^\d+\.\s*', f'{i+1}. ', clean_line)
            
            p.text = clean_line
            p.font.size = Pt(18) if len(content_lines) > 5 else Pt(20)
            p.font.color.rgb = self.colors['dark']
            p.space_after = Pt(12)
            p.level = 0 if not line.startswith('    ') else 1
        
        return slide
    
    def add_two_column_slide(self, title, left_content, right_content):
        """添加双栏对比页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 顶部色条
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # 页面标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # 左栏
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(5.8))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(left_content[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark']
        
        # 右栏
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(6), Inches(5.8))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(right_content[:6]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark']
        
        return slide
    
    def add_section_divider(self, section_title):
        """添加章节分隔页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 全页背景
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_end_slide(self, contact_info):
        """添加结束页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
        
        # 感谢标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢聆听"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 联系方式
        contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(2))
        tf = contact_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(contact_info):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)
        
        return slide
    
    def convert_markdown(self, md_file, output_file, presentation_title=""):
        """转换Markdown文件为PPT"""
        # 读取Markdown
        with open(md_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # 解析标题
        title_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
        main_title = title_match.group(1) if title_match else presentation_title
        
        # 解析各部分
        sections = re.split(r'^##\s+', content, flags=re.MULTILINE)
        
        # 添加标题页
        self.add_title_slide(main_title, "去中心化住宿业交易生态协议")
        
        # 处理每个章节
        for section in sections[1:]:  # 跳过第一个空部分
            lines = section.strip().split('\n')
            section_title = lines[0].strip()
            
            # 添加章节分隔页
            self.add_section_divider(section_title)
            
            # 收集内容
            content_lines = []
            for line in lines[1:]:
                line = line.strip()
                if line and not line.startswith('---'):
                    content_lines.append(line)
            
            # 分批添加内容页
            for i in range(0, len(content_lines), 6):
                batch = content_lines[i:i+6]
                self.add_content_slide(section_title, batch)
        
        # 添加结束页
        self.add_end_slide([
            "AHL智能科技有限公司",
            "联系人：张实",
            "电话/微信：17760348653",
            "邮箱：ericzhangshi@163.com",
            "让住宿交易更简单、更公平、更高效"
        ])
        
        # 保存
        self.prs.save(output_file)
        print(f"[OK] PPT created: {output_file}")
        return output_file


def main():
    """主函数 - 批量转换"""
    import os
    
    source_dir = r'C:\Users\Administrator\Desktop\张实项目总控\06-AHL-去中心化旅行平台'
    output_dir = r'C:\Users\Administrator\Desktop\项目说明书'
    
    os.makedirs(output_dir, exist_ok=True)
    
    # 要转换的文件
    files_to_convert = [
        ('01-政府申请-项目说明书-去中心化协议版.md', 'AHL路演-政府申请.pptx'),
        ('02-商业计划书-投资人版-V2.0.md', 'AHL路演-投资人BP.pptx'),
        ('AHL顶层设计总纲-V2.0.md', 'AHL路演-顶层设计.pptx'),
    ]
    
    print("="*60)
    print("MD2PPT Converter - Markdown转PPT工具")
    print("="*60)
    
    success_count = 0
    for md_name, pptx_name in files_to_convert:
        md_path = os.path.join(source_dir, md_name)
        pptx_path = os.path.join(output_dir, pptx_name)
        
        if os.path.exists(md_path):
            try:
                converter = MD2PPTConverter()
                converter.convert_markdown(md_path, pptx_path)
                success_count += 1
            except Exception as e:
                print(f"[ERROR] {md_name}: {e}")
        else:
            print(f"[SKIP] Not found: {md_name}")
    
    print("="*60)
    print(f"Conversion completed: {success_count}/{len(files_to_convert)} files")
    print(f"Output directory: {output_dir}")
    print("="*60)


if __name__ == '__main__':
    main()
