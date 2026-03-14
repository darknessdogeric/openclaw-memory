#!/usr/bin/env python3
"""
AHL项目说明书PPT - 通俗易懂版
面向所有人，一看就懂
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

class AHLSimplePPT:
    """AHL通俗版PPT生成器"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # 配色
        self.colors = {
            'blue': RGBColor(26, 84, 144),
            'light_blue': RGBColor(66, 133, 244),
            'orange': RGBColor(255, 127, 39),
            'green': RGBColor(52, 168, 83),
            'dark': RGBColor(51, 51, 51),
            'white': RGBColor(255, 255, 255),
            'light_gray': RGBColor(248, 249, 250),
        }
    
    def add_big_title_slide(self, title, subtitle, highlight_text=""):
        """大标题页 - 震撼开场"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 全页背景
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['blue']
        bg.line.fill.background()
        
        # 主标题 - 超大字
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(32)
            p.font.color.rgb = self.colors['orange']
            p.alignment = PP_ALIGN.CENTER
        
        # 高亮文字
        if highlight_text:
            high_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.8))
            tf = high_box.text_frame
            p = tf.paragraphs[0]
            p.text = highlight_text
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_concept_slide(self, title, concept, explanation, analogy=""):
        """概念解释页 - 一句话说明白"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 浅色背景
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['light_gray']
        bg.line.fill.background()
        
        # 顶部色条
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['blue']
        bar.line.fill.background()
        
        # 页面标题
        t = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.5))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # 核心概念 - 大盒子
        concept_box = slide.shapes.add_shape(1, Inches(1), Inches(1.2), Inches(11.333), Inches(1.8))
        concept_box.fill.solid()
        concept_box.fill.fore_color.rgb = self.colors['light_blue']
        concept_box.line.fill.background()
        
        # 概念文字
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(10.333), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = concept
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 解释文字
        exp_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1.5))
        tf = exp_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = explanation
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['dark']
        p.alignment = PP_ALIGN.CENTER
        
        # 类比
        if analogy:
            ana_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1.5))
            tf = ana_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"💡 打个比方：{analogy}"
            p.font.size = Pt(22)
            p.font.color.rgb = self.colors['orange']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_problem_solution_slide(self, problem_title, problems, solution_title, solutions):
        """问题-方案对比页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 左半部分 - 问题（红色系）
        left_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(6.666), Inches(7.5))
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(255, 240, 240)
        left_bg.line.fill.background()
        
        # 左标题
        lt = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(6), Inches(0.8))
        tf = lt.text_frame
        p = tf.paragraphs[0]
        p.text = f"❌ {problem_title}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(200, 50, 50)
        
        # 左内容
        lc = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.666), Inches(5.5))
        tf = lc.text_frame
        tf.word_wrap = True
        for i, prob in enumerate(problems):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = prob
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['dark']
            p.space_after = Pt(16)
        
        # 右半部分 - 方案（绿色系）
        right_bg = slide.shapes.add_shape(1, Inches(6.666), Inches(0), Inches(6.667), Inches(7.5))
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = RGBColor(240, 255, 240)
        right_bg.line.fill.background()
        
        # 右标题
        rt = slide.shapes.add_textbox(Inches(7), Inches(0.3), Inches(6), Inches(0.8))
        tf = rt.text_frame
        p = tf.paragraphs[0]
        p.text = f"✅ {solution_title}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['green']
        
        # 右内容
        rc = slide.shapes.add_textbox(Inches(7.166), Inches(1.2), Inches(5.667), Inches(5.5))
        tf = rc.text_frame
        tf.word_wrap = True
        for i, sol in enumerate(solutions):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = sol
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['dark']
            p.space_after = Pt(16)
        
        return slide
    
    def add_how_it_works_slide(self, steps):
        """工作原理页 - 流程图"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题栏
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['blue']
        bar.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "AHL是怎么工作的？"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # 流程步骤 - 横向排列
        step_width = 3.5
        start_x = 0.5
        y_pos = 1.5
        
        colors = [self.colors['light_blue'], self.colors['orange'], self.colors['green'], self.colors['blue']]
        
        for i, (step_title, step_desc) in enumerate(steps):
            x_pos = start_x + i * (step_width + 0.5)
            
            # 步骤盒子
            box = slide.shapes.add_shape(1, Inches(x_pos), Inches(y_pos), Inches(step_width), Inches(4))
            box.fill.solid()
            box.fill.fore_color.rgb = colors[i % len(colors)]
            box.line.fill.background()
            
            # 步骤编号
            num = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(y_pos + 0.2), Inches(0.8), Inches(0.8))
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            # 步骤标题
            st = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(y_pos + 1.2), Inches(step_width - 0.6), Inches(0.8))
            tf = st.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            
            # 步骤描述
            sd = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(y_pos + 2.2), Inches(step_width - 0.6), Inches(1.5))
            tf = sd.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step_desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['white']
            
            # 箭头（除了最后一个）
            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(13, Inches(x_pos + step_width + 0.1), Inches(y_pos + 1.5), Inches(0.3), Inches(1))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['dark']
                arrow.line.fill.background()
        
        return slide
    
    def add_value_slide(self, values):
        """价值页 - 三方共赢"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题栏
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['blue']
        bar.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "AHL带来什么价值？"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # 三个价值卡片
        card_width = 4
        gap = 0.4
        start_x = 0.5
        y_pos = 1.2
        
        colors = [self.colors['light_blue'], self.colors['orange'], self.colors['green']]
        
        for i, (who, benefits) in enumerate(values):
            x_pos = start_x + i * (card_width + gap)
            
            # 卡片背景
            card = slide.shapes.add_shape(1, Inches(x_pos), Inches(y_pos), Inches(card_width), Inches(5.5))
            card.fill.solid()
            card.fill.fore_color.rgb = colors[i]
            card.line.fill.background()
            
            # 角色名称
            who_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.3), Inches(card_width), Inches(0.8))
            tf = who_box.text_frame
            p = tf.paragraphs[0]
            p.text = who
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            
            # 好处列表
            ben_box = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(y_pos + 1.3), Inches(card_width - 0.6), Inches(4))
            tf = ben_box.text_frame
            tf.word_wrap = True
            for j, benefit in enumerate(benefits):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"✓ {benefit}"
                p.font.size = Pt(18)
                p.font.color.rgb = self.colors['white']
                p.space_after = Pt(12)
        
        return slide
    
    def add_data_slide(self, title, data_items):
        """数据页 - 大数字冲击"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题栏
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['blue']
        bar.line.fill.background()
        
        t = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # 数据展示 - 网格布局
        positions = [(0.5, 1.2), (4.5, 1.2), (8.5, 1.2), (0.5, 4.2), (4.5, 4.2), (8.5, 4.2)]
        
        for i, (number, label) in enumerate(data_items[:6]):
            x, y = positions[i]
            
            # 数字框
            num_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(1.5))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = number
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = self.colors['orange']
            p.alignment = PP_ALIGN.CENTER
            
            # 标签
            label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.3), Inches(4), Inches(0.8))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def create_presentation(self, output_file):
        """创建完整的通俗版PPT"""
        
        # 第1页：封面
        self.add_big_title_slide(
            "AHL",
            "去中心化住宿业交易生态协议",
            "让消费者和酒店直接交易，没有中间商赚差价"
        )
        
        # 第2页：一句话说明白
        self.add_concept_slide(
            "AHL是什么？",
            "住宿业的「智能红娘」",
            "用人工智能帮消费者找到最合适的酒店/民宿，\n让酒店直接接待客人，不需要通过携程美团等平台",
            "就像滴滴连接司机和乘客，AHL连接消费者和酒店，\n但AHL只收5%技术服务费，而不是30%的佣金"
        )
        
        # 第3页：解决什么问题
        self.add_problem_solution_slide(
            "现在的痛点",
            [
                "• 携程美团收20-30%高额佣金",
                "• 酒店利润被大幅压缩",
                "• 消费者订不到真正想要的",
                "• 酒店客户数据被平台垄断"
            ],
            "AHL的方案",
            [
                "• 只收5%技术服务费",
                "• 酒店收入提升35%+",
                "• AI精准匹配，满意率>95%",
                "• 数据归酒店，建立私域资产"
            ]
        )
        
        # 第4页：怎么工作
        self.add_how_it_works_slide([
            ("消费者说需求", "\"我想住大理古城带院子的民宿，预算500\""),
            ("AI管家理解", "提取需求：院子、宠物友好、¥500"),
            ("智能匹配", "从10万+酒店中找到最合适的10家"),
            ("直连交易", "消费者直接付钱给酒店，没有中间商")
        ])
        
        # 第5页：给谁用
        self.add_value_slide([
            ("对消费者", [
                "价格更便宜15-25%",
                "说人话就能找住宿",
                "AI推荐比搜索更准",
                "24小时智能客服"
            ]),
            ("对酒店/民宿", [
                "省下30%佣金",
                "获得精准客户",
                "数据自己掌控",
                "AI自动运营"
            ]),
            ("对行业", [
                "打破OTA垄断",
                "年省1000亿成本",
                "提升交易效率",
                "建立公平生态"
            ])
        ])
        
        # 第6页：核心数据
        self.add_data_slide(
            "AHL的市场潜力",
            [
                ("¥6,000亿", "中国住宿业市场规模"),
                ("¥1,200亿", "每年OTA佣金成本"),
                ("5%", "AHL技术服务费"),
                (">95%", "AI匹配准确率"),
                ("10万+", "目标接入商家"),
                ("¥1,000亿", "可为行业节省")
            ]
        )
        
        # 第7页：技术架构（简化版）
        self.add_concept_slide(
            "AHL的技术核心",
            "双AI管家 + 智能匹配",
            "消费者有个「AI管家」理解需求，\n酒店有个「AI运营官」提供服务，\n中间有个「智能匹配引擎」精准对接",
            "就像买房有中介撮合，但AHL的「中介」是AI，\n不收高额佣金，只收合理服务费"
        )
        
        # 第8页：商业模式
        self.add_problem_solution_slide(
            "传统模式",
            [
                "客户付¥500",
                "平台抽走¥150 (30%)",
                "酒店只收到¥350",
                "客户多花钱，酒店少赚钱"
            ],
            "AHL模式",
            [
                "客户付¥500",
                "AHL收¥25 (5%)",
                "酒店收到¥475",
                "三方共赢，皆大欢喜"
            ]
        )
        
        # 第9页：团队
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['light_gray']
        bg.line.fill.background()
        
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors['blue']
        bar.line.fill.background()
        
        t = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = "我们的团队"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # 团队成员
        team = [
            ("张实", "项目发起人", "24年酒店业经验\n曾任知名酒店高管"),
            ("CTO", "技术负责人", "AI+区块链专家\n10年技术架构经验"),
            ("产品负责人", "双边平台专家", "5年B端产品设计\n精通交易生态"),
        ]
        
        for i, (name, role, desc) in enumerate(team):
            x = 0.5 + i * 4.3
            
            # 头像占位
            avatar = slide.shapes.add_shape(9, Inches(x), Inches(1.5), Inches(1.5), Inches(1.5))
            avatar.fill.solid()
            avatar.fill.fore_color.rgb = self.colors['light_blue']
            avatar.line.fill.background()
            
            # 姓名
            n = slide.shapes.add_textbox(Inches(x), Inches(3.2), Inches(3.8), Inches(0.6))
            tf = n.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.colors['blue']
            
            # 职位
            r = slide.shapes.add_textbox(Inches(x), Inches(3.8), Inches(3.8), Inches(0.5))
            tf = r.text_frame
            p = tf.paragraphs[0]
            p.text = role
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['orange']
            
            # 简介
            d = slide.shapes.add_textbox(Inches(x), Inches(4.4), Inches(3.8), Inches(1.5))
            tf = d.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['dark']
        
        # 第10页：结束
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['blue']
        bg.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "让住宿交易更简单、更公平、更高效"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = "AHL - 去中心化住宿业交易生态协议"
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['orange']
        p.alignment = PP_ALIGN.CENTER
        
        contact = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.333), Inches(1.5))
        tf = contact.text_frame
        tf.word_wrap = True
        lines = [
            "联系人：张实",
            "电话/微信：17760348653",
            "邮箱：ericzhangshi@163.com",
            "我们正在寻找投资人和合作伙伴"
        ]
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['white']
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)
        
        # 保存
        self.prs.save(output_file)
        print(f"[OK] Created: {output_file}")
        print(f"      Total slides: {len(self.prs.slides)}")
        return output_file


def main():
    """生成通俗版PPT"""
    output_dir = r'C:\Users\Administrator\Desktop\项目说明书'
    os.makedirs(output_dir, exist_ok=True)
    
    output_file = os.path.join(output_dir, 'AHL项目说明书-通俗版.pptx')
    
    print("="*60)
    print("Creating AHL Simple Presentation")
    print("="*60)
    
    ppt = AHLSimplePPT()
    ppt.create_presentation(output_file)
    
    print("="*60)
    print(f"Output: {output_file}")
    print("="*60)


if __name__ == '__main__':
    main()
