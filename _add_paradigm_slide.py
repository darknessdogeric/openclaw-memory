# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

bp_path = r'C:\Users\ericz\Desktop\张实项目总控\05-AHL-去中心化旅行平台\项目说明书\AHL-LLM去中心化旅行平台商业计划书V5.1(9).pptx'
prs = Presentation(bp_path)

# 颜色定义
C_DARK  = RGBColor(0x06, 0x16, 0x3A)   # 深蓝黑
C_BLUE  = RGBColor(0x00, 0x5A, 0xB5)   # AHL蓝
C_LIGHT = RGBColor(0x00, 0xA8, 0xFF)   # 亮蓝
C_GOLD  = RGBColor(0xFF, 0xD7, 0x00)   # 金色强调
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x88, 0x99, 0xAA)
C_TEXT  = RGBColor(0x22, 0x33, 0x44)
C_RED   = RGBColor(0xFF, 0x44, 0x44)

# 用 DEFAULT 布局新建幻灯片
layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)

# === 背景 ===
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = C_WHITE

# === 左侧深色区域 ===
left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.0), Inches(7.5))
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = C_DARK
left_bar.line.fill.background()

# === 左上角红色强调点 ===
accent_dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.3), Inches(0.12), Inches(0.12))
accent_dot.fill.solid()
accent_dot.fill.fore_color.rgb = C_RED
accent_dot.line.fill.background()

# === 左上角章节编号 ===
num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(1.5), Inches(0.5))
tf = num_box.text_frame
p = tf.paragraphs[0]
p.text = '第三章'
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = C_LIGHT

# === 左上角主标题 ===
title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(3.5), Inches(1.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '范式革命'
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = C_WHITE

p2 = tf.add_paragraph()
p2.text = '垂直大模型交互的必然性'
p2.font.size = Pt(13)
p2.font.color.rgb = C_LIGHT

# === 左下角核心论点 ===
thesis_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.8), Inches(3.5), Inches(1.2))
tf = thesis_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"中心化搜索"让位于"自然语言理解"'
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = C_GOLD

p2 = tf.add_paragraph()
p2.text = '这不是预测，而是历史的必然。'
p2.font.size = Pt(10)
p2.font.color.rgb = C_GRAY
p2.space_before = Pt(6)

# === 时间线（左侧区域底部） ===
timeline_y = 3.2
# 时间线横线
tl_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(timeline_y + 0.35), Inches(3.4), Pt(2))
tl_line.fill.solid()
tl_line.fill.fore_color.rgb = C_BLUE
tl_line.line.fill.background()

# 时间节点数据
eras = [
    ('2000', '搜索引擎\nYahoo/Google'),
    ('2010', '超级APP\n美团/携程'),
    ('2020', 'LLM\nChatGPT'),
    ('2025', '垂直大模型\n行业落地'),
]
for i, (year, desc) in enumerate(eras):
    cx = 0.4 + i * 0.88
    # 圆点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(timeline_y + 0.22), Inches(0.26), Inches(0.26))
    dot.fill.solid()
    dot.fill.fore_color.rgb = C_LIGHT if i < 3 else C_GOLD
    dot.line.fill.background()
    # 年份
    yr_box = slide.shapes.add_textbox(Inches(cx - 0.1), Inches(timeline_y - 0.35), Inches(0.6), Inches(0.3))
    tf = yr_box.text_frame
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    # 描述
    desc_box = slide.shapes.add_textbox(Inches(cx - 0.2), Inches(timeline_y + 0.55), Inches(0.9), Inches(0.7))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(7.5)
    p.font.color.rgb = C_GRAY
    p.alignment = PP_ALIGN.CENTER

# === 底部结论 ===
bot_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(3.4), Inches(0.5))
tf = bot_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '范式迭代周期：15年 → 10年 → 5年 → ?'
p.font.size = Pt(9)
p.font.color.rgb = C_LIGHT

p2 = tf.add_paragraph()
p2.text = '变革窗口正在压缩，机会稍纵即逝。'
p2.font.size = Pt(8.5)
p2.font.bold = True
p2.font.color.rgb = C_GOLD

# === 右侧五大维度内容区 ===
# 顶部蓝色横条
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(0), Inches(9.33), Inches(0.06))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = C_BLUE
top_bar.line.fill.background()

# 右侧标题
rt_title = slide.shapes.add_textbox(Inches(4.3), Inches(0.25), Inches(8.5), Inches(0.5))
tf = rt_title.text_frame
p = tf.paragraphs[0]
p.text = '知其然 · 知其所以然'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_TEXT

rt_sub = slide.shapes.add_textbox(Inches(4.3), Inches(0.7), Inches(8.5), Inches(0.35))
tf = rt_sub.text_frame
p = tf.paragraphs[0]
p.text = '理解范式转移的五个维度，看清为什么AHL代表必然'
p.font.size = Pt(9.5)
p.font.color.rgb = C_GRAY

# === 五大维度 ===
dims = [
    {
        'num': '01',
        'title': '历史沿革',
        'en': 'HISTORICAL EVOLUTION',
        'color': C_BLUE,
        'points': [
            '第一代：人工目录（Yahoo1994）—— 信息找人',
            '第二代：搜索比价（Google/携程2003）—— 人找信息',
            '第三代：超级APP（美团/抖音2010）—— 中心化入口',
            '第四代：垂直大模型（2024+）—— 自然语言+个性化',
        ],
        'insight': '每代范式存活时间：15年 → 10年 → 5年 → 正在压缩',
    },
    {
        'num': '02',
        'title': '当下情况',
        'en': 'CURRENT STATE',
        'color': RGBColor(0x00, 0x7A, 0xD0),
        'points': [
            '携程模式：中心化搜索 + 人工客服 + 15%佣金',
            '用户痛点：信息过载、多平台比价、等待客服',
            '酒店困境：数据在平台手中，无法直连消费者',
            '已有信号：87%用户希望"说话"而非"搜索"（NLP报告2024）',
        ],
        'insight': '旧范式仍存在，但成本结构已不可持续',
    },
    {
        'num': '03',
        'title': '原因分析',
        'en': 'ROOT CAUSE',
        'color': RGBColor(0x00, 0x96, 0xE0),
        'points': [
            'LLM推理成本：从¥0.5/千token降至¥0.0005（千倍降本）',
            '中文NLP准确率：头部模型达95%+，超越人工客服均值',
            '用户行为：微信对话日均70+条，远超搜索框使用频率',
            '硬件成本：GPU算力成本年降40%，边际成本趋零',
        ],
        'insight': '技术临界点已过，替代加速',
    },
    {
        'num': '04',
        'title': '未来预判',
        'en': 'FUTURE TREND',
        'color': RGBColor(0x00, 0xA8, 0xFF),
        'points': [
            '2025-2027：垂直LLM渗透率从3%→25%（酒店行业）',
            '2027-2030：自然语言交互成为主流预订方式（超50%）',
            '2030+：中心化OTA平台转型为底层数据提供商',
            '替代时间表：先替代客服（2025）→ 再替代搜索（2027）→ 最后替代交易（2029）',
        ],
        'insight': '不是"如果"而是"何时"——已在进行中',
    },
    {
        'num': '05',
        'title': '底层逻辑',
        'en': 'UNDERLYING LOGIC',
        'color': C_GOLD,
        'points': [
            '信息流演进："人找信息" → "信息找人" → "AI替你决策"',
            '交互革命："搜索框" → "对话框" → "自然语言成交"',
            '信任转移："平台背书" → "AI理解个体偏好" → "社区信任"',
            '商业本质：最低交易摩擦的形态必然胜出——自然语言 < 搜索框 < 目录',
        ],
        'insight': '自然语言 = 最低摩擦 → 必然最终形态',
    },
]

dim_w = 2.9
dim_h = 2.1
dim_y_start = 1.1
dim_gap = 0.13
dim_start_x = 4.15

for i, dim in enumerate(dims):
    dx = dim_start_x + i * (dim_w + dim_gap)
    dy = dim_y_start

    # 卡片背景
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(dx), Inches(dy), Inches(dim_w), Inches(dim_h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF9, 0xFF)
    card.line.color.rgb = dim['color']
    card.line.width = Pt(1.2)

    # 顶部色条
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(dx), Inches(dy), Inches(dim_w), Inches(0.05))
    strip.fill.solid()
    strip.fill.fore_color.rgb = dim['color']
    strip.line.fill.background()

    # 编号
    num_box = slide.shapes.add_textbox(Inches(dx + 0.08), Inches(dy + 0.1), Inches(0.4), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = dim['num']
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = dim['color']

    # 标题
    title_box = slide.shapes.add_textbox(Inches(dx + 0.08), Inches(dy + 0.35), Inches(dim_w - 0.16), Inches(0.28))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = dim['title']
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = C_TEXT

    # 英文副标题
    en_box = slide.shapes.add_textbox(Inches(dx + 0.08), Inches(dy + 0.58), Inches(dim_w - 0.16), Inches(0.2))
    tf = en_box.text_frame
    p = tf.paragraphs[0]
    p.text = dim['en']
    p.font.size = Pt(6.5)
    p.font.color.rgb = C_GRAY

    # 要点
    for j, pt in enumerate(dim['points']):
        py = dy + 0.82 + j * 0.26
        pt_box = slide.shapes.add_textbox(Inches(dx + 0.1), Inches(py), Inches(dim_w - 0.2), Inches(0.26))
        tf = pt_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = '· ' + pt
        p.font.size = Pt(7)
        p.font.color.rgb = RGBColor(0x44, 0x55, 0x66)

    # 底部洞察（金字）
    ins_y = dy + dim_h - 0.38
    ins_box = slide.shapes.add_textbox(Inches(dx + 0.08), Inches(ins_y), Inches(dim_w - 0.16), Inches(0.3))
    tf = ins_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = dim['insight']
    p.font.size = Pt(7)
    p.font.color.rgb = dim['color']
    p.font.bold = True

# === 右侧底部总结 ===
sum_y = 3.35

sum_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.15), Inches(sum_y), Inches(9.18), Inches(1.55))
sum_bar.fill.solid()
sum_bar.fill.fore_color.rgb = RGBColor(0x0A, 0x1E, 0x3A)
sum_bar.line.fill.background()

# 三个关键结论
conclusions = [
    ('历史必然', '垂直大模型交互不是"新选项"，而是信息革命的下一站', C_LIGHT),
    ('指数压缩', '范式进步从15年压缩至5年，变革窗口正在关闭', C_GOLD),
    ('就在眼前', '2025年AI预订已完成验证，替代已经开始', RGBColor(0xFF, 0x88, 0x44)),
]

for i, (title, desc, color) in enumerate(conclusions):
    cx = 4.3 + i * 3.0

    # 竖线
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.15), Inches(sum_y + 0.15), Pt(2), Inches(1.25))
    vline.fill.solid()
    vline.fill.fore_color.rgb = color
    vline.line.fill.background()

    # 标题
    t_box = slide.shapes.add_textbox(Inches(cx), Inches(sum_y + 0.15), Inches(2.8), Inches(0.35))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color

    # 描述
    d_box = slide.shapes.add_textbox(Inches(cx), Inches(sum_y + 0.5), Inches(2.8), Inches(0.9))
    tf = d_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0xCC, 0xD9, 0xEC)

# === 右下角CTA ===
cta_box = slide.shapes.add_textbox(Inches(4.3), Inches(5.05), Inches(9.0), Inches(0.4))
tf = cta_box.text_frame
p = tf.paragraphs[0]
p.text = '携程的中心化模式 —— 是工业时代的巅峰，而非永恒。'
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0xFF, 0x66, 0x66)
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = 'AHL代表的是下一个时代的交互形态。先行者，已经在路上了。'
p2.font.size = Pt(9.5)
p2.font.color.rgb = RGBColor(0x55, 0x66, 0x77)

# === 底部时间轴 ===
bot_timeline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.15), Inches(5.8), Inches(9.18), Pt(1))
bot_timeline.fill.solid()
bot_timeline.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
bot_timeline.line.fill.background()

labels = [
    ('2024 LLM元年', 'ChatGPT重塑一切'),
    ('2025 客服替代', 'AI接管预订咨询'),
    ('2027 搜索替代', '自然语言成为入口'),
    ('2029 交易替代', '去中心化交易主流'),
    ('2031+ 新秩序', '平台转型为基础设施'),
]
for i, (yr, desc) in enumerate(labels):
    lx = 4.25 + i * 1.85
    # 年份标签
    yr_box = slide.shapes.add_textbox(Inches(lx), Inches(5.85), Inches(1.7), Inches(0.28))
    tf = yr_box.text_frame
    p = tf.paragraphs[0]
    p.text = yr
    p.font.size = Pt(7.5)
    p.font.bold = True
    p.font.color.rgb = C_BLUE
    # 描述
    desc_box = slide.shapes.add_textbox(Inches(lx), Inches(6.1), Inches(1.7), Inches(0.35))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(7)
    p.font.color.rgb = C_GRAY
    # 连接线
    if i < 4:
        conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(lx + 1.72), Inches(5.97), Inches(0.12), Pt(1))
        conn.fill.solid()
        conn.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        conn.line.fill.background()

# === 底部标语 ===
slogan = slide.shapes.add_textbox(Inches(4.3), Inches(6.55), Inches(9.0), Inches(0.4))
tf = slogan.text_frame
p = tf.paragraphs[0]
p.text = '范式革命的窗口期：5-7年。现在进入，才能定义规则。'
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0x44, 0x44)
p.alignment = PP_ALIGN.CENTER

# === 调整幻灯片顺序：插到第4位（index=3，在slide3市场机会之后）===
sldIdLst = prs.slides._sldIdLst
sldIds = list(sldIdLst)
# 新幻灯片在最后，需要移到第4位（index=3）
new_slide = sldIds[-1]
sldIdLst.remove(new_slide)
sldIdLst.insert(3, new_slide)

prs.save(bp_path)
print('Done. Total slides:', len(prs.slides))
