"""V3.0 - 2026 May Day Hotel Market Report PPT with proper STR classification, city tiers, and regional detail."""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
import os

# === Colors ===
DARK  = RGBColor(0x0D,0x21,0x3E)
BLUE  = RGBColor(0x1B,0x4F,0x8A)
LBLUE = RGBColor(0x3A,0x7C,0xBF)
ORANGE= RGBColor(0xE8,0x6A,0x17)
GREEN = RGBColor(0x1B,0x8A,0x4A)
RED   = RGBColor(0xC4,0x39,0x39)
GRAY  = RGBColor(0x44,0x44,0x44)
LGRAY = RGBColor(0x99,0x99,0x99)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BG    = RGBColor(0xF5,0xF7,0xFA)
HL    = RGBColor(0xFF,0xF0,0xE0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===================== HELPERS =====================
def bg(s, c=BG): s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def title_bar(s, t, sub=None):
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.12), Inches(11.8), Inches(0.55))
    p = tb.text_frame.paragraphs[0]; p.text = t; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        tb2 = s.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.8), Inches(0.3))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = sub; p2.font.size = Pt(12); p2.font.color.rgb = RGBColor(0xBB,0xCC,0xDD)
    accent = s.shapes.add_shape(1, Inches(0), Inches(0.95), prs.slide_width, Inches(0.04))
    accent.fill.solid(); accent.fill.fore_color.rgb = ORANGE; accent.line.fill.background()

def footer(s, n):
    tb = s.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(11.8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"B166ER Research · 2026五一酒店市场预测V2.0 · {n}"
    p.font.size = Pt(8); p.font.color.rgb = LGRAY; p.alignment = PP_ALIGN.RIGHT

def section_divider(s, num, title, body):
    bg(s, DARK)
    for shape in [Inches(0.3),Inches(3.1),Inches(3.5),Inches(3.9)]:
        ln = s.shapes.add_shape(1, shape, Inches(3.5), Inches(0.6), Inches(0.03))
        ln.fill.solid(); ln.fill.fore_color.rgb = ORANGE if shape==Inches(3.5) else RGBColor(0x3A,0x7C,0xBF); ln.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.8), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.text = f"PART {num}"; p.font.size = Pt(14); p.font.color.rgb = ORANGE; p.font.bold = True
    tb2 = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.8), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = title; p2.font.size = Pt(36); p2.font.bold = True; p2.font.color.rgb = WHITE
    tb3 = s.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(11.8), Inches(1.5))
    p3 = tb3.text_frame.paragraphs[0]; p3.text = body; p3.font.size = Pt(14); p3.font.color.rgb = RGBColor(0xBB,0xCC,0xDD); p3.line_spacing = Pt(22)

def kpi_box(s, x, y, w, h, label, val, sub, color=BLUE):
    b = s.shapes.add_shape(1, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    b.shadow.inherit = False
    items = [
        (Inches(0.15), Inches(0.12), label, Pt(11), False, RGBColor(0xCC,0xDD,0xEE)),
        (Inches(0.25), Inches(0.50), val, Pt(26), True, WHITE),
        (Inches(0.35), Inches(1.05), sub, Pt(10), False, RGBColor(0xCC,0xDD,0xEE)),
    ]
    for off, top, txt, sz, bld, clr in items:
        tb = s.shapes.add_textbox(x+off, y+top, w-Inches(0.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = txt; p.font.size = sz; p.font.bold = bld; p.font.color.rgb = clr

def matrix_table(s, x, y, w, h, headers, rows, highlights=None):
    """Clean data table with proper formatting."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = h / n_rows
    tbl = s.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl.table.style = None
    for c, hdr in enumerate(headers):
        cell = tbl.table.cell(0, c)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK
        for p in cell.text_frame.paragraphs: p.font.size=Pt(10); p.font.bold=True; p.font.color.rgb=WHITE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.table.cell(r+1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r%2==0 else RGBColor(0xF0,0xF4,0xF8)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.5)
                if highlights and r in highlights and c in highlights[r]:
                    p.font.bold = True
                    p.font.color.rgb = ORANGE if highlights[r][c]=='up' else RED

def hchart(s, x, y, w, h, title, cats, vals, colors=None):
    """Horizontal bar chart."""
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('', vals)
    cf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, cd)
    ch = cf.chart; ch.has_legend = False; ch.has_title = False
    ch.category_axis.tick_labels.font.size = Pt(9)
    ch.value_axis.visible = False
    if colors:
        for i,c in enumerate(colors):
            ch.plots[0].series[0].points[i].format.fill.solid()
            ch.plots[0].series[0].points[i].format.fill.fore_color.rgb = c
    tb = s.shapes.add_textbox(x, y-Inches(0.3), w, Inches(0.25))
    p = tb.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=DARK

def vchart(s, x, y, w, h, title, cats, vals, colors=None):
    """Vertical bar chart."""
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('', vals)
    cf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    ch = cf.chart; ch.has_legend = False; ch.has_title = False
    ch.category_axis.tick_labels.font.size = Pt(8)
    ch.value_axis.visible = False
    if colors:
        for i,c in enumerate(colors):
            ch.plots[0].series[0].points[i].format.fill.solid()
            ch.plots[0].series[0].points[i].format.fill.fore_color.rgb = c
    tb = s.shapes.add_textbox(x, y-Inches(0.3), w, Inches(0.25))
    p = tb.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=DARK

def doughnut(s, x, y, w, h, title, cats, vals, colors):
    """Donut chart."""
    cd = CategoryChartData(); cd.categories = cats; cd.add_series('', vals)
    cf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd)
    ch = cf.chart; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.font.size = Pt(8); ch.has_title = False
    for i,c in enumerate(colors):
        ch.plots[0].series[0].points[i].format.fill.solid()
        ch.plots[0].series[0].points[i].format.fill.fore_color.rgb = c
    tb = s.shapes.add_textbox(x, y-Inches(0.3), w, Inches(0.25))
    p = tb.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=DARK

def text_block(s, x, y, w, h, title, items, color=GRAY):
    """Text list with title."""
    tb = s.shapes.add_textbox(x, y, w, Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(14); p.font.bold=True; p.font.color.rgb=DARK
    tb2 = s.shapes.add_textbox(x, y+Inches(0.35), w, h-Inches(0.35))
    tf = tb2.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = f"• {item}"; p.font.size = Pt(10.5); p.font.color.rgb = color; p.space_after = Pt(5)

def callout_box(s, x, y, w, h, text, icon="💡"):
    b = s.shapes.add_shape(1, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = HL; b.line.color.rgb = ORANGE; b.line.width = Pt(1.5)
    tb = s.shapes.add_textbox(x+Inches(0.2), y+Inches(0.1), w-Inches(0.4), h-Inches(0.2))
    p = tb.text_frame.paragraphs[0]; p.text = f"{icon} {text}"; p.font.size = Pt(11); p.font.color.rgb = DARK

def icon_card(s, x, y, w, h, icon, title, body, accent_color=BLUE):
    card = s.shapes.add_shape(1, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = RGBColor(0xDD,0xDD,0xDD); card.line.width = Pt(0.5)
    top = s.shapes.add_shape(1, x, y, w, Inches(0.06))
    top.fill.solid(); top.fill.fore_color.rgb = accent_color; top.line.fill.background()
    items = [
        (Inches(0.15), Inches(0.15), icon, Pt(18), False, accent_color),
        (Inches(0.15), Inches(0.45), title, Pt(13), True, DARK),
        (Inches(0.15), Inches(0.85), body, Pt(9.5), False, GRAY),
    ]
    for off, top_y, txt, sz, bld, clr in items:
        tb = s.shapes.add_textbox(x+off, y+top_y, w-Inches(0.3), Inches(0.4))
        p = tb.text_frame.paragraphs[0]; p.text = txt; p.font.size = sz; p.font.bold = bld; p.font.color.rgb = clr
        tb.text_frame.word_wrap = True

# ============================================================
# SLIDE 1: COVER
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
shapes = [
    (Inches(0),Inches(0),Inches(13.333),Inches(7.5), DARK),
    (Inches(0),Inches(7.2),Inches(13.333),Inches(0.3), ORANGE),
]
for x,y,w,h,c in shapes:
    sh = s.shapes.add_shape(1,x,y,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

# Thin decorative lines
for lx in [2.5, 3.5, 5.0]:
    ln = s.shapes.add_shape(1, Inches(lx), Inches(3.3), Inches(1.0), Inches(0.02))
    ln.fill.solid(); ln.fill.fore_color.rgb = ORANGE if lx==3.5 else RGBColor(0x3A,0x7C,0xBF); ln.line.fill.background()

tb = s.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = "2026年五一小长假"; p.font.size=Pt(48); p.font.bold=True; p.font.color.rgb=WHITE
p2 = tb.text_frame.add_paragraph(); p2.text = "全国酒店市场全面预测分析报告"; p2.font.size=Pt(32); p2.font.color.rgb=ORANGE

tb2 = s.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(0.6))
p = tb2.text_frame.paragraphs[0]; p.text = "收益管理驱动 · 业主视角 · STR六档分级 · 七区对比 · 18数据源交叉验证"
p.font.size = Pt(13); p.font.color.rgb = RGBColor(0xBB,0xCC,0xDD)

tb3 = s.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10), Inches(0.4))
p = tb3.text_frame.paragraphs[0]; p.text = "数据截至 2026-05-01 22:00 CST  |  V2.0  |  B166ER Research"
p.font.size = Pt(11); p.font.color.rgb = RGBColor(0x88,0x99,0xAA)

# ============================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title_bar(s,"核心结论：供需共振，拐点确立","宏观出行 → 酒店绩效 → 供给格局 三维交叉验证"); footer(s,"02")

# Key message at top
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(11.9), Inches(0.45))
p = tb.text_frame.paragraphs[0]; p.text = "▌ RevPAR +5~7%（中枢+6%）· ADR全档次转正 · 供给增速+6.5%↓ · 浩华Q2景气指数-9（↑12点）"
p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = DARK

# 6 KPI boxes in 2 rows
kpis = [
    ("跨区域流动", "15.2亿人次", "+4% YoY", BLUE),
    ("酒店RevPAR", "+5~7%", "中枢+6%", ORANGE),
    ("酒店ADR", "+6~7%", "全档次转正", LBLUE),
    ("间夜量", "+6~7%", "跳城游+替代游", GREEN),
    ("供给增速", "+6.5%↓", "年初+7.4%→4月+6.5%", RGBColor(0x8B,0x45,0x13)),
    ("浩华景气指数", "-9", "Q1为-21，↑12点", RGBColor(0x6A,0x1B,0x9A)),
]
for i, (l,v,sv,clr) in enumerate(kpis):
    col = i % 3; row = i // 3
    kpi_box(s, Inches(0.5+col*4.1), Inches(1.75+row*1.65), Inches(3.8), Inches(1.4), l, v, sv, clr)

callout_box(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.65),
    "2026是中国弹性假期制度元年。供给收敛+需求质变+价格修复=行业拐点。首日3.44亿人次验证预测基础，跳城游(60%)+海外替代游+宝藏小城三重超预期。"
)

callout_box(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
    "全年四重支撑：①供给降速(持续) ②春秋假加成(Q2/Q3) ③商旅禁酒令低基数(H2) ④头部分红回购(全年) → 当前是布局全年行情的较好窗口 —— 中信证券"
)

# ============================================================
# SLIDE 3: DIVIDER - PART 1
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "01", "需求端：五大结构性趋势", "弹性假期制度元年 × 跳城游 × 海外替代游 × 宝藏小城 × Color Walk")

# ============================================================
# SLIDE 4: MACRO TRAVEL DEMAND
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title_bar(s,"出行总量与方式结构","交通运输部4/28预测 + 5/1首日确认：3.44亿人次/+3.3%"); footer(s,"04")

# Left: travel mode donut
doughnut(s, Inches(0.3), Inches(1.3), Inches(4.5), Inches(3.0),
    "出行方式结构 (15.2亿人次)",
    ['公路自驾 91.6%', '铁路 7.0%', '民航 0.8%', '水路 0.6%'],
    [91.6, 7.0, 0.8, 0.6],
    [BLUE, LBLUE, ORANGE, LGRAY])

# Middle: traveler demographic donut
doughnut(s, Inches(4.8), Inches(1.3), Inches(4.0), Inches(3.0),
    "客群画像 (航空出行，首日)",
    ['年轻客群 41%', '亲子家庭 25%', '银发一族 13%', '入境/其他 21%'],
    [41, 25, 13, 21],
    [LBLUE, BLUE, GREEN, LGRAY])

# Right: key numbers
text_block(s, Inches(9.0), Inches(1.3), Inches(3.8), Inches(3.0),
    "首日关键数据",
    ["铁路日均1.3万列，+5.2%",
     "高速车流7000万辆次，+3.5%",
     "航空含税均价+13~15%",
     "出境替代·境内转移显著",
     "亲子出行占比从32%→59%",
     "跨省出行占比+15ppts",
     "住宿天数+1~2天"])

# Bottom: 5 trend icon cards
trends = [
    ("🏃", "跳城游", "60%旅客串联2-3城，线性串联订单40%+，租车自驾50%+，租期4.3天", ORANGE),
    ("✈️", "海外替代游", "云南+89% 四川+76% 新疆+73% 广西+70%，跟团客单价+13%", BLUE),
    ("🏘️", "宝藏小城", "县域+128%，丽水+116% 龙岩+85% 赣州+75%，品质酒店+76%", GREEN),
    ("🎨", "Color Walk", "色彩目的地搜索+200%，蓝色系+65%，情绪价值驱动选择", LBLUE),
    ("🤖", "AI预订", "DeepTrip引导预订+80%，规划行程+60%，OTA入口被分流", RGBColor(0x6A,0x1B,0x9A)),
]
for i, (ic, t, bd, clr) in enumerate(trends):
    icon_card(s, Inches(0.3+i*2.55), Inches(4.55), Inches(2.4), Inches(2.3), ic, t, bd, clr)

# ============================================================
# SLIDE 5: DIVIDER - PART 2
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "02", "供给端：增速收敛，供需共振", "全国酒店客房增速持续下行 · 浩华：2025新增中档+约9,800家 · 连锁化率40.09% · 一超多强")

# ============================================================
# SLIDE 6: SUPPLY & COMPETITION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); title_bar(s,"供给收敛与竞争格局","供给增速从+7.4%持续下行至+6.5% | 四大集团2025年报：华住一超、亚朵崛起"); footer(s,"06")

# Supply trend chart
vchart(s, Inches(0.3), Inches(1.4), Inches(4.0), Inches(2.6),
    "全国酒店客房量同比增速(%)",
    ['2026年初', '2月', '3月', '4月底'],
    [7.4, 7.0, 6.8, 6.5],
    [BLUE, BLUE, LBLUE, ORANGE])

# Hotel group revenue/profit
cats = ['华住', '亚朵', '锦江', '首旅']
rev = [253, 98, 138, 76]
profit = [50.8, 16.2, 9.25, 8.1]

cd = CategoryChartData(); cd.categories = cats
cd.add_series('营收(亿元)', rev); cd.add_series('净利润(亿元)', profit)
cf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(4.5), Inches(1.4), Inches(4.0), Inches(2.6), cd)
ch = cf.chart; ch.has_title = False
if ch.legend:
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.font.size = Pt(8)
ch.plots[0].series[0].format.fill.solid(); ch.plots[0].series[0].format.fill.fore_color.rgb = BLUE
ch.plots[0].series[1].format.fill.solid(); ch.plots[0].series[1].format.fill.fore_color.rgb = ORANGE
ch.category_axis.tick_labels.font.size = Pt(8)
ch.value_axis.visible = False

# Ownership structure
doughnut(s, Inches(8.8), Inches(1.4), Inches(4.0), Inches(2.6),
    "中国酒店连锁化率 40.09%\n（2024年末，中国饭店协会）",
    ['连锁', '单体'],
    [40, 60],
    [BLUE, LGRAY])

# Bottom analysis
text_block(s, Inches(0.3), Inches(4.3), Inches(4.0), Inches(2.5),
    "供给端关键判断",
    ["存量改造成为主战场：开业5年+占60%",
     "经济型存量占77.8%，翻牌/轻改造崛起",
     "租金红利→品牌红利+运营红利",
     "回本周期拉长至5.4年，51%项目>5年",
     "龙头储备店较高峰期放缓"])

matrix_table(s, Inches(4.8), Inches(4.3), Inches(8.0), Inches(2.5),
    ['维度', '华住', '亚朵', '锦江', '首旅'],
    [
        ['定位', '一超', '崛起新贵', '规模第一', '中端转型'],
        ['营收(亿)', '253', '98', '138', '76'],
        ['净利润(亿)', '50.8', '16.2', '9.25', '8.1'],
        ['经营客房', '124.6万间', '22万间', '100万+间', '55万间'],
        ['五一策略', '汉庭存量改造', '见野品牌独立', '丽祺度假', '如家4.0焕新'],
    ])

# ============================================================
# SLIDE 7: DIVIDER - PART 3
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "03", "分档次预测 · STR六档标准", "Luxury → Upper Upscale → Upscale → Upper Midscale → Midscale → Economy")

# ============================================================
# SLIDE 8: STR SIX-TIER CLASSIFICATION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s,"STR标准六档分类与五一表现预测","参考标准：酒店行业全景知识库 §2.1 | 数据源：酒店之家4/20周报 + 中信证券预测")
footer(s,"08")

# STR Classification table
matrix_table(s, Inches(0.3), Inches(1.2), Inches(12.7), Inches(2.5),
    ['STR档次', '中文', 'ADR区间(¥)', '代表品牌', '五一ADR同比', '五一OCC预测', 'RevPAR预测', '数据说明'],
    [
        ['Luxury', '奢华', '1500+', '安缦·宝格丽·丽思', '+0.2%', '60-75%', '+1~3%', '豪华型0.4%中包含；奢华修复最慢'],
        ['Upper Upscale', '高端', '800-1500', '洲际·万豪·希尔顿·凯悦', '+0.4%', '65-78%', '+2~4%', '酒店之家「豪华型」数据'],
        ['Upscale', '中高端', '400-800', '美居·桔子水晶·亚朵', '+2.3%', '75-88%', '+4~7%', '酒店之家「高档型」数据 ✓'],
        ['Upper Midscale', '中端', '250-400', '全季·维也纳·麗枫', '+3.9%', '85-95%', '+6~9%', '酒店之家「中档型」数据 ✓'],
        ['Midscale', '经济中端', '150-250', '汉庭·如家·7天', '+3.6%', '90-98%', '+5~7%', '酒店之家「经济型」数据 ✓'],
        ['Economy', '经济', '100-150', '99旅馆·布丁·贝壳', '+2.8%', '85-95%', '+4~6%', '估算：县域低端ADR弹性较弱'],
    ])

# ADR recovery comparison chart
vchart(s, Inches(0.3), Inches(4.0), Inches(6.5), Inches(2.8),
    "各档次ADR同比对比：2025五一 vs 2026五一",
    ['Luxury','U.Upscale','Upscale','U.Midscale','Midscale','Economy'],
    [0.2, 0.4, 2.3, 3.9, 3.6, 2.8],
    [LGRAY, RGBColor(0x66,0x99,0xBB), LBLUE, BLUE, BLUE, RGBColor(0x55,0x88,0xBB)])

# Right: Key insight
text_block(s, Inches(7.2), Inches(4.0), Inches(5.5), Inches(2.8),
    "▌ 核心发现",
    ["2025五一：各档次ADR全面负值(-0.3~-4.5%)",
     "2026五一：各档次ADR全面转正 → 价格修复周期开启",
     "中端(Upper Midscale)是五一最大受益者：",
     "  → ADR+3.9% + OCC 85-95% → RevPAR+6~9%",
     "  → 县域品质酒店+76% + 精品民宿+92%",
     "  → 亲子家庭首选+跳城游中间站",
     "奢华/高端ADR修复滞后：高端商务仍承压",
     "  → 出境替代效应部分弥补（三亚+47%）"])

# ============================================================
# SLIDE 9: DIVIDER - PART 4
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "04", "城市分级与区域市场对比", "四级城市体系 × 七大地理区划 × 首日酒店热度验证")

# ============================================================
# SLIDE 10: CITY TIER COMPARISON
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s,"城市分级酒店表现对比","一线(北上广深) · 新一线(15城) · 二线(30城) · 三四线(县域)")
footer(s,"10")

# City tier RevPAR chart
vchart(s, Inches(0.3), Inches(1.3), Inches(6.5), Inches(2.8),
    "各城市等级 RevPAR 同比预测 (%)",
    ['一线城市', '新一线', '二线城市', '三四线/县域'],
    [5, 8, 12, 22],
    [BLUE, LBLUE, ORANGE, GREEN])

# City tier detail table
matrix_table(s, Inches(7.2), Inches(1.3), Inches(5.5), Inches(2.8),
    ['城市等级', '代表城市', 'RevPAR预测', '核心逻辑'],
    [
        ['一线', '北京·上海·广州·深圳', '+4~6%', '商务筑底+入境增量'],
        ['新一线', '成都·杭州·重庆·武汉·西安\n南京·苏州·天津·长沙·青岛', '+6~10%', '替代游核心+跳城起点'],
        ['二线', '昆明·三亚·厦门·大连·贵阳\n南宁·南昌·福州·合肥·太原\n哈尔滨·长春·呼和浩特·乌鲁木齐\n兰州·银川·西宁·拉萨·海口', '+8~15%', 'Color Walk+度假升级\n长线门户·避暑'],
        ['三四线', '丽水·赣州·龙岩·平潭\n安吉·阳朔·荔波·建水·忻州', '+15~30%', '反向旅游+品质下沉'],
    ])

# Bottom: TOP10 city heat rankings
vchart(s, Inches(0.3), Inches(4.4), Inches(6.5), Inches(2.3),
    "首日酒店入住热度飙升TOP10（同比%）",
    ['丽水', '龙岩', '赣州', '县域', '广州', '青岛', '成都', '武汉', '三亚高品', '长沙'],
    [116, 85, 75, 128, 43, 40, 37, 35, 47, 35],
    [GREEN, GREEN, GREEN, GREEN, BLUE, BLUE, BLUE, BLUE, ORANGE, BLUE])

text_block(s, Inches(7.2), Inches(4.4), Inches(5.5), Inches(2.3),
    "▌ 首日热点城市群（同程大数据）",
    ["华南线：广州+43%—佛山—珠海 岭南文化串联",
     "浙西线：杭州—千岛湖—衢州—丽水+116%",
     "闽地线：福州—泉州—厦门 山海黄金线",
     "西北替代：云南跟团+89% 四川+76% 新疆+73%",
     "三四线城市文旅热度增速首超一线",
     "「高铁抵达+落地自驾」县域深度游主流",
     "县域租车+113%，平均租期3.8天"])

# ============================================================
# SLIDE 11: REGIONAL COMPARISON
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s,"七大区域酒店市场对比","华东 · 华南 · 华北 · 华中 · 西南 · 西北 · 东北")
footer(s,"11")

# Regional RevPAR chart
colors_region = [BLUE, ORANGE, LBLUE, GREEN, BLUE, ORANGE, LGRAY]
hchart(s, Inches(0.3), Inches(1.3), Inches(6.5), Inches(3.2),
    "各区域 RevPAR 同比预测 (%)",
    ['西南', '西北', '华东', '华南', '华中', '东北', '华北'],
    [10, 11.5, 8.5, 7.5, 6.5, 5, 4],
    colors_region[::-1])

# Regional detail matrix
matrix_table(s, Inches(7.2), Inches(1.3), Inches(5.5), Inches(3.2),
    ['区域', 'RevPAR预期', '代表省份/城市', '驱动逻辑'],
    [
        ['华东', '+7~10%', '浙江·江苏·山东·上海\n丽水+116% 青岛+40%', '春假联动+Color Walk'],
        ['华南', '+6~9%', '广东·海南·广西\n广州+43% 三亚+47%', '跳城游起点+蓝色系'],
        ['西南', '+8~12%', '四川·云南·重庆·贵州\n跟团+89%/76%', '海外替代游核心'],
        ['西北', '+8~15%', '新疆·甘肃·陕西\n伊犁+22% 赛里木+58%', '超长线替代·野奢'],
        ['华中', '+5~8%', '湖北·湖南·河南\n武汉+35% 洛阳+30%', 'Color Walk赏花'],
        ['东北', '+4~6%', '辽宁·吉林·黑龙江\n哈尔滨·长白山', '避暑+亲子'],
        ['华北', '+3~5%', '北京·天津·山西\n忻州古城热门', '商旅筑底+入境'],
    ])

# Bottom insight
callout_box(s, Inches(0.3), Inches(4.8), Inches(12.7), Inches(0.65),
    "区域分化核心逻辑：西南/西北受益「海外替代游」最强 → 跟团游暴涨70-89%；华东受益「春假+跳城游」→ 县域/品质酒店领涨；华北/东北修复偏慢 → 商务需求筑底中。"
)
callout_box(s, Inches(0.3), Inches(5.65), Inches(12.7), Inches(0.6),
    "浩华Q2景气验证：华南(-15)与华东(-17)高于全国均值(-9)，三亚(+33)领跑全国，深圳/上海住宿率转正，北京(-25)低于平均 → 南北修复梯度基本吻合。"
)

# ============================================================
# SLIDE 12: DIVIDER - PART 5
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "05", "消费趋势与渠道变革", "客单价·预订渠道·客源来源·「体验复购」取代「流量驱动」")

# ============================================================
# SLIDE 13: CONSUMER TRENDS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s,"消费趋势与客源结构","客单价升级 · 预订渠道迁移 · 入境游持续高景气"); footer(s,"13")

# Spending upgrade chart
vchart(s, Inches(0.3), Inches(1.3), Inches(4.0), Inches(2.4),
    "消费升级指标（同比%）",
    ['跟团客单价', '机票均价', '5晚+连住', '县品质酒店'],
    [13, 10, 10, 76],
    [BLUE, LBLUE, LBLUE, ORANGE])

# Booking channel
vchart(s, Inches(4.6), Inches(1.3), Inches(4.0), Inches(2.4),
    "预订渠道变革信号",
    ['DeepTrip\nAI引导', '传统OTA\n(携程)', '酒店直销\n/会员'],
    [80, 5, 5],
    [ORANGE, BLUE, LBLUE])

# Source market donut
doughnut(s, Inches(8.9), Inches(1.3), Inches(4.0), Inches(2.4),
    "客源TOP10 + 入境TOP6",
    ['北京·上海·成都\n南京·杭州·广州', '天津·西安·苏州\n深圳', '俄·日·韩·马·泰·美'],
    [45, 35, 20],
    [BLUE, LBLUE, ORANGE])

# Bottom text
text_block(s, Inches(0.3), Inches(4.0), Inches(4.0), Inches(2.7),
    "消费行为变革",
    ["跨省出行占比+15ppts（首日）",
     "住宿天数+1~2天",
     "「跳城游」60%旅客串联2-3城",
     "租车自驾50%+，平均租期4.3天",
     "海外替代游客单价+13%",
     "「可以买贵的，不能买贵了」理性消费"])

text_block(s, Inches(4.6), Inches(4.0), Inches(4.0), Inches(2.7),
    "渠道与营销变革",
    ["AI预订入口高速渗透(DeepTrip+80%)",
     "OTA反垄断 → 佣金结构或调整",
     "Color Walk社交驱动选择(+200%)",
     "体验复购 > 流量驱动（同程判断）",
     "酒店数据需结构化，被AI可读取",
     "小红书/抖音内容→OTA成交链路"])

text_block(s, Inches(8.9), Inches(4.0), Inches(4.0), Inches(2.7),
    "消费者画像与来源",
    ["亲子家庭59%（比去年+27ppts）",
     "银发一族16%（高品质长线）",
     "年轻客群41%（跳城+Color Walk）",
     "入境机票预订55万+，+20%",
     "数字化入境卡→机场45分钟通关",
     "入境客消费力强，高星酒店受益"])

# ============================================================
# SLIDE 14: DIVIDER - PART 6
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "06", "风险、策略与展望", "定价窗口 · 产品渠道策略 · 完整复盘时间线")

# ============================================================
# SLIDE 15: RISK & STRATEGY
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s,"风险因素与业主操作策略","定价窗口 · 产品策略 · 渠道布局 · 五一后行动清单"); footer(s,"15")

# Risk matrix
matrix_table(s, Inches(0.3), Inches(1.2), Inches(6.0), Inches(2.3),
    ['风险因素', '概率', '影响', 'Day1验证'],
    [
        ['出境回流超预期(上行)', '🔴高', '利好国内酒店', '✅已显现'],
        ['油价上行抑制航空', '🟡中', '民航弱于预测', '暂未显著'],
        ['消费力压制高客单价', '🟡中', '豪华ADR承压', '替代游高端对冲'],
        ['五一后需求透支回落', '🟡中', '5月中下旬', '⚠️警惕'],
        ['天气突变', '🟢低', '局部影响', '晴好'],
    ])

# Pricing windows
text_block(s, Inches(6.6), Inches(1.2), Inches(3.2), Inches(2.3),
    "定价窗口", [
        "🔴 5/1-2 峰值 | 首日3.44亿",
        "🟡 5/2-3 微调 | 票价-40%",
        "🟠 5/4-5 关注 | 返程退房潮",
        "🟢 5/6-10 错峰 | 延后出行",
    ])

# Strategy cards
icon_card(s, Inches(0.3), Inches(3.8), Inches(4.0), Inches(2.8),
    "🎯", "产品策略", 
    "• 亲子套餐：59%亲子·联通房+活动\n• 连住优惠：5晚+订单+10ppts\n• Color Walk场景：出片率=转化率\n• 县域精品民宿：提前15天爆满", BLUE)

icon_card(s, Inches(4.6), Inches(3.8), Inches(4.0), Inches(2.8),
    "📡", "渠道策略",
    "• OTA为主但监管收紧·加强直销\n• AI预订入口崛起·酒店数据需结构化\n• 小红书/抖音种草→成交链路\n• 落地租车·酒店+租车打包套餐", LBLUE)

icon_card(s, Inches(8.9), Inches(3.8), Inches(4.0), Inches(2.8),
    "📋", "五一后行动清单",
    "• 5/7 快报复盘生成（OTA+交通部）\n• 5/12 完整复盘（酒店之家+券商）\n• 5/25 端午前预测自动生成\n• 跟踪浩华Q3景气指数发布时间", ORANGE)

# ============================================================
# SLIDE 16: CLOSING
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, DARK)
ln = s.shapes.add_shape(1, Inches(1.5), Inches(2.0), Inches(2.0), Inches(0.03))
ln.fill.solid(); ln.fill.fore_color.rgb = ORANGE; ln.line.fill.background()

tb = s.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = "行业拐点 · 供需共振"; p.font.size = Pt(42); p.font.bold = True; p.font.color.rgb = WHITE

tb2 = s.shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(10), Inches(0.8))
p = tb2.text_frame.paragraphs[0]
p.text = "供给增速收敛(+7.4%→+6.5%)  ×  需求结构质变(跳城游+替代游+小城)  ×  价格修复确立(全档次ADR转正)"
p.font.size = Pt(14); p.font.color.rgb = RGBColor(0xBB,0xCC,0xDD)

# 3 final columns
for i, (title, val, desc) in enumerate([
    ("间夜量", "+6~7%", "出行+4%·跳城游\n替代游双引擎"),
    ("ADR", "+6~7%", "全档次转正\n中高端实质性修复"),
    ("供给", "+6.5%↓", "持续收敛\n供需共振格局形成"),
]):
    x = Inches(1.5 + i*3.5)
    card = s.shapes.add_shape(1, x, Inches(4.5), Inches(3.1), Inches(2.0))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x15,0x2E,0x50); card.line.fill.background()
    closings = [
        (Inches(0.2), Inches(0.15), title, Pt(14), RGBColor(0x88,0xAA,0xCC)),
        (Inches(0.2), Inches(0.55), val, Pt(24), ORANGE),
        (Inches(0.2), Inches(1.1), desc, Pt(11), RGBColor(0x99,0xBB,0xDD)),
    ]
    for off, dy, txt, sz, clr in closings:
        tb = s.shapes.add_textbox(x+off, Inches(4.5)+dy, Inches(2.7), Inches(0.4))
        p = tb.text_frame.paragraphs[0]; p.text = txt; p.font.size = sz
        p.font.bold = (dy==Inches(0.55)); p.font.color.rgb = clr

# Source note
tb3 = s.shapes.add_textbox(Inches(1.5), Inches(6.8), Inches(10), Inches(0.35))
p = tb3.text_frame.paragraphs[0]
p.text = "完整复盘：5/12  |  数据日历：DATA-CALENDAR.md  |  建议导出DPI≥300 · 18数据源：交通运输部·酒店之家·浩华·中信·同程·携程·迈点·环球旅讯·STR·戴德梁行等"
p.font.size = Pt(9); p.font.color.rgb = RGBColor(0x77,0x88,0x99); p.alignment = PP_ALIGN.CENTER

# ============================================================
output = r'C:\Users\Administrator\Desktop\2026五一全国酒店市场预测分析报告.pptx'
prs.save(output)
print(f'OK: {output} ({len(prs.slides)} slides)')
