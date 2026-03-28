# -*- coding: utf-8 -*-
import os
from pathlib import Path

folder = Path(r'E:\桌面20250702')

# Get all files sorted by size
files = []
for f in folder.iterdir():
    if f.is_file() and f.stat().st_size > 0:
        files.append((f.stat().st_size, f))

files.sort(reverse=True)

# Write to output
output = []
output.append("# E:\\桌面20250702 文件内容摘要\n")
output.append(f"共 {len(files)} 个文件\n\n")

from docx import Document
import pdfplumber

for size, f in files:
    fname = f.name
    output.append(f"\n## 文件: {fname}\n")
    output.append(f"大小: {size//1024}KB\n")
    
    ext = f.suffix.lower()
    
    try:
        if ext == '.docx':
            doc = Document(str(f))
            text = '\n'.join([p.text for p in doc.paragraphs])
            output.append(f"类型: Word文档\n")
            output.append(f"内容:\n{text[:3000]}\n")
        elif ext == '.doc':
            # Try reading as text with different encoding
            output.append(f"类型: Word97文档\n")
            try:
                # python-docx doesn't support .doc, try textract
                import subprocess
                result = subprocess.run(['antiword', str(f)], capture_output=True, timeout=10)
                if result.returncode == 0:
                    output.append(f"内容:\n{result.stdout.decode('utf-8', errors='replace')[:3000]}\n")
                else:
                    output.append("内容: [无法读取，需要转换格式]\n")
            except:
                output.append("内容: [需要转换为docx格式]\n")
        elif ext == '.pdf':
            output.append(f"类型: PDF文档\n")
            try:
                with pdfplumber.open(str(f)) as pdf:
                    text = ''
                    for page in pdf.pages[:10]:
                        t = page.extract_text()
                        if t:
                            text += t + '\n'
                    output.append(f"内容:\n{text[:3000]}\n")
            except Exception as e:
                output.append(f"内容: [读取失败: {str(e)[:100]}]\n")
        elif ext == '.html':
            output.append(f"类型: HTML\n")
            with open(str(f), 'r', encoding='utf-8', errors='replace') as hf:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(hf.read(), 'html.parser')
                output.append(f"内容:\n{soup.get_text()[:5000]}\n")
        elif ext == '.txt':
            output.append(f"类型: 文本\n")
            with open(str(f), 'r', encoding='utf-8', errors='replace') as tf:
                output.append(f"内容:\n{tf.read()[:3000]}\n")
        elif ext == '.pptx':
            output.append(f"类型: PPTX\n")
            try:
                from pptx import Presentation
                prs = Presentation(str(f))
                text = ''
                for i, slide in enumerate(prs.slides[:10]):
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            text += shape.text + '\n'
                output.append(f"内容:\n{text[:2000]}\n")
            except Exception as e:
                output.append(f"内容: [读取失败: {str(e)[:100]}]\n")
        elif ext in ['.jpg', '.jpeg', '.png', '.pdf']:
            # Could be image
            output.append(f"类型: 可能是图片\n")
        else:
            output.append(f"类型: {ext}\n")
    except Exception as e:
        output.append(f"读取错误: {str(e)[:200]}\n")

# Write output
with open(r'C:\Users\ericz\.openclaw\workspace\docs\E-桌面20250702-文件内容摘要-V1.0.md', 'w', encoding='utf-8') as out:
    out.write('\n'.join(output))

print('Done!')
