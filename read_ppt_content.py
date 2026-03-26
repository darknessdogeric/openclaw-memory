# -*- coding: utf-8 -*-
from pptx import Presentation
import sys
import os

# Set UTF-8
sys.stdout.reconfigure(encoding='utf-8')

file_path = r'C:\Users\ericz\Desktop\张实项目总控\05-AHL-去中心化旅行平台\项目说明书\AHL-LLM去中心化旅行平台商业计划书V5.1(9).pptx'

prs = Presentation(file_path)
print(f"Total pages: {len(prs.slides)}")
print(f"Slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
print()

for i, slide in enumerate(prs.slides, 1):
    print(f"=== Page {i} ===")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()[:150]
            print(f"  {text}")
    print()
