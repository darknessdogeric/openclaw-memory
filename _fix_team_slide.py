# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

final_path = r'C:\Users\ericz\Desktop\AHL酒店智能升级项目 - 合作意向书.pptx'
prs = Presentation(final_path)

slide = prs.slides[11]
spTree = slide.shapes._spTree
for shape in list(slide.shapes):
    spTree.remove(shape._element)

# ========== 背景：白色 ==========
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ========== 顶部蓝色标题栏 ==========
header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.05))
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = RGBColor(0x00, 0x5A, 0xB5)
header_bar.line.fill.background()

title_en = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12), Inches(0.45))
tf = title_en.text_frame
p = tf.paragraphs[0]
p.text = 'CORE TEAM  核心团队'
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

title_cn = slide.shapes.add_textbox(Inches(0.5), Inches(0.58), Inches(12), Inches(0.4))
tf = title_cn.text_frame
p = tf.paragraphs[0]
p.text = '黄金三角组合：产业 + 技术 + 战略'
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(0xCC, 0xE4, 0xFF)

# ========== 引言段落 ==========
intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.55))
tf = intro_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
intro_text = '我们不是传统的创业团队，而是一支范式革命的先遣队。三位核心成员分别代表产业洞察、技术实现、战略与生态，形成完整的"问题-方案-落地"闭环。'
p.text = intro_text
p.font.size = Pt(10.5)
p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p.font.italic = True

# ========== 颜色常量 ==========
BLUE_DARK = RGBColor(0x00, 0x5A, 0xB5)
BLUE_MID  = RGBColor(0x00, 0x8A, 0xD4)
BLUE_LIGHT= RGBColor(0x00, 0xA8, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_GRAY = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT= RGBColor(0x88, 0x88, 0x88)
GOLD = RGBColor(0xE8, 0xA8, 0x00)
LINE_COLOR= RGBColor(0xDD, 0xDD, 0xDD)

# ========== 三个成员 ==========
members = [
    {
        'name': '张实',
        'en': 'Eric Zhang',
        'title': '项目总控 / 发起人',
        'core': '住宿业范式革命的提出者与总设计师',
        'bg': BLUE_DARK,
        'bullets': [
            '北京第二外国语学院旅游管理专业',
            '24年酒店业深耕经验',
            '曾任岷山集团副总（14年）',
            '创立四川远途酒店管理公司',
            '维景酒店副总/总经理',
        ],
        'comp': ['24年酒店经验', '产业洞察', '运营Know-how'],
    },
    {
        'name': '李源',
        'en': 'Li Yuan',
        'title': '首席技术官 (CTO)',
        'core': 'AI技术架构的构建者与算法壁垒的建立者',
        'bg': BLUE_MID,
        'bullets': [
            '华中科技大学人工智能博士',
            '10年+AI技术研发经验',
            '主导20+AI工程落地',
            '核心技术：大模型微调/AGENT编排/知识图谱',
        ],
        'comp': ['10年+AI经验', '技术壁垒', '算法能力'],
    },
    {
        'name': '陈思序',
        'en': 'Chen Sixu',
        'title': '首席战略与生态官 (CSO)',
        'core': '范式革命的战略推动者与生态建设者',
        'bg': BLUE_LIGHT,
        'bullets': [
            '世界500强企业战略高管',
            '深耕PE/VC投资领域',
            '具备IPO上市全流程经验',
            '负责战略制定和生态建设',
        ],
        'comp': ['资本运作', '战略规划', '生态资源'],
    },
]

col_w = 3.9
col_h = 4.55
start_x = 0.45
start_y = 1.8
gap = 0.3

for i, m in enumerate(members):
    cx = start_x + i * (col_w + gap)
    cy = start_y

    # 卡片背景
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(col_w), Inches(col_h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFF)
    card.line.color.rgb = LINE_COLOR
    card.line.width = Pt(0.5)

    # 顶部色条
    top_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(col_w), Inches(0.08))
    top_strip.fill.solid()
    top_strip.fill.fore_color.rgb = m['bg']
    top_strip.line.fill.background()

    # 头像圆
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + col_w/2 - 0.42), Inches(cy + 0.18), Inches(0.84), Inches(0.84))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = m['bg']
    avatar.line.fill.background()

    # 头像内首字
    initial = slide.shapes.add_textbox(Inches(cx + col_w/2 - 0.42), Inches(cy + 0.35), Inches(0.84), Inches(0.5))
    tf = initial.text_frame
    p = tf.paragraphs[0]
    p.text = m['name'][0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 名字 + 英文
    name_box = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(cy + 1.1), Inches(col_w - 0.24), Inches(0.38))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = m['name'] + '  ' + m['en']
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK

    # 职位
    title_box = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(cy + 1.46), Inches(col_w - 0.24), Inches(0.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = m['title']
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = m['bg']

    # 核心定位
    core_box = slide.shapes.add_textbox(Inches(cx + 0.12), Inches(cy + 1.76), Inches(col_w - 0.24), Inches(0.38))
    tf = core_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = m['core']
    p.font.size = Pt(9.5)
    p.font.color.rgb = GOLD
    p.font.bold = True

    # 分隔线
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx + 0.12), Inches(cy + 2.17), Inches(col_w - 0.24), Pt(0.75))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    sep.line.fill.background()

    # 要点
    for j, bullet in enumerate(m['bullets']):
        by = cy + 2.25 + j * 0.37
        b_box = slide.shapes.add_textbox(Inches(cx + 0.15), Inches(by), Inches(col_w - 0.3), Inches(0.37))
        tf = b_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(8.5)
        p.font.color.rgb = TEXT_GRAY

    # 底部能力标签
    tag_y = cy + col_h - 0.48
    tag_gap = col_w / 3
    for k, tag in enumerate(m['comp']):
        tx = cx + k * tag_gap + 0.08
        tag_box = slide.shapes.add_textbox(Inches(tx), Inches(tag_y), Inches(tag_gap - 0.1), Inches(0.38))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(7.5)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

# ========== 底部总结 ==========
sum_y = 6.52

bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(sum_y - 0.05), Inches(12.4), Pt(0.75))
bottom_line.fill.solid()
bottom_line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
bottom_line.line.fill.background()

dims = [
    ('产业', '24年酒店经验', BLUE_DARK),
    ('技术', '10年+AI研发+工程落地', BLUE_MID),
    ('战略', 'PE/VC+IPO+生态建设', BLUE_LIGHT),
]

for i, (dim, desc, color) in enumerate(dims):
    bx = 0.5 + i * 4.2
    dim_box = slide.shapes.add_textbox(Inches(bx), Inches(sum_y + 0.08), Inches(1.2), Inches(0.38))
    tf = dim_box.text_frame
    p = tf.paragraphs[0]
    p.text = dim
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color

    desc_box = slide.shapes.add_textbox(Inches(bx + 1.25), Inches(sum_y + 0.12), Inches(2.8), Inches(0.35))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_GRAY

    if i < 2:
        vbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bx + 4.0), Inches(sum_y + 0.1), Pt(1), Inches(0.3))
        vbar.fill.solid()
        vbar.fill.fore_color.rgb = LINE_COLOR
        vbar.line.fill.background()

prs.save(final_path)
print('Done')
