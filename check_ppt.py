from pptx import Presentation
import sys

file = "C:/Users/ericz/Desktop/张实项目总控/05-AHL-去中心化旅行平台/项目说明书/AHL-LLM去中心化旅行平台商业计划书V5.1(9).pptx"
prs = Presentation(file)

print(f"总页数: {len(prs.slides)}")
print(f"幻灯片尺寸: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")
print()

for i, slide in enumerate(prs.slides[:20], 1):
    print(f"=== 第{i}页 ===")
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()[:120]
            texts.append(text)
    if texts:
        for t in texts[:6]:
            print(f"  {t}")
    else:
        print("  [图片页/空白页]")
    print()

if len(prs.slides) > 20:
    print(f"... 还有 {len(prs.slides) - 20} 页")
