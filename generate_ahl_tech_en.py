# -*- coding: utf-8 -*-
"""
AHL-LLM Tech PPT Generator - English Version
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

sys.stdout.reconfigure(encoding='utf-8')

# Colors
DARK_BG = RGBColor(5, 15, 35)
PURPLE = RGBColor(124, 58, 237)
BLUE = RGBColor(59, 130, 246)
ORANGE = RGBColor(249, 115, 22)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 200)
CARD_BG = RGBColor(20, 35, 60)
RED_DARK = RGBColor(50, 30, 30)
PINK = RGBColor(236, 72, 153)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

def add_bar(slide):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    top.fill.solid()
    top.fill.fore_color.rgb = PURPLE
    top.line.fill.background()
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.04), prs.slide_width, Inches(0.04))
    bot.fill.solid()
    bot.fill.fore_color.rgb = ORANGE
    bot.line.fill.background()

def add_text(slide, text, left, top, width, height, size=20, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

def add_card(slide, left, top, width, height, title, lines, border_color=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = border_color
    
    add_text(slide, title, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.5), size=16, bold=True, color=border_color, align=PP_ALIGN.CENTER)
    
    content_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.6), width - Inches(0.3), height - Inches(0.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

def add_page_num(slide, num, total):
    add_text(slide, f"{num}/{total}", prs.slide_width - Inches(1), prs.slide_height - Inches(0.5), Inches(0.9), Inches(0.4), size=12, color=LIGHT, align=PP_ALIGN.RIGHT)

# Page 1: Cover
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)
add_bar(s1)

c1 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
c1.fill.solid()
c1.fill.fore_color.rgb = PURPLE
c1.line.fill.background()

c2 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
c2.fill.solid()
c2.fill.fore_color.rgb = BLUE
c2.line.fill.background()

add_text(s1, "AHL-LLM", 0, Inches(2), prs.slide_width, Inches(1), size=72, bold=True, color=WHITE)
add_text(s1, "Decentralized Travel Platform", 0, Inches(3), prs.slide_width, Inches(0.8), size=36, bold=True, color=BLUE)
add_text(s1, "Tech Project Overview", 0, Inches(4), prs.slide_width, Inches(0.6), size=28, color=LIGHT)
add_text(s1, "First LLM + Dual-AGENT Smart Hosting Platform", 0, Inches(5), prs.slide_width, Inches(0.5), size=20, color=ORANGE)

# Page 2: Overview
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
add_bar(s2)
add_page_num(s2, 2, 14)

add_text(s2, "Project Overview", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s2, "AHL: First LLM + Dual-AGENT Smart Hosting Platform", Inches(0.5), Inches(1.2), Inches(12), Inches(0.5), size=18, color=ORANGE)

items = [
    ("LLM Layer", ["DeepSeek/Qwen fine-tuned", "Hotel knowledge injected", "Multi-modal interaction"]),
    ("AGENT Layer", ["C-end AI Butler", "B-end AI Operations", "7x24h service"]),
    ("SKILL Layer", ["80+ pluggable SKILLs", "Like LEGO", "Fast adaptation"]),
    ("Vector Engine", ["<3s response", "95%+ accuracy", "Real-time inference"])
]

for i, (title, content) in enumerate(items):
    left = Inches(0.7) + i * Inches(3.1)
    add_card(s2, left, Inches(2), Inches(2.9), Inches(3.5), title, content, PURPLE if i%2==0 else BLUE)

# Page 3: Architecture
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
add_bar(s3)
add_page_num(s3, 3, 14)

add_text(s3, "4-Layer Tech Architecture", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

layers = [
    ("Layer 1: LLM", PURPLE, ["DeepSeek/Qwen fine-tuned", "12,000+ KB PP&SOP knowledge"]),
    ("Layer 2: AGENT", BLUE, ["C-end AI Butler - Guest Service", "B-end AI Operations - Autonomous"]),
    ("Layer 3: SKILL", ORANGE, ["80+ pluggable SKILLs", "Like LEGO - Fast adaptation"]),
    ("Layer 4: Vector", PINK, ["<3s response speed", "95%+ matching accuracy"])
]

for i, (title, color, content) in enumerate(layers):
    top = Inches(1.3) + i * Inches(1.4)
    layer_bg = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12.3), Inches(1.3))
    layer_bg.fill.solid()
    layer_bg.fill.fore_color.rgb = CARD_BG
    layer_bg.line.color.rgb = color
    
    add_text(s3, title, Inches(0.7), top + Inches(0.1), Inches(4), Inches(0.5), size=18, bold=True, color=color)
    content_box = s3.shapes.add_textbox(Inches(0.7), top + Inches(0.5), Inches(11), Inches(0.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = " | ".join(content)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

# Page 4: C-end
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4)
add_bar(s4)
add_page_num(s4, 4, 14)

add_text(s4, "C-end AI Butler", 0, Inches(0.3), prs.slide_width, Inches(0.6), size=40, bold=True)
add_text(s4, "7x24h Smart Travel Assistant", 0, Inches(1), Inches(12), Inches(0.4), size=20, color=ORANGE)

features = [
    ("Smart Booking", ["Natural language", "3-min booking", "No price compare"]),
    ("Trip Planning", ["Personalized recommend", "Real-time adjust", "Travel advisor"]),
    ("Room Service", ["Voice control", "One-click service", "Full access"]),
    ("Local Guide", ["Food recommend", "Attractions", "Local experience"]),
    ("Membership", ["Points mgmt", "Offers push", "Personal care"])
]

for i, (title, content) in enumerate(features):
    left = Inches(0.4) + i * Inches(2.5)
    add_card(s4, left, Inches(1.6), Inches(2.4), Inches(4.5), title, content, PURPLE if i%2==0 else BLUE)

# Page 5: B-end
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5)
add_bar(s5)
add_page_num(s5, 5, 14)

add_text(s5, "B-end AI Operations Officer", 0, Inches(0.3), prs.slide_width, Inches(0.6), size=40, bold=True)
add_text(s5, "Autonomous Operations Brain", 0, Inches(1), Inches(12), Inches(0.4), size=20, color=ORANGE)

b_features = [
    ("Revenue MGMT", ["Dynamic pricing", "RevPAR +15-30%", "Inventory opt"]),
    ("Channel OPS", ["OTA optimization", "Private traffic", "Commission -50%+"]),
    ("Guest Service", ["AI chatbot", "Complaint handling", "Staff -30%"]),
    ("Marketing", ["AI content gen", "Cost -40%", "Social media"]),
    ("Data Analytics", ["Reports", "Competitor monitoring", "Decision support"])
]

for i, (title, content) in enumerate(b_features):
    left = Inches(0.4) + i * Inches(2.5)
    add_card(s5, left, Inches(1.6), Inches(2.4), Inches(4.5), title, content, ORANGE if i%2==0 else PURPLE)

# Page 6: SKILLs
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6)
add_bar(s6)
add_page_num(s6, 6, 14)

add_text(s6, "Pluggable SKILLs System", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s6, "Like LEGO - Fast adaptation to different business types", 0, Inches(1.1), Inches(12), Inches(0.4), size=18, color=ORANGE)

add_text(s6, "C-end AI Butler SKILLs", Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.4), size=20, bold=True, color=BLUE)
add_text(s6, "Room 7 | F&B 13 | Banquet 6 | Front 6 | 4th Space 5", Inches(0.5), Inches(2.2), Inches(6), Inches(1.5), size=14, color=WHITE)

add_text(s6, "B-end AI Operations SKILLs", Inches(7), Inches(1.7), Inches(5.5), Inches(0.4), size=20, bold=True, color=ORANGE)
add_text(s6, "Revenue 5 | Front 5 | Room 5 | F&B 8 | MKT 5 | B2B 5 | Finance 5 | Energy 4", Inches(7), Inches(2.2), Inches(6), Inches(1.5), size=14, color=WHITE)

metrics = [("80+", "SKILLs"), ("95%+", "Accuracy"), ("<3s", "Response"), ("7x24h", "Always-on")]
for i, (num, label) in enumerate(metrics):
    left = Inches(1) + i * Inches(3)
    add_text(s6, num, left, Inches(4), Inches(2.5), Inches(1), size=48, bold=True, color=PURPLE if i%2==0 else BLUE, align=PP_ALIGN.CENTER)
    add_text(s6, label, left, Inches(5), Inches(2.5), Inches(0.5), size=16, color=LIGHT, align=PP_ALIGN.CENTER)

# Page 7: Knowledge Base
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7)
add_bar(s7)
add_page_num(s7, 7, 14)

add_text(s7, "PP&SOP Knowledge Base", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s7, "Core Data Foundation of AHL", 0, Inches(1.1), Inches(12), Inches(0.4), size=18, color=ORANGE)

kb_items = [
    ("Industry Knowledge", ["77 docs, 2000+KB", "12-layer architecture", "10+ brand standards"]),
    ("Operations SOP", ["Front/Room/F&B/MKT", "Procure/QC/Engineering", "Finance/HR/Admin"]),
    ("Scenario Knowledge", ["12 customer types", "10 product services", "10 operation modes"]),
    ("Vector Framework", ["Universal dimensions", "Specific core", "Real-time dynamics"])
]

for i, (title, content) in enumerate(kb_items):
    left = Inches(0.5) + i * Inches(3.2)
    add_card(s7, left, Inches(1.8), Inches(3), Inches(4.5), title, content, PURPLE if i%2==0 else BLUE)

# Page 8: Comparison
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8)
add_bar(s8)
add_page_num(s8, 8, 14)

add_text(s8, "Traditional vs AHL", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

left_bg = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.5))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = RED_DARK
left_bg.line.color.rgb = RGBColor(200, 100, 100)

add_text(s8, "Traditional OTA", Inches(0.5), Inches(1.5), Inches(5.4), Inches(0.5), size=24, bold=True, color=RGBColor(255, 150, 150), align=PP_ALIGN.CENTER)

left_items = ["X 15-25% high commission", "X Data silos", "X Passive response", "X Single function", "X Hotels = workers"]
for i, item in enumerate(left_items):
    add_text(s8, item, Inches(0.8), Inches(2.2) + i * Inches(0.8), Inches(5), Inches(0.7), size=16, color=RGBColor(255, 200, 200))

add_text(s8, "VS", Inches(6.3), Inches(3.5), Inches(0.8), Inches(0.8), size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

right_bg = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.8), Inches(5.5))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = RGBColor(20, 40, 60)
right_bg.line.color.rgb = BLUE

add_text(s8, "AHL Decentralized", Inches(7.4), Inches(1.5), Inches(5.4), Inches(0.5), size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

right_items = ["V Only 2% tech fee", "V Data connected", "V 7x24h active", "V 80+ SKILLs", "V Hotel autonomy"]
for i, item in enumerate(right_items):
    add_text(s8, item, Inches(7.6), Inches(2.2) + i * Inches(0.8), Inches(5), Inches(0.7), size=16, color=WHITE)

# Page 9: Tech Flow
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9)
add_bar(s9)
add_page_num(s9, 9, 14)

add_text(s9, "AHL Tech Flow", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

steps = [
    ("1", "User Input", "Natural language"),
    ("2", "NLP Processing", "Intent recognition"),
    ("3", "Vector Match", "SKILL retrieval"),
    ("4", "AGENT Execute", "Butler/Ops"),
    ("5", "Output", "Learning feedback")
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.5) + i * Inches(2.5)
    
    circle = s9.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(1.8), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PURPLE if i%2==0 else BLUE
    circle.line.fill.background()
    
    add_text(s9, num, left + Inches(0.8), Inches(1.9), Inches(0.7), Inches(0.5), size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_card(s9, left, Inches(2.7), Inches(2.3), Inches(2.5), title, [desc], PURPLE if i%2==0 else BLUE)
    
    if i < 4:
        arrow = s9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.35), Inches(3.8), Inches(0.15), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# Page 10: Metrics
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10)
add_bar(s10)
add_page_num(s10, 10, 14)

add_text(s10, "Core Performance Metrics", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

metrics = [
    ("<3s", "Response", PURPLE),
    ("95%+", "Accuracy", BLUE),
    ("7x24h", "Always-on", ORANGE),
    ("80+", "SKILLs", PINK),
    ("15-30%", "RevPAR+", PURPLE),
    ("50%+", "Commission-", BLUE)
]

for i, (num, label, color) in enumerate(metrics):
    row = i // 3
    col = i % 3
    left = Inches(1.5) + col * Inches(3.5)
    top = Inches(1.8) + row * Inches(2.5)
    
    card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3), Inches(2))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = color
    
    add_text(s10, num, left, top + Inches(0.3), Inches(3), Inches(1), size=48, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s10, label, left, top + Inches(1.3), Inches(3), Inches(0.5), size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Page 11: Roadmap
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11)
add_bar(s11)
add_page_num(s11, 11, 14)

add_text(s11, "Implementation Roadmap", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

phases = [
    ("Phase 1", "Core SKILL Dev", ["Room booking+Revenue", "Front+Room ops", "1-2 months"], PURPLE),
    ("Phase 2", "Scenario Expansion", ["F&B+Banquet+B2B", "MKT+Membership", "3-6 months"], BLUE),
    ("Phase 3", "Ecosystem Open", ["SKILL market", "3rd party access", "7-12 months"], ORANGE)
]

for i, (phase, title, content, color) in enumerate(phases):
    left = Inches(0.6) + i * Inches(4.2)
    
    card_bg = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(4.5))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = CARD_BG
    card_bg.line.color.rgb = color
    
    add_text(s11, phase, left, Inches(1.7), Inches(3.8), Inches(0.5), size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s11, title, left, Inches(2.3), Inches(3.8), Inches(0.5), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    content_box = s11.shapes.add_textbox(left + Inches(0.2), Inches(3), Inches(3.4), Inches(2.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(content):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    if i < 2:
        arrow = s11.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(3.85), Inches(3.5), Inches(0.35), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# Page 12: Team
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12)
add_bar(s12)
add_page_num(s12, 12, 14)

add_text(s12, "Core Team", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s12, "Golden Triangle: Industry + Tech + Capital", 0, Inches(1.1), Inches(12), Inches(0.4), size=18, color=ORANGE)

team = [
    ("Zhang Shi", "Project Lead", ["24 yrs hotel ind.", "Multi-group exec", "Paradigm designer"], PURPLE),
    ("Li Yuan", "CTO", ["AI PhD HUST", "10+ yrs AI R&D", "20+ AI projects"], BLUE),
    ("Chen Sixu", "CSO", ["Fortune 500", "PE/VC background", "IPO expert"], ORANGE)
]

for i, (name, role, desc, color) in enumerate(team):
    left = Inches(0.8) + i * Inches(4.2)
    
    card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(3.8), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = color
    
    avatar = s12.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.4), Inches(1.9), Inches(1), Inches(1))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = color
    avatar.line.fill.background()
    
    add_text(s12, name, left, Inches(3), Inches(3.8), Inches(0.5), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s12, role, left, Inches(3.5), Inches(3.8), Inches(0.4), size=14, color=color, align=PP_ALIGN.CENTER)
    
    content_box = s12.shapes.add_textbox(left + Inches(0.2), Inches(4), Inches(3.4), Inches(2.3))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT
        p.alignment = PP_ALIGN.CENTER

# Page 13: Summary
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s13)
add_bar(s13)
add_page_num(s13, 13, 14)

add_text(s13, "AHL Core Value", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

values = [
    ("For C-end", ["7x24h Butler", "Experience innov.", "Direct connect"], BLUE),
    ("For B-end", ["Autonomous ops", "Cost efficiency", "OTA-free"], PURPLE),
    ("For Platform", ["2% vs 15%", "New paradigm", "Tech revolution"], ORANGE),
    ("Tech Barrier", ["LLM+KB+SKILLs", "Trinity", "Moat"], PINK)
]

for i, (title, content, color) in enumerate(values):
    left = Inches(0.5) + i * Inches(3.2)
    
    card_bg = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.9), Inches(4.5))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = CARD_BG
    card_bg.line.color.rgb = color
    
    add_text(s13, title, left, Inches(1.8), Inches(2.9), Inches(0.6), size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
    
    content_box = s13.shapes.add_textbox(left + Inches(0.2), Inches(2.6), Inches(2.5), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(content):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# Page 14: Contact
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s14)
add_bar(s14)

c1 = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
c1.fill.solid()
c1.fill.fore_color.rgb = PURPLE
c1.line.fill.background()

c2 = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
c2.fill.solid()
c2.fill.fore_color.rgb = BLUE
c2.line.fill.background()

add_text(s14, "Welcome to AI New Era", 0, Inches(1.5), prs.slide_width, Inches(0.8), size=44, color=WHITE)
add_text(s14, "Contact Us", 0, Inches(2.8), prs.slide_width, Inches(0.6), size=32, bold=True, color=ORANGE)

contacts = ["Zhang Shi (Eric)", "Project Lead", "17760348653", "ericzhangshi@163.com", "AHL Team"]
for i, info in enumerate(contacts):
    add_text(s14, info, 0, Inches(4) + i * Inches(0.5), prs.slide_width, Inches(0.5), size=20, color=WHITE)

# Save
output = r"C:\Users\ericz\.openclaw\workspace\ppt_output\AHL-Tech\AHL_Tech_Project.pptx"
os.makedirs(os.path.dirname(output), exist_ok=True)
prs.save(output)
print(f"Generated: {output}")
print(f"Total pages: {len(prs.slides)}")
