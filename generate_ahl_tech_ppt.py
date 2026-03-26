# -*- coding: utf-8 -*-
"""
AHL-LLM去中心化旅行平台 - 技术项目说明PPT生成器
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 设置UTF-8
import sys
sys.stdout.reconfigure(encoding='utf-8')

# 颜色定义
GRADIENT_DARK = RgbColor(10, 22, 40)      # 深蓝黑
GRADIENT_PURPLE = RgbColor(124, 58, 237)   # 紫色
GRADIENT_ORANGE = RgbColor(249, 115, 22)    # 橙色
GRADIENT_BLUE = RgbColor(59, 130, 246)      # 蓝色
GRADIENT_PINK = RgbColor(236, 72, 153)     # 粉色
WHITE = RgbColor(255, 255, 255)
LIGHT_GRAY = RgbColor(200, 200, 200)
DARK_TEXT = RgbColor(30, 30, 30)

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

def add_title_shape(slide, text, left, top, width, height, font_size=44, bold=True, color=WHITE):
    """添加标题形状"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_subtitle_shape(slide, text, left, top, width, height, font_size=24, color=LIGHT_GRAY):
    """添加副标题形状"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_content_text(slide, lines, left, top, width, height, font_size=18, color=WHITE, align=PP_ALIGN.LEFT):
    """添加内容文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(12)
    return shape

def add_card(slide, left, top, width, height, title, content_lines):
    """添加卡片形状"""
    # 卡片背景
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(30, 40, 80)  # 深蓝卡片
    shape.line.color.rgb = GRADIENT_BLUE
    
    # 标题
    title_shape = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.5))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRADIENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 内容
    content_shape = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.6), width - Inches(0.3), height - Inches(0.7))
    tf = content_shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)
    
    return shape

def add_background(slide):
    """添加渐变背景"""
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(5, 15, 35)  # 深蓝黑背景
    background.line.fill.background()
    # 将背景移到最底层
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_decorative_elements(slide):
    """添加装饰元素"""
    # 顶部渐变条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GRADIENT_PURPLE
    top_bar.line.fill.background()
    
    # 底部渐变条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.05), prs.slide_width, Inches(0.05))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = GRADIENT_ORANGE
    bottom_bar.line.fill.background()

def add_page_number(slide, num, total):
    """添加页码"""
    shape = slide.shapes.add_textbox(prs.slide_width - Inches(1), prs.slide_height - Inches(0.5), Inches(0.9), Inches(0.4))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT

# ========== 第1页：封面 ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_background(slide1)
add_decorative_elements(slide1)

# 装饰圆形
circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
circle1.fill.solid()
circle1.fill.fore_color.rgb = GRADIENT_PURPLE
circle1.line.fill.background()

circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
circle2.fill.solid()
circle2.fill.fore_color.rgb = GRADIENT_BLUE
circle2.line.fill.background()

add_title_shape(slide1, "AHL-LLM", Inches(0), Inches(2), prs.slide_width, Inches(1), font_size=72, bold=True, color=WHITE)
add_title_shape(slide1, "去中心化旅行平台", Inches(0), Inches(3), prs.slide_width, Inches(0.8), font_size=48, bold=True, color=GRADIENT_BLUE)
add_subtitle_shape(slide1, "技术项目说明", Inches(0), Inches(4), prs.slide_width, Inches(0.6), font_size=28, color=LIGHT_GRAY)
add_subtitle_shape(slide1, "住宿业首个"大模型+双AGENT"智能托管平台", Inches(0), Inches(5), prs.slide_width, Inches(0.5), font_size=20, color=GRADIENT_ORANGE)
add_subtitle_shape(slide1, "基于DeepSeek/Qwen开源模型微调", Inches(0), Inches(5.8), prs.slide_width, Inches(0.4), font_size=16, color=LIGHT_GRAY)

# ========== 第2页：项目概述 ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2)
add_decorative_elements(slide2)
add_page_number(slide2, 2, 14)

add_title_shape(slide2, "项目概述", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)
add_subtitle_shape(slide2, "一句话定位：AHL是住宿业首个"大模型+双AGENT"智能托管平台", Inches(0.5), Inches(1.2), Inches(12), Inches(0.5), font_size=18, color=GRADIENT_ORANGE)

# 技术栈卡片
tech_items = [
    ("🤖 大模型层", ["基于DeepSeek/Qwen开源模型微调", "注入酒店行业知识库", "多模态交互(文本/语音/图像)"]),
    ("🎯 AGENT层", ["C端AI管家 | B端AI运营官", "7×24h全天候服务", "自主决策与执行"]),
    ("🧩 SKILL层", ["80+可插拔专业SKILLs", "覆盖前厅/客房/餐饮/营销", "灵活组合像乐高"]),
    ("⚡ 向量引擎", ["<3秒响应速度", "95%+匹配准确率", "实时推理与决策"])
]

card_width = Inches(2.8)
card_height = Inches(3)
start_left = Inches(0.7)
gap = Inches(0.3)

for i, (title, content) in enumerate(tech_items):
    left = start_left + i * (card_width + gap)
    card = add_card(slide2, left, Inches(2), card_width, card_height, title, content)

# ========== 第3页：四层技术架构 ==========
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3)
add_decorative_elements(slide3)
add_page_number(slide3, 3, 14)

add_title_shape(slide3, "四层技术架构", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)

# 架构层
layers = [
    ("第一层：大模型层 (AHL-LLM)", GRADIENT_PURPLE, ["基于DeepSeek/Qwen开源模型微调", "注入12,000+字PP&SOP知识库", "多模态交互(文本/语音/图像)"]),
    ("第二层：AGENT层", GRADIENT_BLUE, ["C端AI管家 - 对客服务智能助手", "B端AI运营官 - 自主运营大脑"]),
    ("第三层：SKILL层", GRADIENT_ORANGE, ["80+可插拔专业SKILLs", "像乐高一样灵活组合", "快速适配不同业态"]),
    ("第四层：向量匹配引擎", GRADIENT_PINK, ["<3秒响应速度", "95%+匹配准确率", "实时推理"])
]

layer_height = Inches(1.3)
for i, (title, color, content) in enumerate(layers):
    top = Inches(1.3) + i * (layer_height + Inches(0.15))
    
    # 层背景
    layer_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12.3), layer_height)
    layer_bg.fill.solid()
    layer_bg.fill.fore_color.rgb = RgbColor(20, 30, 60)
    layer_bg.line.color.rgb = color
    
    # 层标题
    title_box = slide3.shapes.add_textbox(Inches(0.7), top + Inches(0.1), Inches(4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    
    # 内容
    content_box = slide3.shapes.add_textbox(Inches(0.7), top + Inches(0.5), Inches(11), Inches(0.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(content):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

# ========== 第4页：C端AI管家 ==========
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4)
add_decorative_elements(slide4)
add_page_number(slide4, 4, 14)

add_title_shape(slide4, "C端AI管家", Inches(0), Inches(0.3), prs.slide_width, Inches(0.6), font_size=40)
add_subtitle_shape(slide4, "消费者的7×24小时智能旅行助手", Inches(0), Inches(1), Inches(12), Inches(0.4), font_size=20, color=GRADIENT_ORANGE)

# 功能卡片
features = [
    ("智能预订", "自然语言交互\n3分钟完成预订\n无需比价"),
    ("行程规划", "基于用户偏好\n个性化行程推荐\n实时调整优化"),
    ("客房服务", "语音控制设备\n服务请求一键处理\n酒店服务全触达"),
    ("本地向导", "美食推荐\n景点介绍\n交通规划\n在地体验"),
    ("会员运营", "积分管理\n优惠推送\n个性化关怀")
]

card_width = Inches(2.3)
card_height = Inches(4)
start_left = Inches(0.4)
gap = Inches(0.2)

for i, (title, content) in enumerate(features):
    left = start_left + i * (card_width + gap)
    lines = content.split('\n')
    card = add_card(slide4, left, Inches(1.6), card_width, card_height, title, lines)

# ========== 第5页：B端AI运营官 ==========
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5)
add_decorative_elements(slide5)
add_page_number(slide5, 5, 14)

add_title_shape(slide5, "B端AI运营官", Inches(0), Inches(0.3), prs.slide_width, Inches(0.6), font_size=40)
add_subtitle_shape(slide5, "酒店/民宿的自主运营大脑", Inches(0), Inches(1), Inches(12), Inches(0.4), font_size=20, color=GRADIENT_ORANGE)

features = [
    ("收益管理", "动态定价算法\nRevPAR提升15-30%\n库存优化"),
    ("渠道运营", "OTA优化\n私域流量导流\n降低佣金50%+"),
    ("客户服务", "智能客服系统\n投诉处理\n好评引导\n人工减少30%"),
    ("营销推广", "AI内容生成\n社媒运营\n获客成本降低40%"),
    ("数据分析", "经营报表\n竞品监控\n决策支持")
]

for i, (title, content) in enumerate(features):
    left = start_left + i * (card_width + gap)
    lines = content.split('\n')
    card = add_card(slide5, left, Inches(1.6), card_width, card_height, title, lines)

# ========== 第6页：SKILLs体系 ==========
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6)
add_decorative_elements(slide6)
add_page_number(slide6, 6, 14)

add_title_shape(slide6, "SKILLs可插拔体系", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)
add_subtitle_shape(slide6, "像乐高一样灵活组合，快速适配不同业态", Inches(0), Inches(1.1), Inches(12), Inches(0.4), font_size=18, color=GRADIENT_ORANGE)

# C端SKILL卡片
c_end_title = slide6.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.4))
tf = c_end_title.text_frame
p = tf.paragraphs[0]
p.text = "C端AI管家 SKILL群"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GRADIENT_BLUE

c_items = "客房服务7个 | 餐饮服务13个 | 宴会服务6个 | 前厅礼宾6个 | 第四空间5个"
c_content = slide6.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6), Inches(1.5))
tf = c_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = c_items
p.font.size = Pt(14)
p.font.color.rgb = WHITE

# B端SKILL卡片
b_end_title = slide6.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.5), Inches(0.4))
tf = b_end_title.text_frame
p = tf.paragraphs[0]
p.text = "B端AI运营官 SKILL群"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GRADIENT_ORANGE

b_items = "收益管理5个 | 前厅运营5个 | 客房管理5个 | 餐饮运营8个 | 市场营销5个 | 闲置空间6个 | B2B市场5个 | 财务5个 | 能耗4个"
b_content = slide6.shapes.add_textbox(Inches(7), Inches(2.2), Inches(6), Inches(1.5))
tf = b_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = b_items
p.font.size = Pt(14)
p.font.color.rgb = WHITE

# 核心指标
metrics = [
    ("80+", "细分SKILLs"),
    ("95%+", "匹配准确率"),
    ("<3秒", "响应速度"),
    ("7×24h", "全天候服务")
]

for i, (num, label) in enumerate(metrics):
    left = Inches(1) + i * Inches(3)
    
    # 数字
    num_box = slide6.shapes.add_textbox(left, Inches(4), Inches(2.5), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GRADIENT_PURPLE
    p.alignment = PP_ALIGN.CENTER
    
    # 标签
    label_box = slide6.shapes.add_textbox(left, Inches(5), Inches(2.5), Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

# ========== 第7页：PP&SOP知识库 ==========
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7)
add_decorative_elements(slide7)
add_page_number(slide7, 7, 14)

add_title_shape(slide7, "PP&SOP知识库", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)
add_subtitle_shape(slide7, "AHL的核心数据底座", Inches(0), Inches(1.1), Inches(12), Inches(0.4), font_size=18, color=GRADIENT_ORANGE)

# 知识库内容
knowledge = [
    ("行业知识库", ["77个文档，2000+KB", "12层全景架构体系", "万豪/希尔顿/洲际/华住等10+品牌", "消防/食品/卫生/劳动/工商税务"]),
    ("运营SOP", ["前厅/客房/餐饮/营销", "采购/质检/工程/安保", "财务/人力/总经办", "八部一室/扁平化/酒管公司"]),
    ("场景知识", ["12种客户群体", "10种产品服务", "10种运营模式", "10种收费模式"]),
    ("向量框架", ["通用维度基底", "特异化核心", "实时动态系数", "AHL事实双向向量"])
]

for i, (title, content) in enumerate(knowledge):
    left = Inches(0.5) + i * Inches(3.2)
    card = add_card(slide7, left, Inches(1.8), Inches(3), Inches(4.5), title, content)

# ========== 第8页：传统方案 vs AHL ==========
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8)
add_decorative_elements(slide8)
add_page_number(slide8, 8, 14)

add_title_shape(slide8, "传统方案 vs AHL", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)

# VS 布局
# 左侧：传统
left_bg = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.5))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = RgbColor(50, 30, 30)
left_bg.line.color.rgb = RgbColor(200, 100, 100)

left_title = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.4), Inches(0.5))
tf = left_title.text_frame
p = tf.paragraphs[0]
p.text = "传统OTA平台"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RgbColor(255, 150, 150)
p.alignment = PP_ALIGN.CENTER

left_items = ["❌ 15-25%高佣金剥削", "❌ 数据孤岛，无法协同", "❌ 被动响应，效率低下", "❌ 单一功能，无法适配", "❌ 酒店沦为打工者"]
for i, item in enumerate(left_items):
    item_box = slide8.shapes.add_textbox(Inches(0.8), Inches(2.2) + i * Inches(0.8), Inches(5), Inches(0.7))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(255, 200, 200)

# 中间 VS
vs_box = slide8.shapes.add_textbox(Inches(6.3), Inches(3.5), Inches(0.8), Inches(0.8))
tf = vs_box.text_frame
p = tf.paragraphs[0]
p.text = "VS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 右侧：AHL
right_bg = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.8), Inches(5.5))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = RgbColor(20, 40, 60)
right_bg.line.color.rgb = GRADIENT_BLUE

right_title = slide8.shapes.add_textbox(Inches(7.4), Inches(1.5), Inches(5.4), Inches(0.5))
tf = right_title.text_frame
p = tf.paragraphs[0]
p.text = "AHL去中心化平台"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GRADIENT_BLUE
p.alignment = PP_ALIGN.CENTER

right_items = ["✅ 仅2%技术服务费", "✅ 数据互联，智能协同", "✅ 7×24h主动服务", "✅ 80+SKILLs灵活组合", "✅ 酒店自主运营"]
for i, item in enumerate(right_items):
    item_box = slide8.shapes.add_textbox(Inches(7.6), Inches(2.2) + i * Inches(0.8), Inches(5), Inches(0.7))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

# ========== 第9页：技术流程 ==========
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9)
add_decorative_elements(slide9)
add_page_number(slide9, 9, 14)

add_title_shape(slide9, "AHL技术流程", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)

# 流程步骤
steps = [
    ("1", "用户需求输入", "自然语言/语音/图像"),
    ("2", "自然语言处理", "意图识别 & 实体提取"),
    ("3", "向量匹配引擎", "SKILL检索 & 选择"),
    ("4", "AGENT执行", "C端管家/B端运营官"),
    ("5", "结果输出", "学习反馈 & 持续优化")
]

step_width = Inches(2.3)
step_height = Inches(3)
start_left = Inches(0.5)
gap = Inches(0.3)

for i, (num, title, content) in enumerate(steps):
    left = start_left + i * (step_width + gap)
    
    # 步骤圆圈
    circle = slide9.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(1.8), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GRADIENT_PURPLE if i % 2 == 0 else GRADIENT_BLUE
    circle.line.fill.background()
    
    # 步骤编号
    num_box = slide9.shapes.add_textbox(left + Inches(0.8), Inches(1.9), Inches(0.7), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 步骤卡片
    card = add_card(slide9, left, Inches(2.7), step_width, step_height, title, [content])
    
    # 连接箭头
    if i < len(steps) - 1:
        arrow = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + step_width + Inches(0.05), Inches(4.1), Inches(0.2), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRADIENT_ORANGE
        arrow.line.fill.background()

# ========== 第10页：核心性能指标 ==========
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10)
add_decorative_elements(slide10)
add_page_number(slide10, 10, 14)

add_title_shape(slide10, "核心性能指标", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)

metrics = [
    ("<3秒", "响应速度", GRADIENT_PURPLE),
    ("95%+", "匹配准确率", GRADIENT_BLUE),
    ("7×24h", "全天候服务", GRADIENT_ORANGE),
    ("80+", "专业SKILLs", GRADIENT_PINK),
    ("15-30%", "RevPAR提升", GRADIENT_PURPLE),
    ("50%+", "佣金成本降低", GRADIENT_BLUE)
]

# 2行3列布局
for i, (num, label, color) in enumerate(metrics):
    row = i // 3
    col = i % 3
    left = Inches(1.5) + col * Inches(3.5)
    top = Inches(1.8) + row * Inches(2.5)
    
    # 指标卡片
    card_bg = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3), Inches(2))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RgbColor(20, 35, 60)
    card_bg.line.color.rgb = color
    
    # 数字
    num_box = slide10.shapes.add_textbox(left, top + Inches(0.3), Inches(3), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    # 标签
    label_box = slide10.shapes.add_textbox(left, top + Inches(1.3), Inches(3), Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ========== 第11页：实施路径 ==========
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11)
add_decorative_elements(slide11)
add_page_number(slide11, 11, 14)

add_title_shape(slide11, "实施路径", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)

phases = [
    ("Phase 1", "核心SKILL开发", ["客房预订+收益管理", "前厅运营+客房管理", "1-2个月", GRADIENT_PURPLE]),
    ("Phase 2", "场景SKILL扩展", ["餐饮+宴会+B2B", "市场营销+会员运营", "3-6个月", GRADIENT_BLUE]),
    ("Phase 3", "生态开放", ["SKILL市场+第三方接入", "开放API", "7-12个月", GRADIENT_ORANGE])
]

phase_width = Inches(3.8)
phase_height = Inches(4.5)
start_left = Inches(0.6)
gap = Inches(0.4)

for i, (phase, title, content, color) in enumerate(phases):
    left = start_left + i * (phase_width + gap)
    
    # 阶段卡片
    card_bg = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), phase_width, phase_height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RgbColor(15, 30, 50)
    card_bg.line.color.rgb = color
    
    # 阶段标签
    phase_box = slide11.shapes.add_textbox(left, Inches(1.7), phase_width, Inches(0.5))
    tf = phase_box.text_frame
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    # 标题
    title_box = slide11.shapes.add_textbox(left, Inches(2.3), phase_width, Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 内容
    content_box = slide11.shapes.add_textbox(left + Inches(0.2), Inches(3), phase_width - Inches(0.4), Inches(2.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(content[:3]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # 连接箭头
    if i < len(phases) - 1:
        arrow = slide11.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + phase_width + Inches(0.1), Inches(3.5), Inches(0.2), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRADIENT_ORANGE
        arrow.line.fill.background()

# ========== 第12页：核心技术团队 ==========
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide12)
add_decorative_elements(slide12)
add_page_number(slide12, 12, 14)

add_title_shape(slide12, "核心技术团队", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)
add_subtitle_shape(slide12, "黄金三角组合：产业+技术+资本", Inches(0), Inches(1.1), Inches(12), Inches(0.4), font_size=18, color=GRADIENT_ORANGE)

team = [
    ("张实", "项目总控/发起人", ["24年酒店业深耕", "北京第二外国语学院", "多家知名酒店集团高管", "住宿业范式革命总设计师"]),
    ("李源", "首席技术官 (CTO)", ["华中科技大学AI博士", "10年+AI技术研发", "20+AI工程落地", "AI架构与算法壁垒"]),
    ("陈思序", "首席战略与生态官 (CSO)", ["世界500强战略高管", "深耕PE/VC投资", "IPO全流程经验", "战略推动与生态建设"])
]

card_width = Inches(3.8)
card_height = Inches(4.8)
start_left = Inches(0.8)
gap = Inches(0.5)

for i, (name, role, desc) in enumerate(team):
    left = start_left + i * (card_width + gap)
    
    # 卡片
    card_bg = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), card_width, card_height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RgbColor(20, 35, 60)
    card_bg.line.color.rgb = GRADIENT_BLUE if i == 1 else GRADIENT_PURPLE
    
    # 头像占位（圆形）
    avatar = slide12.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.4), Inches(1.9), Inches(1), Inches(1))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = GRADIENT_PURPLE if i == 0 else (GRADIENT_BLUE if i == 1 else GRADIENT_ORANGE)
    avatar.line.fill.background()
    
    # 名字
    name_box = slide12.shapes.add_textbox(left, Inches(3), card_width, Inches(0.5))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 角色
    role_box = slide12.shapes.add_textbox(left, Inches(3.5), card_width, Inches(0.4))
    tf = role_box.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(14)
    p.font.color.rgb = GRADIENT_BLUE if i == 1 else GRADIENT_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # 描述
    desc_box = slide12.shapes.add_textbox(left + Inches(0.2), Inches(4), card_width - Inches(0.4), Inches(2.3))
    tf = desc_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

# ========== 第13页：总结 ==========
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide13)
add_decorative_elements(slide13)
add_page_number(slide13, 13, 14)

add_title_shape(slide13, "AHL核心价值", Inches(0), Inches(0.3), prs.slide_width, Inches(0.8), font_size=40)

values = [
    ("对C端", "7×24h个性化AI管家\n革新旅行体验\n直连商家无中间商", GRADIENT_BLUE),
    ("对B端", "自主运营AGENT\n降本增效\n摆脱OTA高佣金依赖", GRADIENT_PURPLE),
    ("对平台", "去中心化2% vs OTA 15%\n构建交易新范式\n技术服务费革命", GRADIENT_ORANGE),
    ("技术壁垒", "大模型+知识库+SKILLs\n三位一体\n不可复制的技术护城河", GRADIENT_PINK)
]

card_width = Inches(2.9)
card_height = Inches(4.5)
start_left = Inches(0.5)
gap = Inches(0.3)

for i, (title, content, color) in enumerate(values):
    left = start_left + i * (card_width + gap)
    
    card_bg = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), card_width, card_height)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RgbColor(15, 30, 50)
    card_bg.line.color.rgb = color
    
    title_box = slide13.shapes.add_textbox(left, Inches(1.8), card_width, Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    content_box = slide13.shapes.add_textbox(left + Inches(0.2), Inches(2.6), card_width - Inches(0.4), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True
    lines = content.split('\n')
    for j, line in enumerate(lines):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# ========== 第14页：联系我们 ==========
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide14)
add_decorative_elements(slide14)
add_page_number(slide14, 14, 14)

# 装饰
circle1 = slide14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
circle1.fill.solid()
circle1.fill.fore_color.rgb = GRADIENT_PURPLE
circle1.line.fill.background()

circle2 = slide14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = GRADIENT_BLUE
circle2.line.fill.background()

add_title_shape(slide14, "开启住宿业的AI新纪元", Inches(0), Inches(1.5), prs.slide_width, Inches(0.8), font_size=44, color=WHITE)
add_subtitle_shape(slide14, "联系我们", Inches(0), Inches(2.8), prs.slide_width, Inches(0.6), font_size=32, color=GRADIENT_ORANGE)

contact_info = [
    "张实 (Eric Zhang)",
    "项目总控",
    "☎️ 17760348653",
    "📧 ericzhangshi@163.com",
    "🌐 AHL团队"
]

for i, info in enumerate(contact_info):
    info_box = slide14.shapes.add_textbox(Inches(0), Inches(4) + i * Inches(0.5), prs.slide_width, Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = info
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# 保存文件
output_path = r"C:\Users\ericz\.openclaw\workspace\ppt_output\AHL-Tech\AHL技术项目说明_v1.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"共 {len(prs.slides)} 页")
