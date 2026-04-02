# -*- coding: utf-8 -*-
import shutil, os, glob
from pathlib import Path
from pptx import Presentation

# Image directory
img_dir = Path(r"C:\Users\ericz\.openclaw\media\tool-image-generation")
out_dir = Path(r"C:\Users\ericz\.openclaw\workspace\ppt_output\ahl_roadshow\slides")
out_dir.mkdir(parents=True, exist_ok=True)

# Get all images in the directory
all_images = sorted(img_dir.glob("*.png"))
print(f"Found {len(all_images)} images")

# Copy as slide-01.png through slide-14.png
for i, img in enumerate(all_images, start=1):
    dst = out_dir / f"slide-{i:02d}.png"
    shutil.copy(img, dst)
    print(f"  slide-{i:02d}.png <- {img.name}")

# Build PPTX
prs = Presentation()
prs.slide_width = int(13.333 * 914400)   # 16:9 widescreen
prs.slide_height = int(7.5 * 914400)
blank_layout = prs.slide_layouts[6]

slide_files = sorted(out_dir.glob("slide-*.png"))
print(f"\nBuilding PPTX with {len(slide_files)} slides...")
for sf in slide_files:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(sf), 0, 0, width=prs.slide_width, height=prs.slide_height)
    print(f"  Added: {sf.name}")

pptx_path = Path(r"C:\Users\ericz\.openclaw\workspace\ppt_output\ahl_roadshow\AHL_商业计划书_v1.0.pptx")
prs.save(str(pptx_path))
print(f"\nSaved: {pptx_path}")
print(f"Total slides: {len(slide_files)}")
