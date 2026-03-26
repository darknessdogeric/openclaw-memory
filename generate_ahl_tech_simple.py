# -*- coding: utf-8 -*-
"""
AHL-LLM PPT Generator - Simple Version
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

# UTF-8 support
sys.stdout.reconfigure(encoding='utf-8')

# Colors
DARK_BG = RGBColor(5, 15, 35)
PURPLE = RGBColor(124, 58, 237)
BLUE = RGBColor(59, 130, 246)
ORANGE = RGBColor(249, 115, 22)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 200)
CARD_BG = RGBColor(20, 35, 60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

def add_bar(slide):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    top.fill.solid()
    top.fill.fore_color.rgb = PURPLE
    top.line.fill.background()
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.04), prs.slide_width, Inches(0.04))
    bot.fill.solid()
    bot.fill.fore_color.rgb = ORANGE
    bot.line.fill.background()

def add_text(slide, text, left, top, width, height, size=20, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

def add_card(slide, left, top, width, height, title, lines, border_color=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = border_color
    
    add_text(slide, title, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.5), size=16, bold=True, color=border_color, align=PP_ALIGN.CENTER)
    
    content_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.6), width - Inches(0.3), height - Inches(0.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

def add_page_num(slide, num, total):
    add_text(slide, f"{num}/{total}", prs.slide_width - Inches(1), prs.slide_height - Inches(0.5), Inches(0.9), Inches(0.4), size=12, color=LIGHT, align=PP_ALIGN.RIGHT)

# ========== Page 1: Cover ==========
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)
add_bar(s1)

c1 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
c1.fill.solid()
c1.fill.fore_color.rgb = PURPLE
c1.line.fill.background()

c2 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
c2.fill.solid()
c2.fill.fore_color.rgb = BLUE
c2.line.fill.background()

add_text(s1, "AHL-LLM", 0, Inches(2), prs.slide_width, Inches(1), size=72, bold=True, color=WHITE)
add_text(s1, "去中心化旅行平台", 0, Inches(3), prs.slide_width, Inches(0.8), size=48, bold=True, color=BLUE)
add_text(s1, "技术项目说�?, 0, Inches(4), prs.slide_width, Inches(0.6), size=28, color=LIGHT)
add_text(s1, "住宿业首个\"大模�?双AGENT\"智能托管平台", 0, Inches(5), prs.slide_width, Inches(0.5), size=20, color=ORANGE)

# ========== Page 2: Overview ==========
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
add_bar(s2)
add_page_num(s2, 2, 14)

add_text(s2, "项目概述", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s2, "一句话定位：AHL是住宿业首个\"大模�?双AGENT\"智能托管平台", Inches(0.5), Inches(1.2), Inches(12), Inches(0.5), size=18, color=ORANGE)

items = [
    ("大模型层", ["基于DeepSeek/Qwen微调", "注入酒店知识�?, "多模态交�?]),
    ("AGENT�?, ["C端AI管家", "B端AI运营�?, "7x24h服务"]),
    ("SKILL�?, ["80+可插拔SKILLs", "像乐高一样组�?, "快速适配业�?]),
    ("向量引擎", ["<3秒响�?, "95%+匹配准确�?, "实时推理决策"])
]

for i, (title, content) in enumerate(items):
    left = Inches(0.7) + i * Inches(3.1)
    add_card(s2, left, Inches(2), Inches(2.9), Inches(3.5), title, content, PURPLE if i%2==0 else BLUE)

# ========== Page 3: Architecture ==========
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
add_bar(s3)
add_page_num(s3, 3, 14)

add_text(s3, "四层技术架�?, 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

layers = [
    ("第一层：大模型层", PURPLE, ["基于DeepSeek/Qwen开源模型微�?, "注入12,000+字PP&SOP知识�?]),
    ("第二层：AGENT�?, BLUE, ["C端AI管家 - 对客服务智能助手", "B端AI运营�?- 自主运营大脑"]),
    ("第三层：SKILL�?, ORANGE, ["80+可插拔专业SKILLs", "像乐高一样灵活组合快速适配"]),
    ("第四层：向量匹配引擎", RGBColor(236, 72, 153), ["<3秒响应速度", "95%+匹配准确�?])
]

for i, (title, color, content) in enumerate(layers):
    top = Inches(1.3) + i * Inches(1.4)
    layer_bg = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12.3), Inches(1.3))
    layer_bg.fill.solid()
    layer_bg.fill.fore_color.rgb = CARD_BG
    layer_bg.line.color.rgb = color
    
    add_text(s3, title, Inches(0.7), top + Inches(0.1), Inches(4), Inches(0.5), size=18, bold=True, color=color)
    content_box = s3.shapes.add_textbox(Inches(0.7), top + Inches(0.5), Inches(11), Inches(0.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = " | ".join(content)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

# ========== Page 4: C-end ==========
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4)
add_bar(s4)
add_page_num(s4, 4, 14)

add_text(s4, "C端AI管家", 0, Inches(0.3), prs.slide_width, Inches(0.6), size=40, bold=True)
add_text(s4, "消费者的7x24小时智能旅行助手", 0, Inches(1), Inches(12), Inches(0.4), size=20, color=ORANGE)

features = [
    ("智能预订", ["自然语言交互", "3分钟完成预订", "无需比价"]),
    ("行程规划", ["基于偏好个性化", "实时调整优化", "专属旅行顾问"]),
    ("客房服务", ["语音控制设备", "服务一键触�?, "酒店服务全链�?]),
    ("本地向导", ["美食推荐", "景点介绍", "在地体验"]),
    ("会员运营", ["积分管理", "优惠推�?, "个性化关怀"])
]

for i, (title, content) in enumerate(features):
    left = Inches(0.4) + i * Inches(2.5)
    add_card(s4, left, Inches(1.6), Inches(2.4), Inches(4.5), title, content, PURPLE if i%2==0 else BLUE)

# ========== Page 5: B-end ==========
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5)
add_bar(s5)
add_page_num(s5, 5, 14)

add_text(s5, "B端AI运营�?, 0, Inches(0.3), prs.slide_width, Inches(0.6), size=40, bold=True)
add_text(s5, "酒店/民宿的自主运营大�?, 0, Inches(1), Inches(12), Inches(0.4), size=20, color=ORANGE)

b_features = [
    ("收益管理", ["动态定价算�?, "RevPAR提升15-30%", "库存优化"]),
    ("渠道运营", ["OTA优化", "私域流量导流", "降低佣金50%+"]),
    ("客户服务", ["智能客服系统", "投诉处理", "人工减少30%"]),
    ("营销推广", ["AI内容生成", "获客成本�?0%", "社媒运营"]),
    ("数据分析", ["经营报表", "竞品监控", "决策支持"])
]

for i, (title, content) in enumerate(b_features):
    left = Inches(0.4) + i * Inches(2.5)
    add_card(s5, left, Inches(1.6), Inches(2.4), Inches(4.5), title, content, ORANGE if i%2==0 else PURPLE)

# ========== Page 6: SKILLs ==========
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6)
add_bar(s6)
add_page_num(s6, 6, 14)

add_text(s6, "SKILLs可插拔体�?, 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s6, "像乐高一样灵活组合，快速适配不同业�?, 0, Inches(1.1), Inches(12), Inches(0.4), size=18, color=ORANGE)

# C-end & B-end sections
add_text(s6, "C端AI管家 SKILL�?, Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.4), size=20, bold=True, color=BLUE)
add_text(s6, "客房7 | 餐饮13 | 宴会6 | 前厅6 | 第四空间5", Inches(0.5), Inches(2.2), Inches(6), Inches(1.5), size=14, color=WHITE)

add_text(s6, "B端AI运营�?SKILL�?, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4), size=20, bold=True, color=ORANGE)
add_text(s6, "收益5 | 前厅5 | 客房5 | 餐饮8 | 营销5 | B2B5 | 财务5 | 能�?", Inches(7), Inches(2.2), Inches(6), Inches(1.5), size=14, color=WHITE)

# Metrics
metrics = [("80+", "细分SKILLs"), ("95%+", "匹配准确�?), ("<3�?, "响应速度"), ("7x24h", "全天�?)]
for i, (num, label) in enumerate(metrics):
    left = Inches(1) + i * Inches(3)
    add_text(s6, num, left, Inches(4), Inches(2.5), Inches(1), size=48, bold=True, color=PURPLE if i%2==0 else BLUE, align=PP_ALIGN.CENTER)
    add_text(s6, label, left, Inches(5), Inches(2.5), Inches(0.5), size=16, color=LIGHT, align=PP_ALIGN.CENTER)

# ========== Page 7: Knowledge Base ==========
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7)
add_bar(s7)
add_page_num(s7, 7, 14)

add_text(s7, "PP&SOP知识�?, 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s7, "AHL的核心数据底�?, 0, Inches(1.1), Inches(12), Inches(0.4), size=18, color=ORANGE)

kb_items = [
    ("行业知识�?, ["77个文档，2000+KB", "12层全景架�?, "10+品牌标准"]),
    ("运营SOP", ["前厅/客房/餐饮/营销", "采购/质检/工程", "财务/人力/总经�?]),
    ("场景知识", ["12种客户群�?, "10种产品服�?, "10种运营模�?]),
    ("向量框架", ["通用维度基底", "特异化核�?, "实时动态系�?])
]

for i, (title, content) in enumerate(kb_items):
    left = Inches(0.5) + i * Inches(3.2)
    add_card(s7, left, Inches(1.8), Inches(3), Inches(4.5), title, content, PURPLE if i%2==0 else BLUE)

# ========== Page 8: Comparison ==========
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8)
add_bar(s8)
add_page_num(s8, 8, 14)

add_text(s8, "传统方案 vs AHL", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

# Left: Traditional
left_bg = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.5))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = RGBColor(50, 30, 30)
left_bg.line.color.rgb = RGBColor(200, 100, 100)

add_text(s8, "传统OTA平台", Inches(0.5), Inches(1.5), Inches(5.4), Inches(0.5), size=24, bold=True, color=RGBColor(255, 150, 150), align=PP_ALIGN.CENTER)

left_items = ["X 15-25%高佣金剥�?, "X 数据孤岛无法协同", "X 被动响应效率�?, "X 单一功能无法适配", "X 酒店沦为打工�?]
for i, item in enumerate(left_items):
    add_text(s8, item, Inches(0.8), Inches(2.2) + i * Inches(0.8), Inches(5), Inches(0.7), size=16, color=RGBColor(255, 200, 200))

# VS
add_text(s8, "VS", Inches(6.3), Inches(3.5), Inches(0.8), Inches(0.8), size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Right: AHL
right_bg = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.8), Inches(5.5))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = RGBColor(20, 40, 60)
right_bg.line.color.rgb = BLUE

add_text(s8, "AHL去中心化平台", Inches(7.4), Inches(1.5), Inches(5.4), Inches(0.5), size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

right_items = ["V  �?%技术服务费", "V 数据互联智能协同", "V 7x24h主动服务", "V 80+SKILLs灵活组合", "V 酒店自主运营"]
for i, item in enumerate(right_items):
    add_text(s8, item, Inches(7.6), Inches(2.2) + i * Inches(0.8), Inches(5), Inches(0.7), size=16, color=WHITE)

# ========== Page 9: Tech Flow ==========
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9)
add_bar(s9)
add_page_num(s9, 9, 14)

add_text(s9, "AHL技术流�?, 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

steps = [
    ("1", "用户需�?, "自然语言输入"),
    ("2", "NLP处理", "意图识别"),
    ("3", "向量匹配", "SKILL检�?),
    ("4", "AGENT执行", "管家/运营�?),
    ("5", "结果输出", "学习反馈")
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.5) + i * Inches(2.5)
    
    circle = s9.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(1.8), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PURPLE if i%2==0 else BLUE
    circle.line.fill.background()
    
    add_text(s9, num, left + Inches(0.8), Inches(1.9), Inches(0.7), Inches(0.5), size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_card(s9, left, Inches(2.7), Inches(2.3), Inches(2.5), title, [desc], PURPLE if i%2==0 else BLUE)
    
    if i < 4:
        arrow = s9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.35), Inches(3.8), Inches(0.15), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# ========== Page 10: Metrics ==========
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10)
add_bar(s10)
add_page_num(s10, 10, 14)

add_text(s10, "核心性能指标", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

metrics = [
    ("<3�?, "响应速度", PURPLE),
    ("95%+", "匹配准确�?, BLUE),
    ("7x24h", "全天�?, ORANGE),
    ("80+", "专业SKILLs", RGBColor(236, 72, 153)),
    ("15-30%", "RevPAR提升", PURPLE),
    ("50%+", "佣金降低", BLUE)
]

for i, (num, label, color) in enumerate(metrics):
    row = i // 3
    col = i % 3
    left = Inches(1.5) + col * Inches(3.5)
    top = Inches(1.8) + row * Inches(2.5)
    
    card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3), Inches(2))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = color
    
    add_text(s10, num, left, top + Inches(0.3), Inches(3), Inches(1), size=48, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s10, label, left, top + Inches(1.3), Inches(3), Inches(0.5), size=18, color=WHITE, align=PP_ALIGN.CENTER)

# ========== Page 11: Roadmap ==========
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11)
add_bar(s11)
add_page_num(s11, 11, 14)

add_text(s11, "实施路径", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

phases = [
    ("Phase 1", "核心SKILL开�?, ["客房预订+收益", "前厅+客房", "1-2个月"], PURPLE),
    ("Phase 2", "场景SKILL扩展", ["餐饮+宴会+B2B", "营销+会员", "3-6个月"], BLUE),
    ("Phase 3", "生态开�?, ["SKILL市场", "第三方接�?, "7-12个月"], ORANGE)
]

for i, (phase, title, content, color) in enumerate(phases):
    left = Inches(0.6) + i * Inches(4.2)
    
    card_bg = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(4.5))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = CARD_BG
    card_bg.line.color.rgb = color
    
    add_text(s11, phase, left, Inches(1.7), Inches(3.8), Inches(0.5), size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s11, title, left, Inches(2.3), Inches(3.8), Inches(0.5), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    content_box = s11.shapes.add_textbox(left + Inches(0.2), Inches(3), Inches(3.4), Inches(2.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(content):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    if i < 2:
        arrow = s11.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(3.85), Inches(3.5), Inches(0.35), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# ========== Page 12: Team ==========
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12)
add_bar(s12)
add_page_num(s12, 12, 14)

add_text(s12, "核心技术团�?, 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)
add_text(s12, "黄金三角组合：产�?技�?资本", 0, Inches(1.1), Inches(12), Inches(0.4), size=18, color=ORANGE)

team = [
    ("张实", "项目总控", ["24年酒店业", "多家集团高管", "范式革命设计�?], PURPLE),
    ("李源", "CTO", ["华中科大AI博士", "10�?AI研发", "20+工程落地"], BLUE),
    ("陈思序", "CSO", ["500强战�?, "PE/VC背景", "IPO全流�?], ORANGE)
]

for i, (name, role, desc, color) in enumerate(team):
    left = Inches(0.8) + i * Inches(4.2)
    
    card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(3.8), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = color
    
    avatar = s12.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.4), Inches(1.9), Inches(1), Inches(1))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = color
    avatar.line.fill.background()
    
    add_text(s12, name, left, Inches(3), Inches(3.8), Inches(0.5), size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s12, role, left, Inches(3.5), Inches(3.8), Inches(0.4), size=14, color=color, align=PP_ALIGN.CENTER)
    
    content_box = s12.shapes.add_textbox(left + Inches(0.2), Inches(4), Inches(3.4), Inches(2.3))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT
        p.alignment = PP_ALIGN.CENTER

# ========== Page 13: Summary ==========
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s13)
add_bar(s13)
add_page_num(s13, 13, 14)

add_text(s13, "AHL核心价�?, 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True)

values = [
    ("对C�?, ["7x24h管家", "革新体验", "直连无中间商"], BLUE),
    ("对B�?, ["自主运营", "降本增效", "摆脱OTA依赖"], PURPLE),
    ("对平�?, ["2% vs 15%", "新范�?, "技术革�?], ORANGE),
    ("技术壁�?, ["大模�?知识�?, "+SKILLs三位一�?, "护城�?], RGBColor(236, 72, 153))
]

for i, (title, content, color) in enumerate(values):
    left = Inches(0.5) + i * Inches(3.2)
    
    card_bg = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.9), Inches(4.5))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = CARD_BG
    card_bg.line.color.rgb = color
    
    add_text(s13, title, left, Inches(1.8), Inches(2.9), Inches(0.6), size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
    
    content_box = s13.shapes.add_textbox(left + Inches(0.2), Inches(2.6), Inches(2.5), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(content):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# ========== Page 14: Contact ==========
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s14)
add_bar(s14)

c1 = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
c1.fill.solid()
c1.fill.fore_color.rgb = PURPLE
c1.line.fill.background()

c2 = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
c2.fill.solid()
c2.fill.fore_color.rgb = BLUE
c2.line.fill.background()

add_text(s14, "开启住宿业的AI新纪�?, 0, Inches(1.5), prs.slide_width, Inches(0.8), size=44, color=WHITE)
add_text(s14, "联系我们", 0, Inches(2.8), prs.slide_width, Inches(0.6), size=32, bold=True, color=ORANGE)

contacts = ["张实 (Eric Zhang)", "项目总控", "17760348653", "ericzhangshi@163.com", "AHL团队"]
for i, info in enumerate(contacts):
    add_text(s14, info, 0, Inches(4) + i * Inches(0.5), prs.slide_width, Inches(0.5), size=20, color=WHITE)

# Save
output = r"C:\Users\ericz\.openclaw\workspace\ppt_output\AHL-Tech\AHL技术项目说�?pptx"
os.makedirs(os.path.dirname(output), exist_ok=True)
prs.save(output)
print(f"Generated: {output}")
print(f"Total pages: {len(prs.slides)}")
