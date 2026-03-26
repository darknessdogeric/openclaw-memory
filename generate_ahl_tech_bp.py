# -*- coding: utf-8 -*-
"""
AHL-LLM BP Generator - 科技属性政府/投资人版
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

sys.stdout.reconfigure(encoding='utf-8')

# Colors
DARK_BG = RGBColor(5, 15, 35)
DARK_BLUE = RGBColor(10, 25, 50)
PURPLE = RGBColor(124, 58, 237)
BLUE = RGBColor(59, 130, 246)
CYAN = RGBColor(34, 211, 238)
GREEN = RGBColor(34, 197, 94)
ORANGE = RGBColor(249, 115, 22)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(180, 180, 180)
CARD_BG = RGBColor(15, 30, 60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

def add_decoration(slide):
    # Top gradient bar
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    top.fill.solid()
    top.fill.fore_color.rgb = CYAN
    top.line.fill.background()
    
    # Bottom accent
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.04), prs.slide_width, Inches(0.04))
    bot.fill.solid()
    bot.fill.fore_color.rgb = PURPLE
    bot.line.fill.background()

def add_text(slide, text, left, top, width, height, size=20, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

def add_card(slide, left, top, width, height, title, lines, border_color=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    
    # Title
    add_text(slide, title, left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), Inches(0.5), size=16, bold=True, color=border_color, align=PP_ALIGN.CENTER)
    
    # Content
    content_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.6), width - Inches(0.3), height - Inches(0.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

def add_page_num(slide, num, total):
    add_text(slide, f"{num}/{total}", prs.slide_width - Inches(1), prs.slide_height - Inches(0.45), Inches(0.9), Inches(0.4), size=12, color=LIGHT, align=PP_ALIGN.RIGHT)

def add_bullet_text(slide, items, left, top, width, height, size=14, color=WHITE):
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)

# ========== Page 1: Cover ==========
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)

# Decorative elements
c1 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1.5), Inches(5), Inches(5))
c1.fill.solid()
c1.fill.fore_color.rgb = RGBColor(20, 40, 80)
c1.line.fill.background()

c2 = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(4), Inches(4))
c2.fill.solid()
c2.fill.fore_color.rgb = RGBColor(30, 60, 100)
c2.line.fill.background()

# Tech grid lines
for i in range(5):
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5 + i * 0.08), prs.slide_width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(60, 60, 80)
    line.line.fill.background()

add_text(s1, "AHL-LLM", 0, Inches(1.8), prs.slide_width, Inches(1), size=72, bold=True, color=CYAN)
add_text(s1, "AI Hotels Language Model", 0, Inches(2.9), prs.slide_width, Inches(0.6), size=28, bold=True, color=BLUE)
add_text(s1, "大模型驱动的新一代住宿业智能运营平台", 0, Inches(3.6), prs.slide_width, Inches(0.5), size=22, color=WHITE)

# Tech tags
tags = ["国家高新技术企业", "AI大模型", "智能酒店", "科技创新"]
for i, tag in enumerate(tags):
    left = Inches(2.5) + i * Inches(2.2)
    tag_bg = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.5), Inches(2), Inches(0.45))
    tag_bg.fill.solid()
    tag_bg.fill.fore_color.rgb = RGBColor(30, 50, 90)
    tag_bg.line.color.rgb = CYAN
    add_text(s1, tag, left, Inches(4.5), Inches(2), Inches(0.45), size=12, color=CYAN)

add_text(s1, "科技属性 · 创业创新 · AI大模型应用", 0, Inches(5.5), prs.slide_width, Inches(0.4), size=16, color=ORANGE)

# ========== Page 2: Project Overview ==========
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
add_decoration(s2)
add_page_num(s2, 2, 14)

add_text(s2, "项目概述", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

# Core positioning
pos_bg = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.2))
pos_bg.fill.solid()
pos_bg.fill.fore_color.rgb = RGBColor(20, 40, 70)
pos_bg.line.color.rgb = CYAN

add_text(s2, "AHL是国内首个基于大语言模型的住宿业智能运营平台", Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.5), size=20, bold=True, color=WHITE)
add_text(s2, "专注于用AI技术重构酒店、民宿的运营全链路，而非单纯的预订交易平台", Inches(0.7), Inches(2.0), Inches(11.9), Inches(0.4), size=14, color=LIGHT)

# Tech attributes
add_text(s2, "核心科技属性", Inches(0.5), Inches(2.8), Inches(4), Inches(0.4), size=18, bold=True, color=PURPLE)

tech_items = [
    ("🤖 AI大模型", "基于DeepSeek/Qwen开源模型微调\n酒店行业知识库注入\nLoRA低成本领域适配"),
    ("🔧 AGENT架构", "C端AI管家 + B端AI运营官\n7×24h全天候服务\n自主决策与执行"),
    ("🧩 SKILL体系", "80+可插拔专业能力\n像乐高一样灵活组合\n周级迭代速度"),
    ("⚡ 向量引擎", "<3秒响应速度\n95%+匹配准确率\n实时推理与决策")
]

for i, (title, content) in enumerate(tech_items):
    left = Inches(0.5) + i * Inches(3.15)
    lines = content.split('\n')
    add_card(s2, left, Inches(3.3), Inches(3), Inches(3), title, lines, PURPLE if i%2==0 else BLUE)

# ========== Page 3: Tech Architecture ==========
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
add_decoration(s3)
add_page_num(s3, 3, 14)

add_text(s3, "四层技术架构", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

layers = [
    ("第一层：LLM大模型", "底层智能", PURPLE, ["DeepSeek/Qwen开源模型", "LoRA微调: 低成本领域适配", "PP&SOP知识库: 12,000+KB", "多模态: 文本/语音/图像"]),
    ("第二层：双AGENT", "执行智能", BLUE, ["C端AI管家: 对客服务", "B端AI运营官: 运营决策", "7×24h全天候服务", "自主学习持续进化"]),
    ("第三层：SKILLs", "能力智能", CYAN, ["80+可插拔专业能力", "覆盖前厅/客房/餐饮", "营销/收益/财务等", "周级迭代快速上线"]),
    ("第四层：向量引擎", "匹配智能", GREEN, ["Milvus向量数据库", "<3秒响应速度", "95%+匹配准确率", "实时推理决策"])
]

for i, (title, subtitle, color, content) in enumerate(layers):
    top = Inches(1.3) + i * Inches(1.5)
    
    # Layer bar
    layer_bg = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12.3), Inches(1.4))
    layer_bg.fill.solid()
    layer_bg.fill.fore_color.rgb = RGBColor(15, 30, 55)
    layer_bg.line.color.rgb = color
    layer_bg.line.width = Pt(2)
    
    # Layer number
    num_circle = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.3), Inches(0.8), Inches(0.8))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = color
    num_circle.line.fill.background()
    add_text(s3, str(i+1), Inches(0.7), top + Inches(0.4), Inches(0.8), Inches(0.6), size=24, bold=True, color=WHITE)
    
    # Title
    add_text(s3, title, Inches(1.7), top + Inches(0.15), Inches(4), Inches(0.5), size=18, bold=True, color=color)
    add_text(s3, subtitle, Inches(1.7), top + Inches(0.7), Inches(4), Inches(0.4), size=12, color=LIGHT)
    
    # Content
    content_box = s3.shapes.add_textbox(Inches(6), top + Inches(0.2), Inches(6.5), Inches(1.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = " • ".join(content[:2])
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

# ========== Page 4: Market Opportunity ==========
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4)
add_decoration(s4)
add_page_num(s4, 4, 14)

add_text(s4, "市场机会：AI赋能传统行业", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=36, bold=True, color=CYAN)
add_text(s4, "万亿级传统行业升级机遇", 0, Inches(1.1), prs.slide_width, Inches(0.4), size=18, color=ORANGE)

# Market data
market_data = [
    ("5000亿+", "中国酒店市场规模", "传统行业升级需求旺盛"),
    ("750亿", "年OTA佣金池", "降本增效空间巨大"),
    ("45万+", "酒店/民宿数量", "中小商户亟需赋能"),
    ("260%", "5年获客成本涨幅", "行业痛点亟待解决")
]

for i, (num, title, desc) in enumerate(market_data):
    left = Inches(0.5) + i * Inches(3.15)
    
    card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(3), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(20, 40, 70)
    card.line.color.rgb = CYAN
    
    add_text(s4, num, left, Inches(1.9), Inches(3), Inches(0.8), size=36, bold=True, color=CYAN)
    add_text(s4, title, left, Inches(2.7), Inches(3), Inches(0.4), size=14, bold=True, color=WHITE)
    add_text(s4, desc, left, Inches(3.1), Inches(3), Inches(0.6), size=11, color=LIGHT)

# Policy support
add_text(s4, "政策利好", Inches(0.5), Inches(4.2), Inches(4), Inches(0.4), size=18, bold=True, color=GREEN)

policies = [
    ("国家高新技术企业认定", "15%优惠税率 + 研发加计扣除"),
    ("AI大模型扶持政策", "算力补贴 + 场景开放"),
    ("科技创新专项资金", "最高3000万无偿资助"),
    ("创业创新支持", "人才引进 + 场地补贴")
]

for i, (title, desc) in enumerate(policies):
    left = Inches(0.5) + i * Inches(3.15)
    add_card(s4, left, Inches(4.7), Inches(3), Inches(1.6), title, [desc], GREEN if i%2==0 else BLUE)

# ========== Page 5: Business Model ==========
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5)
add_decoration(s5)
add_page_num(s5, 5, 14)

add_text(s5, "商业模式：电信运营模式", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)
add_text(s5, "类比电信运营商，商户按需订阅，平台稳定营收", 0, Inches(1.1), Inches(12), Inches(0.4), size=16, color=LIGHT)

# Three layers
layers_bm = [
    ("C端用户", "酒店/民宿", "AHL平台", "类比", "手机用户", "基础运营商", "设备提供商"),
    ("按次/套餐", "月订阅+TOKEN", "2%交易费", "收费", "语音分钟", "月租+流量", "设备折旧")
]

# B2B model
add_text(s5, "商户订阅套餐（核心收入）", Inches(0.5), Inches(1.8), Inches(6), Inches(0.4), size=18, bold=True, color=PURPLE)

suits = [
    ("入门版", "¥999/月", ["基础C端SKILLs", "100K TOKENS", "1个B端SKILL"], PURPLE),
    ("专业版", "¥2,999/月", ["全部37个C端SKILLs", "500K TOKENS", "8个B端SKILLs"], BLUE),
    ("企业版", "¥9,999/月", ["无限C端SKILLs", "2M TOKENS", "全部25个B端SKILLs"], CYAN)
]

for i, (name, price, features, color) in enumerate(suits):
    left = Inches(0.5) + i * Inches(4.2)
    
    card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.3), Inches(4), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    add_text(s5, name, left, Inches(2.5), Inches(4), Inches(0.5), size=20, bold=True, color=color)
    add_text(s5, price, left, Inches(3.0), Inches(4), Inches(0.6), size=28, bold=True, color=WHITE)
    
    content_box = s5.shapes.add_textbox(left + Inches(0.2), Inches(3.7), Inches(3.6), Inches(1.3))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, feat in enumerate(features):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + feat
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

# Revenue model
add_text(s5, "收入结构", Inches(0.5), Inches(5.3), Inches(4), Inches(0.4), size=18, bold=True, color=GREEN)

revenue_items = [
    "商户月订阅费: 40%",
    "超额TOKEN费: 20%", 
    "交易服务费: 30%",
    "SKILL增值: 10%"
]
add_bullet_text(s5, revenue_items, Inches(0.5), Inches(5.8), Inches(6), Inches(1.5), size=14, color=WHITE)

# Comparison
add_text(s5, "vs 传统OTA", Inches(7), Inches(5.3), Inches(5), Inches(0.4), size=18, bold=True, color=ORANGE)

comparison_items = [
    "OTA费率: 15-25%",
    "AHL费率: 2%+套餐",
    "节省成本: 80%+"
]
add_bullet_text(s5, comparison_items, Inches(7), Inches(5.8), Inches(5), Inches(1.5), size=14, color=WHITE)

# ========== Page 6: Tech Barrier ==========
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6)
add_decoration(s6)
add_page_num(s6, 6, 14)

add_text(s6, "技术壁垒与竞争优势", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

barriers = [
    ("数据飞轮", "商户越多 → 数据越丰富 → 匹配越精准 → 用户越多", PURPLE),
    ("网络效应", "双边市场：商户端供给丰富 + 用户端需求活跃", BLUE),
    ("知识壁垒", "PP&SOP知识库: 77个文档, 2000+KB, 持续积累", CYAN),
    ("技术壁垒", "开源模型 + 酒店微调 + 向量引擎 + AGENT架构", GREEN)
]

for i, (title, desc, color) in enumerate(barriers):
    top = Inches(1.3) + i * Inches(1.4)
    
    card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12.3), Inches(1.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # Icon circle
    icon = s6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.25), Inches(0.8), Inches(0.8))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()
    
    add_text(s6, str(i+1), Inches(0.7), top + Inches(0.35), Inches(0.8), Inches(0.6), size=20, bold=True, color=WHITE)
    add_text(s6, title, Inches(1.7), top + Inches(0.2), Inches(3), Inches(0.5), size=18, bold=True, color=color)
    add_text(s6, desc, Inches(1.7), top + Inches(0.7), Inches(10.5), Inches(0.5), size=14, color=WHITE)

# Comparison table
add_text(s6, "竞品对比", Inches(0.5), Inches(6.9), Inches(4), Inches(0.4), size=16, bold=True, color=ORANGE)

comparison = [
    ("传统OTA", "15-25%费率", "数据垄断", "有限AI"),
    ("国际品牌", "10-15%费率", "系统封闭", "定制昂贵"),
    ("AHL", "2%+套餐", "数据自有", "80+SKILLs")
]

for i, (name, fee, data, ai) in enumerate(comparison):
    left = Inches(0.5) + i * Inches(4.2)
    color = ORANGE if i == 2 else LIGHT
    add_text(s6, name, left, Inches(7.3), Inches(4), Inches(0.3), size=12, bold=True, color=color)
    add_text(s6, fee + " | " + data + " | " + ai, left, Inches(7.55), Inches(4), Inches(0.3), size=10, color=WHITE)

# ========== Page 7: Implementation ==========
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7)
add_decoration(s7)
add_page_num(s7, 7, 14)

add_text(s7, "实施路径", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

phases = [
    ("Phase 1", "M1-M3", "产品孵化", ["核心LLM微调完成", "MVP上线测试", "种子商户10+"], PURPLE),
    ("Phase 2", "M4-M6", "模式验证", ["套餐模式验证", "商户50+上线", "月GMV100万"], BLUE),
    ("Phase 3", "M7-M12", "规模复制", ["商户200+", "区域覆盖", "月GMV500万"], CYAN),
    ("Phase 4", "M13-M18", "生态开放", ["SKILL市场开放", "商户500+", "A轮启动"], GREEN)
]

for i, (phase, time, title, goals, color) in enumerate(phases):
    left = Inches(0.5) + i * Inches(3.15)
    
    card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    add_text(s7, phase, left, Inches(1.7), Inches(3), Inches(0.4), size=14, bold=True, color=color)
    add_text(s7, time, left, Inches(2.1), Inches(3), Inches(0.4), size=20, bold=True, color=WHITE)
    add_text(s7, title, left, Inches(2.6), Inches(3), Inches(0.5), size=16, bold=True, color=color)
    
    content_box = s7.shapes.add_textbox(left + Inches(0.15), Inches(3.2), Inches(2.7), Inches(2.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, goal in enumerate(goals):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + goal
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
    
    if i < 3:
        arrow = s7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(3.05), Inches(3.5), Inches(0.1), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# ========== Page 8: Use of Funds ==========
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8)
add_decoration(s8)
add_page_num(s8, 8, 14)

add_text(s8, "融资计划", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

# Funding details
add_text(s8, "种子轮", Inches(0.5), Inches(1.3), Inches(4), Inches(0.4), size=18, bold=True, color=PURPLE)

funding = [
    ("融资规模", "500-800万"),
    ("出让股份", "15-20%"),
    ("投后估值", "3000-5000万"),
    ("投资形式", "天使/机构")
]

for i, (label, value) in enumerate(funding):
    top = Inches(1.8) + i * Inches(0.7)
    add_text(s8, label + ":", Inches(0.5), top, Inches(2), Inches(0.5), size=16, color=LIGHT)
    add_text(s8, value, Inches(2.5), top, Inches(3), Inches(0.5), size=18, bold=True, color=WHITE)

# Use of funds
add_text(s8, "资金用途", Inches(0.5), Inches(4.7), Inches(4), Inches(0.4), size=18, bold=True, color=GREEN)

uses = [
    ("技术研发", "40%", "200-320万", ["LLM微调", "AGENT开发", "向量引擎"]),
    ("数据体系", "30%", "150-240万", ["知识库建设", "数据标注", "模型训练"]),
    ("市场拓展", "20%", "100-160万", ["BD团队", "试点城市", "品牌推广"]),
    ("运营储备", "10%", "50-80万", ["法务合规", "资质申请", "团队建设"])
]

for i, (title, pct, amount, items) in enumerate(uses):
    left = Inches(0.5) + i * Inches(3.15)
    
    card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.2), Inches(3), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = GREEN if i == 0 else (BLUE if i == 1 else (CYAN if i == 2 else ORANGE))
    
    add_text(s8, title, left, Inches(5.3), Inches(3), Inches(0.4), size=14, bold=True, color=WHITE)
    add_text(s8, pct + " | " + amount, left, Inches(5.7), Inches(3), Inches(0.4), size=12, color=LIGHT)
    
    content_box = s8.shapes.add_textbox(left + Inches(0.1), Inches(6.1), Inches(2.8), Inches(0.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE

# ========== Page 9: Team ==========
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9)
add_decoration(s9)
add_page_num(s9, 9, 14)

add_text(s9, "核心团队", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

team = [
    ("张实", "项目总控", "24年酒店业深耕", "北京第二外国语学院", "多家知名酒店集团高管", "住宿业范式革命设计师", PURPLE),
    ("李源", "首席技术官 CTO", "华中科技大学AI博士", "10年+ AI技术研发", "20+ AI工程落地", "AI架构与算法专家", BLUE),
    ("陈思序", "首席战略官 CSO", "世界500强战略高管", "深耕PE/VC投资", "IPO全流程经验", "战略推动与生态建设", CYAN)
]

for i, (name, role, line1, line2, line3, line4, color) in enumerate(team):
    left = Inches(0.5) + i * Inches(4.2)
    
    # Card
    card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(4), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # Avatar
    avatar = s9.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.4), Inches(1.7), Inches(1.2), Inches(1.2))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = color
    avatar.line.fill.background()
    
    add_text(s9, name, left, Inches(3.0), Inches(4), Inches(0.5), size=24, bold=True, color=WHITE)
    add_text(s9, role, left, Inches(3.5), Inches(4), Inches(0.4), size=14, bold=True, color=color)
    
    lines = [line1, line2, line3, line4]
    content_box = s9.shapes.add_textbox(left + Inches(0.2), Inches(4.0), Inches(3.6), Inches(2.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(lines):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT

# ========== Page 10: Achievements ==========
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10)
add_decoration(s10)
add_page_num(s10, 10, 14)

add_text(s10, "已取得成果", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

achievements = [
    ("技术成果", ["LLM大模型微调完成", "PP&SOP知识库77个文档", "80+ SKILLs架构设计", "向量匹配引擎原型"]),
    ("商业成果", ["完整商业计划书 V3.0", "融资材料完备", "60+投资人清单", "合作模式设计完成"]),
    ("行业积累", ["酒店行业全景知识库", "12层架构体系", "77个SOP文档", "多品牌对标分析"]),
    ("团队组建", ["创始团队3人到位", "技术顾问资源", "行业专家顾问", "战略合作伙伴"])
]

for i, (title, items) in enumerate(achievements):
    left = Inches(0.5) + i * Inches(3.15)
    add_card(s10, left, Inches(1.4), Inches(3), Inches(4.5), title, items, PURPLE if i%2==0 else BLUE)

# Metrics
add_text(s10, "关键数据", Inches(0.5), Inches(6.2), Inches(4), Inches(0.4), size=18, bold=True, color=GREEN)

metrics = [("知识库", "2000+KB"), ("SKILLs", "80+个"), ("SOP文档", "77个"), ("品牌对标", "10+个")]

for i, (label, value) in enumerate(metrics):
    left = Inches(0.5) + i * Inches(3.15)
    add_text(s10, value, left, Inches(6.6), Inches(3), Inches(0.5), size=24, bold=True, color=CYAN)
    add_text(s10, label, left, Inches(7.05), Inches(3), Inches(0.3), size=12, color=LIGHT)

# ========== Page 11: Gov Support ==========
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11)
add_decoration(s11)
add_page_num(s11, 11, 14)

add_text(s11, "可申报扶持与资质", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)
add_text(s11, "科技型中小企业 + 高新技术企业 + 专精特新", 0, Inches(1.1), Inches(12), Inches(0.4), size=16, color=ORANGE)

supports = [
    ("科技型中小企业", "¥5-20万", ["研发费用加计扣除", "技术创新专项", "人才引进补贴"], PURPLE),
    ("高新技术企业", "¥10-50万", ["15%优惠税率", "研发费用补贴", "人才安居政策"], BLUE),
    ("专精特新企业", "¥20-100万", ["梯度培育支持", "融资担保服务", "市场拓展补贴"], CYAN),
    ("AI大模型专项", "最高¥500万", ["算力补贴政策", "场景开放支持", "模型备案奖励"], GREEN)
]

for i, (title, amount, features, color) in enumerate(supports):
    left = Inches(0.5) + i * Inches(3.15)
    
    card = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(3), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    add_text(s11, title, left, Inches(1.9), Inches(3), Inches(0.5), size=16, bold=True, color=color)
    add_text(s11, amount, left, Inches(2.4), Inches(3), Inches(0.6), size=24, bold=True, color=WHITE)
    
    content_box = s11.shapes.add_textbox(left + Inches(0.15), Inches(3.1), Inches(2.7), Inches(2))
    tf = content_box.text_frame
    tf.word_wrap = True
    for j, feat in enumerate(features):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + feat
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

# Local support
add_text(s11, "地方配套支持", Inches(0.5), Inches(5.5), Inches(4), Inches(0.4), size=18, bold=True, color=ORANGE)

local = [
    "创业孵化载体: 免租1-3年",
    "人才公寓: 最高100%租金补贴",
    "场地租金: 最高50%补贴",
    "个税返还: 地方留成部分补贴"
]
add_bullet_text(s11, local, Inches(0.5), Inches(6.0), Inches(6), Inches(1.5), size=14, color=WHITE)

# ========== Page 12: Risk ==========
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12)
add_decoration(s12)
add_page_num(s12, 12, 14)

add_text(s12, "风险与对策", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

risks = [
    ("巨头竞争", "互联网巨头可能入局", "中小商户差异化定位", "专注垂直领域深耕"),
    ("技术迭代", "大模型技术快速演进", "保持技术敏感度", "模块化架构快速适配"),
    ("市场接受度", "传统商户认知转变慢", "标杆案例轻资产验证", "逐步渗透建立信任"),
    ("政策监管", "AI和数据监管趋严", "合规前置主动申请资质", "数据安全优先设计")
]

for i, (risk, desc, solution_title, solution) in enumerate(risks):
    top = Inches(1.3) + i * Inches(1.5)
    color = ORANGE
    
    card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(12.3), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    
    # Risk
    add_text(s12, "风险: " + risk, Inches(0.7), top + Inches(0.1), Inches(4), Inches(0.4), size=14, bold=True, color=ORANGE)
    add_text(s12, desc, Inches(0.7), top + Inches(0.5), Inches(4), Inches(0.6), size=11, color=LIGHT)
    
    # Arrow
    arrow = s12.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5), top + Inches(0.5), Inches(0.4), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()
    
    # Solution
    add_text(s12, "对策: " + solution_title, Inches(5.6), top + Inches(0.1), Inches(6.5), Inches(0.4), size=14, bold=True, color=GREEN)
    add_text(s12, solution, Inches(5.6), top + Inches(0.5), Inches(6.5), Inches(0.6), size=11, color=WHITE)

# ========== Page 13: Milestones ==========
s13 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s13)
add_decoration(s13)
add_page_num(s13, 13, 14)

add_text(s13, "里程碑承诺", 0, Inches(0.3), prs.slide_width, Inches(0.8), size=40, bold=True, color=CYAN)

milestones = [
    ("M3", "产品上线", "注册商户50+", "技术验证完成", PURPLE),
    ("M6", "模式验证", "月GMV100万", "种子用户验证", BLUE),
    ("M12", "A轮启动", "商户200+", "月GMV500万", CYAN),
    ("M18", "区域领先", "商户500+", "战略盈利", GREEN),
    ("M36", "行业标杆", "商户2000+", "市场占有率10%", ORANGE)
]

for i, (time, title, metric1, metric2, color) in enumerate(milestones):
    left = Inches(0.5) + i * Inches(2.5)
    
    card = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.4), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(15, 30, 55)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    add_text(s13, time, left, Inches(1.7), Inches(2.4), Inches(0.5), size=24, bold=True, color=color)
    add_text(s13, title, left, Inches(2.3), Inches(2.4), Inches(0.5), size=16, bold=True, color=WHITE)
    
    content_box = s13.shapes.add_textbox(left + Inches(0.1), Inches(3.0), Inches(2.2), Inches(2.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [metric1, metric2]
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
    
    if i < 4:
        arrow = s13.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.4), Inches(3.5), Inches(0.1), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT
        arrow.line.fill.background()

# ========== Page 14: Contact ==========
s14 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s14)

# Decorative
c1 = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
c1.fill.solid()
c1.fill.fore_color.rgb = RGBColor(20, 40, 80)
c1.line.fill.background()

c2 = s14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
c2.fill.solid()
c2.fill.fore_color.rgb = RGBColor(30, 60, 100)
c2.line.fill.background()

# Tech lines
for i in range(3):
    line = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5 + i * 0.08), prs.slide_width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(60, 60, 80)
    line.line.fill.background()

add_text(s14, "开启AI新时代", 0, Inches(1.8), prs.slide_width, Inches(0.8), size=52, bold=True, color=CYAN)
add_text(s14, "用科技重新定义住宿业", 0, Inches(2.7), prs.slide_width, Inches(0.6), size=28, color=WHITE)

# Contact info
contact_items = [
    "张实 (Eric Zhang)",
    "项目总控 / 创始人",
    "📞 17760348653",
    "📧 ericzhangshi@163.com",
    "AHL团队"
]

for i, item in enumerate(contact_items):
    is_title = (i == 0)
    size = 24 if is_title else 16
    color = CYAN if is_title else WHITE
    add_text(s14, item, 0, Inches(3.8) + i * Inches(0.6), prs.slide_width, Inches(0.5), size=size, bold=is_title, color=color)

# QR placeholder
qr = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(4.5), Inches(2), Inches(2))
qr.fill.solid()
qr.fill.fore_color.rgb = RGBColor(40, 60, 100)
qr.line.color.rgb = CYAN
add_text(s14, "扫码联系", Inches(10.5), Inches(5.3), Inches(2), Inches(0.4), size=14, color=CYAN)

# Save
output = r"C:\Users\ericz\.openclaw\workspace\ppt_output\AHL-Tech\AHL-BP-科技属性版.pptx"
os.makedirs(os.path.dirname(output), exist_ok=True)
prs.save(output)
print(f"Generated: {output}")
print(f"Total pages: {len(prs.slides)}")
