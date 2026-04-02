from pptx import Presentation
import os

files = [
    r'C:\Users\ericz\Desktop\商业计划书\AHL智能升级项目试点合作建议书.pptx',
    r'C:\Users\ericz\Desktop\商业计划书\AHL智能升级项目试点合作建议书2.pptx'
]

for f in files:
    print('='*60)
    print(os.path.basename(f))
    print('='*60)
    prs = Presentation(f)
    print(f'总页数: {len(prs.slides)}')
    print(f'尺寸: {prs.slide_width.inches:.1f}" x {prs.slide_height.inches:.1f}"')
    
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text[:60].replace('\n', ' '))
        title = texts[0] if texts else '(无标题)'
        print(f'  第{i:2d}页: {title}')
    print()
