from pptx import Presentation
from pptx.util import Inches
import os

files = [
    (r'C:\Users\ericz\Desktop\商业计划书\AHL智能升级项目试点合作建议书.pptx', '版本1 (精简版)'),
    (r'C:\Users\ericz\Desktop\商业计划书\AHL智能升级项目试点合作建议书2.pptx', '版本2 (完整版)')
]

for fpath, fname in files:
    print('='*70)
    print(f'文件: {fname}')
    print(f'大小: {os.path.getsize(fpath)/1024/1024:.1f} MB')
    print('='*70)
    prs = Presentation(fpath)
    print(f'页数: {len(prs.slides)}')
    print(f'尺寸: {prs.slide_width.inches:.1f}" x {prs.slide_height.inches:.1f}"')
    
    for i, slide in enumerate(prs.slides, 1):
        # Count shapes
        shape_count = len(slide.shapes)
        
        # Get all text
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                t = shape.text.strip()
                if t:
                    all_text.append(t[:100])
        
        # Count images
        pic_count = sum(1 for s in slide.shapes if s.shape_type == 13)  # MSO_SHAPE_TYPE.PICTURE
        
        print(f'\n第{i}页:')
        print(f'  元素数: {shape_count}, 图片: {pic_count}')
        if all_text:
            print(f'  内容: {all_text[0][:80]}')
            if len(all_text) > 1:
                print(f'  ...共{len(all_text)}个文本块')
        else:
            print(f'  内容: (无文本，可能为纯图片页)')
    print()
